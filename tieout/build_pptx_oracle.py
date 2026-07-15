"""Generate the PPTX subsystem parity oracle for the Rust `fm-pptx` port.

Mirrors `tieout/build_ev_bridge_oracle.py`. Produces, from the reference
Python `src/research/pptx_*.py` implementations, the deterministic fixtures
and JSON snapshots the Rust parity gates diff against:

  6.1 inspector  — `inspect_pptx_json` over fixture decks
                   -> tieout/excel_snapshots/PPTX_inspect_<name>.json
  6.2 editor     — duplicate/delete/reorder/recolor/replace-text round-trips;
                   normalized control-file members + edit-log
                   -> tieout/excel_snapshots/PPTX_edit_<op>.json
  6.3 writer     — pure functions over an input matrix
                   -> tieout/excel_snapshots/PPTX_pure.json
  6.4 drawingml  — write_ev_bridge_deck / write_ifrs_bridge_deck reference
                   decks + their inspected shape trees
                   -> tests/fixtures/pptx/*.pptx
                   -> tieout/excel_snapshots/PPTX_deck_<name>.json

Fixture decks are committed to tests/fixtures/pptx/ so the Rust inspector
walks the exact same bytes the Python oracle inspected.

Run:  py tieout/build_pptx_oracle.py
"""
import json
import sys
from pathlib import Path

REPO = Path(__file__).parent.parent.resolve()
sys.path.insert(0, str(REPO))

SNAP_DIR = REPO / "tieout/excel_snapshots"
FIX_DIR = REPO / "tests/fixtures/pptx"
SNAP_DIR.mkdir(parents=True, exist_ok=True)
FIX_DIR.mkdir(parents=True, exist_ok=True)

# Pinned tiny 8x8 red PNG (from tests/test_pptx_editor.py::_make_test_png).
_PNG_HEX = (
    "89504e470d0a1a0a0000000d4948445200000008000000080806000000c40fbe"
    "8b0000001b49444154789c63fcffff3f032003200a30c00100c40c0440c30202"
    "00f5170106e7e7c0c30000000049454e44ae426082"
)


def _make_test_png(path: Path) -> Path:
    path.write_bytes(bytes.fromhex(_PNG_HEX))
    return path


def build_three_slide_deck(path: Path) -> Path:
    """Slide 0: title. Slide 1: text+image. Slide 2: text only.

    Verbatim from tests/test_pptx_editor.py::_make_three_slide_deck so the
    Rust editor gate round-trips the exact fixture the Python tests use.
    """
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    s0 = prs.slides.add_slide(blank)
    tb0 = s0.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    tb0.name = "Title"
    tb0.text_frame.text = "Cover Title"

    s1 = prs.slides.add_slide(blank)
    tb1 = s1.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    tb1.name = "Body"
    tb1.text_frame.text = "Q1 revenue grew strongly"
    img = _make_test_png(FIX_DIR / "_logo.png")
    pic = s1.shapes.add_picture(str(img), Inches(8), Inches(0.5), Inches(1.5), Inches(1.5))
    pic.name = "Logo"

    s2 = prs.slides.add_slide(blank)
    tb2 = s2.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    tb2.name = "Outro"
    tb2.text_frame.text = "Thank you — Q1 wrap-up"

    prs.save(str(path))
    return path


# ── 6.1 inspector ─────────────────────────────────────────────────────────────

def oracle_inspect() -> None:
    from src.research.pptx_inspector import inspect_pptx_json

    deck = FIX_DIR / "deck.pptx"
    build_three_slide_deck(deck)
    js = inspect_pptx_json(deck)  # include_raw_xml=False, max_xml_chars=600
    (SNAP_DIR / "PPTX_inspect_deck.json").write_text(js, encoding="utf-8")
    data = json.loads(js)
    print(f"  6.1 inspect deck.pptx: slides={data['slideCount']} "
          f"layouts={data['layoutCount']} masters={data['masterCount']} "
          f"-> PPTX_inspect_deck.json")


# ── 6.2 editor ────────────────────────────────────────────────────────────────

_CTRL = ["ppt/presentation.xml", "ppt/_rels/presentation.xml.rels", "[Content_Types].xml"]


def _members(path, names):
    import zipfile
    out = {}
    with zipfile.ZipFile(path) as z:
        have = set(z.namelist())
        for n in names:
            if n in have:
                out[n] = z.read(n).decode("utf-8")
    return out


def _slide_parts(path):
    import zipfile
    with zipfile.ZipFile(path) as z:
        return sorted(
            (n for n in z.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml")),
            key=lambda n: int(n.rsplit("slide", 1)[1].rsplit(".", 1)[0]),
        )


def _last_log(deck_out):
    from src.research.pptx_editor import get_edit_history
    hist = get_edit_history(deck_out, last_n=1)
    if not hist:
        return None
    e = dict(hist[-1])
    e.pop("ts", None)
    params = dict(e.get("params", {}))
    params.pop("output_path", None)
    e["params"] = params
    return e


def oracle_edit() -> None:
    import shutil
    from src.research import pptx_editor as ed

    src = FIX_DIR / "deck.pptx"
    work = FIX_DIR / "_edit_work"
    work.mkdir(exist_ok=True)

    def run(tag, op, control_only=False):
        deck = work / f"{tag}.pptx"
        shutil.copy2(src, deck)
        ed.clear_edit_history(str(deck))
        out = op(str(deck))
        snap = {
            "op": _last_log(out).get("op") if _last_log(out) else tag,
            "log": _last_log(out),
            "slide_parts": _slide_parts(out),
            "control": _members(out, _CTRL),
        }
        if not control_only:
            snap["slides"] = _members(out, _slide_parts(out))
        (SNAP_DIR / f"PPTX_edit_{tag}.json").write_text(
            json.dumps(snap, indent=2), encoding="utf-8"
        )
        print(f"  6.2 {tag}: slides={len(snap['slide_parts'])} -> PPTX_edit_{tag}.json")

    run("duplicate", lambda d: ed.duplicate_slide(d, 1))
    run("delete", lambda d: ed.delete_slide(d, 1))
    run("reorder", lambda d: ed.reorder_slides(d, [2, 0, 1]))
    # recolor: only the theme part changes; capture theme members.
    def _recolor(d):
        return ed.recolor_theme(d, {"accent1": "#255BE3", "accent2": "#0F1632"})
    deck = work / "recolor.pptx"
    shutil.copy2(src, deck)
    ed.clear_edit_history(str(deck))
    out = _recolor(str(deck))
    import zipfile
    with zipfile.ZipFile(out) as z:
        themes = [n for n in z.namelist() if n.startswith("ppt/theme/") and n.endswith(".xml")]
    snap = {"op": "recolor_theme", "log": _last_log(out), "theme": _members(out, themes)}
    (SNAP_DIR / "PPTX_edit_recolor.json").write_text(json.dumps(snap, indent=2), encoding="utf-8")
    print(f"  6.2 recolor: themes={len(themes)} -> PPTX_edit_recolor.json")
    # replace_text: gated behaviourally (inspector), capture only the log.
    deck = work / "replace.pptx"
    shutil.copy2(src, deck)
    ed.clear_edit_history(str(deck))
    out = ed.replace_text_in_deck(str(deck), {"Q1": "Q2"})
    snap = {"op": "replace_text_in_deck", "log": _last_log(out)}
    (SNAP_DIR / "PPTX_edit_replace.json").write_text(json.dumps(snap, indent=2), encoding="utf-8")
    print("  6.2 replace_text: -> PPTX_edit_replace.json")
    shutil.rmtree(work, ignore_errors=True)

# ── 6.3 writer pure functions ─────────────────────────────────────────────────

_MD_SAMPLE = """\
type: cover
title: Sandvik AB Investment Memo
subtitle: Industrials | Long
date: April 2026
---
type: bar_chart
action_title: Sandvik trades at a discount to peer median
labels: [SAND.ST, CAT, KMT, ITW]
values: [10.5, 11.0, 9.2, 11.8]
target_label: SAND.ST
value_format: "{:.1f}x"
x_label: EV / LTM EBITDA
source: Bloomberg, Apr 30 2026
"""


def oracle_pure() -> None:
    from src.research.pptx_writer import (
        pick_slide_archetype, split_into_chunks, parse_deck_markdown,
        PPTXDeckWriter,
    )
    from dataclasses import asdict

    # pick_slide_archetype matrix.
    arch_cases = [
        ("quotes", 4, 1, False, False),
        ("comparison", 5, 3, True, False),
        ("events", 6, 1, False, False),
        ("comparison", 5, 3, False, True),
        ("framework", 3, 1, False, False),
        ("strategy", 6, 1, False, False),
        ("process", 9, 1, False, False),
        ("structure", 3, 1, False, False),
        ("single_stat", 1, 5, False, False),
        ("comparison", 1, 4, False, False),
        ("comparison", 10, 3, False, False),
        ("breakdown", 5, 5, False, False),
        ("comparison", 8, 8, False, False),
    ]
    arch_out = []
    for ds, ne, nm, hq, isd in arch_cases:
        d = pick_slide_archetype(
            data_shape=ds, n_entities=ne, n_metrics=nm, has_quotes=hq, is_dated=isd)
        arch_out.append({"in": [ds, ne, nm, hq, isd], "out": asdict(d)})

    # split_into_chunks matrix.
    chunk_cases = [
        (list(range(20)), "comparison_matrix"),
        (["a", "b", "c", "d", "e"], "scorecard"),
        (list(range(10)), "timeline"),
        (list(range(9)), "process_diagram"),
    ]
    chunk_out = [{"in": [items, arch], "out": split_into_chunks(items, arch)}
                 for items, arch in chunk_cases]

    # normalize_heading matrix.
    head_cases = [
        "sandvik trades at a discount to peers",
        "Investment Of The Year",
        "EBITDA margin expands 200bps in FY26",
        "the iPhone drives Apple's Q1 growth",
        "McKinsey and Company vs Bain",
        "u.s. market share of the top 3 players",
    ]
    head_out = [{"in": t, "out": PPTXDeckWriter._normalize_heading(t)} for t in head_cases]

    # fmt_to_numfmt matrix.
    fmt_cases = ["{:+,.0f}", "{:,.0f}", "{:.1f}", "{:.1%}", "{:.1f}%", "{:.1f}x",
                 "{:,.1f}", "{:,.2f}", "{:.0f}", "{:.2%}", "{:.0%}", "{:,}", "{:.2f}x", "weird"]
    fmt_out = [{"in": f, "out": PPTXDeckWriter._fmt_to_numfmt(f)} for f in fmt_cases]

    # format_value matrix (typed).
    fv_cases = [
        ("null", None), ("str", "Hello"), ("float", 0.25), ("float", 0.0),
        ("float", 1234.5), ("float", -0.5), ("float", 5.0), ("float", -1234.6),
        ("int", 1000000), ("int", -42), ("float", 12.0), ("float", 0.125),
    ]
    fv_out = [{"kind": k, "in": v, "out": PPTXDeckWriter._format_value(v)} for k, v in fv_cases]

    specs = parse_deck_markdown(_MD_SAMPLE)

    snap = {
        "pick_slide_archetype": arch_out,
        "split_into_chunks": chunk_out,
        "normalize_heading": head_out,
        "fmt_to_numfmt": fmt_out,
        "format_value": fv_out,
        "parse_deck_markdown": specs,
    }
    (SNAP_DIR / "PPTX_pure.json").write_text(json.dumps(snap, indent=2), encoding="utf-8")
    print(f"  6.3 pure fns: archetype={len(arch_out)} chunks={len(chunk_out)} "
          f"headings={len(head_out)} fmt={len(fmt_out)} format_value={len(fv_out)} "
          f"specs={len(specs)} -> PPTX_pure.json")


# ── 6.4 DrawingML decks ───────────────────────────────────────────────────────

def _pin_date():
    """Pin date.today() -> 2026-01-01 so cover date text is deterministic."""
    import datetime as _dt
    import src.research.pptx_writer as W

    class _FixedDate(_dt.date):
        @classmethod
        def today(cls):
            return _dt.date(2026, 1, 1)
    W.date = _FixedDate


def _ev_input():
    from kb.ev_bridge import EVBridgeInput
    return EVBridgeInput(
        company="DemoCo", period="LTM Sep-25", currency="USD",
        share_price=150.0, shares_outstanding=1_000_000_000.0,
        total_debt=50_000_000_000.0, finance_leases=5_000_000_000.0,
        operating_leases=8_000_000_000.0, underfunded_pension=2_000_000_000.0,
        minority_interest=1_000_000_000.0, preferred_stock=500_000_000.0,
        cash=20_000_000_000.0, short_term_investments=10_000_000_000.0,
        equity_investments=3_000_000_000.0,
        ltm_revenue=100_000_000_000.0, ltm_ebitda=30_000_000_000.0,
    )


def oracle_decks() -> None:
    from types import SimpleNamespace
    from src.research.pptx_inspector import inspect_pptx
    from src.research.pptx_output import ResearchPPTXWriter

    _pin_date()
    writer = ResearchPPTXWriter(output_dir=str(FIX_DIR))

    ev_path = writer.write_ev_bridge_deck(_ev_input(), filename="ref_ev_bridge")
    ifrs_inputs = SimpleNamespace(
        accounting_standard="IFRS", reported_ebitda=30_000_000_000.0,
        rou_depreciation=4_000_000_000.0, lease_interest=1_200_000_000.0,
        short_term_rent=300_000_000.0,
    )
    ifrs_out = SimpleNamespace(adjusted_ebitda=24_500_000_000.0)
    ifrs_path = writer.write_ifrs_bridge_deck(
        ifrs_inputs, ifrs_out, company="DemoCo", period="FY2025",
        revenue=100_000_000_000.0, filename="ref_ifrs_bridge",
    )

    import os
    for tag, path in (("ev", ev_path), ("ifrs", ifrs_path)):
        data = inspect_pptx(path, include_raw_xml=False)
        (SNAP_DIR / f"PPTX_deck_{tag}.json").write_text(
            json.dumps(data, indent=2, ensure_ascii=False), encoding="utf-8")
        print(f"  6.4 {tag}: slides={data['slideCount']} -> PPTX_deck_{tag}.json")
        os.remove(path)  # reference deck is throwaway; the JSON is the oracle


def main() -> None:
    oracle_inspect()
    oracle_edit()
    oracle_pure()
    oracle_decks()


if __name__ == "__main__":
    main()
