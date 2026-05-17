"""
Build TSLA 3-Statement Model PowerPoint (3 slides: IS, BS, CF)
Layout: 75% left = financial table | 25% right = analyst commentary
"""
from __future__ import annotations
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ── Palette ──────────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x0D, 0x1F, 0x3C)   # dark navy header
ACCENT      = RGBColor(0x00, 0x53, 0xA0)   # mid-blue subheader / rule
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x0F, 0x0F, 0x0F)
GRAY_LIGHT  = RGBColor(0xF2, 0xF4, 0xF7)   # alternating row bg
GRAY_MID    = RGBColor(0xD8, 0xDE, 0xE8)   # subtotal row bg
GRAY_TEXT   = RGBColor(0x55, 0x65, 0x7A)   # secondary text
GREEN_POS   = RGBColor(0x18, 0x7A, 0x36)
RED_NEG     = RGBColor(0xC0, 0x1C, 0x1C)
DIVIDER_CLR = RGBColor(0xBD, 0xC5, 0xD4)

FONT_BODY   = "Calibri"
FONT_HEADER = "Calibri"

W = Inches(13.33)
H = Inches(7.5)

# ── Layout constants ──────────────────────────────────────────────────────────
HDR_H    = Inches(0.85)
FTR_Y    = Inches(7.05)
FTR_H    = Inches(0.38)
TABLE_X  = Inches(0.28)
TABLE_Y  = Inches(0.95)
TABLE_W  = Inches(9.65)
TABLE_H  = Inches(5.95)
COMM_X   = Inches(10.05)
COMM_Y   = Inches(0.95)
COMM_W   = Inches(3.05)
COMM_H   = Inches(5.95)


# ── Helpers ───────────────────────────────────────────────────────────────────
def _rgb(r: RGBColor):
    return f"{r[0]:02X}{r[1]:02X}{r[2]:02X}"


def set_cell_bg(cell, color: RGBColor):
    """Set solid fill on a table cell via OOXML."""
    tc = cell._tc
    tcPr = tc.find(qn("a:tcPr"))
    if tcPr is None:
        tcPr = etree.SubElement(tc, qn("a:tcPr"))
    # remove existing fills
    for old in tcPr.findall(qn("a:solidFill")):
        tcPr.remove(old)
    solidFill = etree.SubElement(tcPr, qn("a:solidFill"))
    srgbClr   = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgbClr.set("val", _rgb(color))


def set_cell_text(cell, text: str, bold=False, italic=False,
                  font_size=7.5, color: RGBColor = BLACK,
                  align=PP_ALIGN.RIGHT):
    tf = cell.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = align
    # clear existing runs
    for run in p.runs:
        pass
    p.clear()
    run = p.add_run()
    run.text = text
    run.font.name = FONT_BODY
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def _fmt(v, pct=False, mn=False, dp=1, show_neg_red=False):
    """Format number for display."""
    if v is None:
        return "—"
    if pct:
        return f"{v*100:.{dp}f}%"
    if mn:
        return f"({abs(v):,.{dp}f})" if v < 0 else f"{v:,.{dp}f}"
    return f"{v:,.{dp}f}"


def add_slide(prs: Presentation) -> object:
    blank = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(blank)


def add_header(slide, title: str, subtitle: str = ""):
    """Dark navy header bar."""
    shp = slide.shapes.add_shape(1, 0, 0, W, HDR_H)
    shp.fill.solid()
    shp.fill.fore_color.rgb = NAVY
    shp.line.fill.background()

    tf = shp.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name  = FONT_HEADER
    run.font.size  = Pt(18)
    run.font.bold  = True
    run.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT

    shp.left  = 0
    shp.top   = 0
    shp.width = W
    shp.height = HDR_H

    # Reposition text with internal margin via ooxml
    txBody = tf._txBody
    bodyPr = txBody.find(qn("a:bodyPr"))
    bodyPr.set("lIns", str(int(Inches(0.3))))
    bodyPr.set("tIns", str(int(Inches(0.20))))

    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name  = FONT_HEADER
        r2.font.size  = Pt(9)
        r2.font.bold  = False
        r2.font.color.rgb = RGBColor(0xA8, 0xBC, 0xD8)


def add_footer(slide, source: str = "Source: SEC EDGAR, Company Filings | USD millions | Model: Virtual Analyst"):
    shp = slide.shapes.add_textbox(0, FTR_Y, W, FTR_H)
    tf  = shp.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]

    # thin top rule via line above textbox — use a line shape
    line = slide.shapes.add_shape(1, 0, FTR_Y, W, Inches(0.015))
    line.fill.solid()
    line.fill.fore_color.rgb = DIVIDER_CLR
    line.line.fill.background()

    run = p.add_run()
    run.text = source
    run.font.name  = FONT_BODY
    run.font.size  = Pt(6.5)
    run.font.color.rgb = GRAY_TEXT
    p.alignment = PP_ALIGN.LEFT

    txBody = tf._txBody
    bodyPr = txBody.find(qn("a:bodyPr"))
    bodyPr.set("lIns", str(int(Inches(0.28))))
    bodyPr.set("tIns", str(int(Inches(0.06))))


def add_vertical_divider(slide):
    """Thin vertical line between table and commentary."""
    x = COMM_X - Inches(0.08)
    line = slide.shapes.add_shape(1, x, TABLE_Y, Inches(0.012), TABLE_H)
    line.fill.solid()
    line.fill.fore_color.rgb = DIVIDER_CLR
    line.line.fill.background()


def add_commentary_box(slide, header: str, bullets: list[str]):
    """25% right panel — section header + bullet list."""
    # light background panel
    panel = slide.shapes.add_shape(1, COMM_X - Inches(0.05), TABLE_Y,
                                   COMM_W + Inches(0.05), TABLE_H)
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(0xF7, 0xF9, 0xFC)
    panel.line.color.rgb = DIVIDER_CLR
    panel.line.width = Pt(0.5)

    # header label
    hdr = slide.shapes.add_textbox(COMM_X, TABLE_Y + Inches(0.08),
                                   COMM_W, Inches(0.35))
    tf  = hdr.text_frame
    p   = tf.paragraphs[0]
    run = p.add_run()
    run.text = header.upper()
    run.font.name  = FONT_HEADER
    run.font.size  = Pt(7.5)
    run.font.bold  = True
    run.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.LEFT
    txBody = tf._txBody
    bodyPr = txBody.find(qn("a:bodyPr"))
    bodyPr.set("lIns", str(int(Inches(0.12))))

    # thin underline under header
    ul = slide.shapes.add_shape(1, COMM_X, TABLE_Y + Inches(0.43),
                                 COMM_W, Inches(0.012))
    ul.fill.solid()
    ul.fill.fore_color.rgb = ACCENT
    ul.line.fill.background()

    # bullets textbox
    tb = slide.shapes.add_textbox(COMM_X, TABLE_Y + Inches(0.50),
                                  COMM_W, TABLE_H - Inches(0.60))
    tf = tb.text_frame
    tf.word_wrap = True
    txBody = tf._txBody
    bodyPr = txBody.find(qn("a:bodyPr"))
    bodyPr.set("lIns", str(int(Inches(0.10))))
    bodyPr.set("rIns", str(int(Inches(0.08))))
    bodyPr.set("tIns", str(int(Inches(0.04))))

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT

        # custom bullet via tab stop
        pPr = p._p.get_or_add_pPr()
        pPr.set("indent", str(int(Inches(-0.10))))
        pPr.set("marL",   str(int(Inches(0.15))))

        run = p.add_run()
        run.text = f"•  {bullet}"
        run.font.name  = FONT_BODY
        run.font.size  = Pt(7.5)
        run.font.color.rgb = BLACK
        run.font.bold  = False

        # add spacing after each bullet
        spcAft = etree.SubElement(pPr, qn("a:spcAft"))
        spcPts = etree.SubElement(spcAft, qn("a:spcPts"))
        spcPts.set("val", "600")  # 6pt spacing


def add_table_to_slide(slide, rows: list[dict], col_headers: list[str],
                        col_widths_in: list[float]):
    """
    rows: list of dicts with keys:
      label, values (list matching col_headers),
      style: 'normal' | 'subtotal' | 'section' | 'margin'
      bold: bool
    """
    from pptx.util import Inches as I
    n_rows = len(rows) + 1  # +1 for header
    n_cols = 1 + len(col_headers)

    table_shape = slide.shapes.add_table(
        n_rows, n_cols, TABLE_X, TABLE_Y, TABLE_W, TABLE_H
    ).table

    # Column widths
    label_w = I(col_widths_in[0])
    data_w  = [(TABLE_W - label_w) // len(col_headers)] * len(col_headers)
    table_shape.columns[0].width = label_w
    for j, w in enumerate(data_w):
        table_shape.columns[j + 1].width = w

    # Header row
    header_cells = table_shape.rows[0].cells
    set_cell_bg(header_cells[0], NAVY)
    set_cell_text(header_cells[0], "USD millions", bold=True,
                  font_size=7.5, color=WHITE, align=PP_ALIGN.LEFT)
    for j, hdr in enumerate(col_headers):
        set_cell_bg(header_cells[j + 1], NAVY)
        is_proj = hdr.endswith("E")
        clr = RGBColor(0xA8, 0xBC, 0xD8) if is_proj else WHITE
        set_cell_text(header_cells[j + 1], hdr, bold=True,
                      font_size=7.5, color=clr, align=PP_ALIGN.CENTER)

    # Data rows
    for i, row in enumerate(rows):
        r_idx = i + 1
        cells = table_shape.rows[r_idx].cells
        style = row.get("style", "normal")
        bold  = row.get("bold", False)

        if style == "section":
            bg = NAVY
            lbl_clr = WHITE
            val_clr = WHITE
            bold = True
        elif style == "subtotal":
            bg = GRAY_MID
            lbl_clr = BLACK
            val_clr = BLACK
            bold = True
        elif style == "margin":
            bg = WHITE
            lbl_clr = GRAY_TEXT
            val_clr = GRAY_TEXT
            bold = False
        else:
            bg = GRAY_LIGHT if i % 2 == 0 else WHITE
            lbl_clr = BLACK
            val_clr = BLACK

        set_cell_bg(cells[0], bg)
        set_cell_text(cells[0], row["label"], bold=bold,
                      font_size=7.5 if style != "section" else 7,
                      color=lbl_clr, align=PP_ALIGN.LEFT)

        for j, val in enumerate(row.get("values", [])):
            set_cell_bg(cells[j + 1], bg)
            txt = val if isinstance(val, str) else _fmt(val, dp=0)
            # detect negatives for color hint in normal rows
            clr = val_clr
            if style == "normal" and isinstance(val, (int, float)) and val < 0:
                clr = RED_NEG
            set_cell_text(cells[j + 1], txt, bold=bold,
                          font_size=7.5, color=clr, align=PP_ALIGN.RIGHT)

    # Remove all cell borders (clean look)
    tbl = table_shape._tbl
    for tr in tbl.findall(qn("a:tr")):
        for tc in tr.findall(qn("a:tc")):
            tcPr = tc.find(qn("a:tcPr"))
            if tcPr is None:
                tcPr = etree.SubElement(tc, qn("a:tcPr"))
            for side in ("lnL", "lnR", "lnT", "lnB"):
                existing = tcPr.find(qn(f"a:{side}"))
                if existing is not None:
                    tcPr.remove(existing)
                ln = etree.SubElement(tcPr, qn(f"a:{side}"))
                etree.SubElement(ln, qn("a:noFill"))


# ─────────────────────────────────────────────────────────────────────────────
# Data definitions
# ─────────────────────────────────────────────────────────────────────────────
PERIODS = ["2023A", "2024A", "2025A", "2026E", "2027E", "2028E"]
PROJ_COLS = {"2026E", "2027E", "2028E"}

IS_ROWS = [
    {"label": "INCOME STATEMENT", "values": [""] * 6, "style": "section"},
    {"label": "Revenue",
     "values": [96773, 97690, 94827, 93887, 92956, 92034],
     "bold": True},
    {"label": "  Growth %",
     "values": ["—", "0.9%", "(2.9%)", "(1.0%)", "(1.0%)", "(1.0%)"],
     "style": "margin"},
    {"label": "Cost of Revenue",
     "values": [79113, 80240, 77733, 76944, 76181, 75426]},
    {"label": "Gross Profit",
     "values": [17660, 17450, 17094, 16943, 16775, 16608],
     "style": "subtotal"},
    {"label": "  Gross Margin %",
     "values": ["18.2%", "17.9%", "18.0%", "18.0%", "18.0%", "18.0%"],
     "style": "margin"},
    {"label": "R&D Expense",
     "values": [3969, 4540, 6411, 4854, 4806, 4758]},
    {"label": "SG&A Expense",
     "values": [4800, 5150, 5834, 5128, 5077, 5026]},
    {"label": "Other / Non-recurring",
     "values": [0, 583, 390, 390, 390, 390]},
    {"label": "EBIT (Operating Income)",
     "values": [8891, 7076, 4355, 2515, 2486, 2457],
     "style": "subtotal"},
    {"label": "  EBIT Margin %",
     "values": ["9.2%", "7.2%", "4.6%", "2.7%", "2.7%", "2.7%"],
     "style": "margin"},
    {"label": "D&A",
     "values": [3330, 4120, 5030, 4057, 4017, 3977]},
    {"label": "EBITDA",
     "values": [12221, 11196, 9385, 6572, 6502, 6434],
     "style": "subtotal"},
    {"label": "  EBITDA Margin %",
     "values": ["12.6%", "11.5%", "9.9%", "7.0%", "7.0%", "7.0%"],
     "style": "margin"},
    {"label": "Interest Expense",
     "values": [45, 21, 12, 230, 230, 230]},
    {"label": "Interest Income",
     "values": [1066, 1569, 1680, 330, 270, 205]},
    {"label": "EBT",
     "values": [9912, 8624, 6023, 2614, 2525, 2432],
     "style": "subtotal"},
    {"label": "Income Tax",
     "values": [-5001, 1837, 1423, 626, 604, 582]},
    {"label": "Net Income",
     "values": [14997, 7091, 3794, 1989, 1921, 1850],
     "style": "subtotal", "bold": True},
    {"label": "  Net Margin %",
     "values": ["15.5%", "7.3%", "4.0%", "2.1%", "2.1%", "2.0%"],
     "style": "margin"},
    {"label": "EPS – Diluted ($)",
     "values": ["—", "—", "—", "0.55", "0.53", "0.51"],
     "style": "margin"},
]

BS_ROWS = [
    {"label": "ASSETS", "values": [""] * 6, "style": "section"},
    {"label": "Cash & Equivalents",
     "values": [16398, 16139, 16513, 13500, 10238, 6956],
     "bold": True},
    {"label": "Accounts Receivable",
     "values": [3508, 4418, 4576, 4060, 4020, 3980]},
    {"label": "Inventory",
     "values": [13626, 12017, 12392, 12347, 12225, 12104]},
    {"label": "Total Current Assets",
     "values": [49616, 58360, 68642, 65069, 61644, 58201],
     "style": "subtotal"},
    {"label": "PP&E, net",
     "values": [29725, 35836, "n/a", 5269, 10485, 15650]},
    {"label": "Goodwill & Intangibles",
     "values": [431, 394, 381, 381, 381, 381]},
    {"label": "Total Assets",
     "values": [106618, 122070, 137806, 139501, 141293, 143014],
     "style": "subtotal", "bold": True},
    {"label": "LIABILITIES & EQUITY", "values": [""] * 6, "style": "section"},
    {"label": "Accounts Payable",
     "values": [14431, 12474, 13371, 13077, 12948, 12819]},
    {"label": "Total Current Liabilities",
     "values": [28748, 28821, 31714, 31420, 31291, 31162],
     "style": "subtotal"},
    {"label": "Long-Term Debt",
     "values": [2682, 5535, 6584, 6584, 6584, 6584]},
    {"label": "Total Liabilities",
     "values": [43009, 48390, 54941, 54647, 54518, 54389],
     "style": "subtotal"},
    {"label": "Retained Earnings",
     "values": [27882, 35209, 39003, 40931, 42791, 44580]},
    {"label": "Total Equity",
     "values": [63367, 73617, 82807, 84796, 86717, 88567],
     "style": "subtotal", "bold": True},
    {"label": "Total Liab + Equity",
     "values": [106618, 122070, 137806, 139501, 141293, 143014],
     "style": "subtotal"},
]

CF_ROWS = [
    {"label": "OPERATING ACTIVITIES", "values": [""] * 6, "style": "section"},
    {"label": "Net Income",
     "values": [14997, 7091, 3794, 1989, 1921, 1850]},
    {"label": "D&A Add-back",
     "values": [3330, 4120, 5030, 4057, 4017, 3977]},
    {"label": "Change in Working Capital",
     "values": [-5071, 3712, 5923, 267, 33, 33]},
    {"label": "Cash from Operations (CFO)",
     "values": [13256, 14923, 14747, 6313, 5971, 5859],
     "style": "subtotal", "bold": True},
    {"label": "INVESTING ACTIVITIES", "values": [""] * 6, "style": "section"},
    {"label": "Capital Expenditures",
     "values": [-8899, -11342, -8527, -9325, -9233, -9141]},
    {"label": "Net Purchases of Investments",
     "values": [-19112, -35955, -37109, 0, 0, 0]},
    {"label": "Other Investing",
     "values": [12427, 28510, 30158, 0, 0, 0]},
    {"label": "Cash from Investing (CFI)",
     "values": [-15584, -18787, -15478, -9325, -9233, -9141],
     "style": "subtotal"},
    {"label": "FINANCING ACTIVITIES", "values": [""] * 6, "style": "section"},
    {"label": "Other Financing (Debt ± / Issuances)",
     "values": [2589, 3853, 1139, 0, 0, 0]},
    {"label": "Cash from Financing (CFF)",
     "values": [2589, 3853, 1139, 0, 0, 0],
     "style": "subtotal"},
    {"label": "Net Change in Cash",
     "values": [265, -152, 579, -3013, -3262, -3282],
     "style": "subtotal"},
    {"label": "Ending Cash",
     "values": [16398, 16139, 16513, 13500, 10238, 6956],
     "bold": True},
    {"label": "Free Cash Flow (CFO − CapEx)",
     "values": [4357, 3581, 6220, -3012, -3262, -3282],
     "style": "subtotal", "bold": True},
]

# ── Commentary ────────────────────────────────────────────────────────────────
IS_BULLETS = [
    "Revenue declined 2.9% in FY2025 to $94.8B after near-flat FY2024; automotive demand softness reflects EV cycle normalization and pricing pressure",
    "Gross margin recovered to 18.0% (FY2025) from trough of 17.9% (FY2024), but remains ~220bps below FY2023 peak — sustained cost reduction critical",
    "R&D surged to 6.8% of revenue (FY2025 vs. 4.1% FY2023) — signals major pivot toward autonomy (FSD/Robotaxi) and energy, at the cost of near-term margins",
    "EBIT margin collapsed to 4.6% (FY2025) from 9.2% (FY2023); projected to compress further to ~2.7% through FY2028E as R&D burden persists",
    "Net income fell from $15.0B (FY2023) to $3.8B (FY2025); FY2023 included a $5.0B deferred tax asset benefit — normalized earnings trajectory is materially lower",
    "EBITDA margin at 9.9% (FY2025); D&A expanding ($3.3B→$5.0B) as Gigafactory assets ramp — high operating leverage, but fixed cost base growing",
]

BS_BULLETS = [
    "Total assets expanded to $137.8B (FY2025) from $106.6B (FY2023) — driven by current asset growth, reflecting large investment securities portfolio",
    "Cash & equivalents stable at $16.5B (FY2025); strong liquidity supports capex and R&D runway even as FCF turns negative in projections",
    "LT debt nearly tripled: $2.7B → $6.6B (FY2023–FY2025); debt-to-equity of ~8% remains very conservative; balance sheet not leveraged",
    "PP&E grew ~21% (FY2023–FY2024) to $35.8B; FY2025 data missing — likely reflects continued Gigafactory Texas/Nevada/Germany ramp",
    "Retained earnings grew consistently ($27.9B→$39.0B) despite margin compression, as cumulative profits compound on the equity base",
    "Total equity at $82.8B (FY2025) up from $63.4B in FY2023 — no material dilution; zero dividends paid, reinforcing full reinvestment posture",
]

CF_BULLETS = [
    "CFO remained robust at $14.7B (FY2025) despite NI of $3.8B — D&A add-back ($5.0B) and WC release ($5.9B) drive the divergence; quality of earnings high",
    "FY2023 CFO sanity flag: CFO ($13.3B) < NI ($15.0B) despite $3.3B D&A — unusually large WC drag; likely inventory build and deferred revenue timing",
    "CapEx: $8.5B (FY2025), ~9.0% of revenue — down from $11.3B in FY2024 as Gigafactory ramp cycle matures; expect uptick as new capacity is greenlit",
    "FCF improved sharply to $6.2B (FY2025) vs $3.6B (FY2024) on lower capex; however, projections show FCF turning negative (~$3B deficit) FY2026E+ as growth reinvestment accelerates",
    "Investment securities activity dominates CFI ($37.1B net purchases FY2025); systematic cash deployment into short-duration instruments — not M&A",
    "No dividends or buybacks modelled; Tesla retains 100% of cash for organic growth — consistent with founder-led high-growth posture",
]


# ─────────────────────────────────────────────────────────────────────────────
# Build deck
# ─────────────────────────────────────────────────────────────────────────────
def build():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slides_data = [
        ("IS", "Income Statement",  IS_ROWS,  IS_BULLETS,  "Key Insights"),
        ("BS", "Balance Sheet",     BS_ROWS,  BS_BULLETS,  "Key Insights"),
        ("CF", "Cash Flow Statement", CF_ROWS, CF_BULLETS, "Key Insights"),
    ]

    for tag, title, rows, bullets, comm_hdr in slides_data:
        slide = add_slide(prs)
        add_header(slide,
                   f"Tesla, Inc. — {title}",
                   "USD millions  |  FY2023A – FY2028E  |  Base Case  |  Source: SEC EDGAR, Company Filings")
        add_table_to_slide(slide, rows, PERIODS, col_widths_in=[3.2])
        add_commentary_box(slide, comm_hdr, bullets)
        add_vertical_divider(slide)
        add_footer(slide,
                   "Source: Tesla 10-K (SEC EDGAR) • Projected: Virtual Analyst 3-Statement Model (Base Case) • Values in USD millions • As of May 2026")

    out = Path(__file__).parent / "decks" / "TSLA_3Statement_Model.pptx"
    out.parent.mkdir(exist_ok=True)
    prs.save(str(out))
    print(f"Saved: {out}")
    return str(out)


if __name__ == "__main__":
    build()
