"""PPTX audit post-pass — attach source links to numbers in a generated deck.

Mirrors the Excel post-pass (audit_pipeline.annotate_workbook_with_links): scan
every text run (in text boxes and table cells), parse a number out of it, match
it against the same filing-provenance and market-citation indexes, and set the
run's hyperlink. Filing numbers -> finmodelaudit: page link; market numbers ->
provider https URL.

Because decks often display rounded figures ($172.7B), only numbers shown at
their stored precision will match — coverage is necessarily partial and honest.

Public API:
    annotate_pptx_with_links(pptx_path, *, cache_path) -> dict
"""
from __future__ import annotations

import json
import re
from pathlib import Path
from typing import Any, Optional

from .audit_pipeline import build_link_indexes, _label_tokens
from .audit_open import build_uri


_NUM_TOKEN = re.compile(r"\(?-?\d[\d.,   ]*\)?%?")


def _parse_number(text: str) -> Optional[float]:
    """Extract the first numeric value from a run's text, or None.

    Handles thousands separators (comma / space / NBSP), parenthesised and signed
    negatives, and trailing percent (-> value/100). Decimal point is '.'.
    """
    if not text:
        return None
    m = _NUM_TOKEN.search(text)
    if not m:
        return None
    tok = m.group(0)
    neg = tok.startswith("(") and tok.endswith(")")
    pct = tok.endswith("%")
    core = tok.strip("()%").strip()
    core = core.replace(",", "").replace(" ", "").replace(" ", "").replace(" ", "")
    if core in ("", "-", ".", "-."):
        return None
    try:
        val = float(core)
    except ValueError:
        return None
    if neg:
        val = -abs(val)
    if pct:
        val = val / 100.0
    return val


def _link_text_frame(tf, context: str, value_index, market_index, counts) -> None:
    ctx_tokens = _label_tokens(context)
    for para in tf.paragraphs:
        for run in para.runs:
            num = _parse_number(run.text)
            if num is None:
                continue
            fv = float(num)

            cands = value_index.get(fv)
            if cands:
                best = max(cands, key=lambda c: len(c["label_tokens"] & ctx_tokens))
                if len(cands) > 1 and not (best["label_tokens"] & ctx_tokens):
                    continue
                if not best["pdf"] or not Path(best["pdf"]).exists():
                    continue
                page_idx = None if best["low_confidence"] else best["page_index"]
                page_1based = (page_idx + 1) if page_idx is not None else None
                run.hyperlink.address = build_uri(best["pdf"], page_1based)
                if page_idx is not None:
                    counts["linked_page"] += 1
                else:
                    counts["linked_doc"] += 1
                continue

            mcands = market_index.get(round(fv, 6))
            if not mcands:
                continue
            mbest = max(mcands, key=lambda c: len(c["label_tokens"] & ctx_tokens))
            if len(mcands) > 1 and not (mbest["label_tokens"] & ctx_tokens):
                continue
            run.hyperlink.address = mbest["url"]
            counts["linked_market"] += 1


def annotate_pptx_with_links(
    pptx_path: str | Path,
    *,
    cache_path: Optional[str | Path] = None,
) -> dict[str, int]:
    """Attach source hyperlinks to numeric text runs in a .pptx deck.

    Returns {"linked_page", "linked_doc", "linked_market", "total"}.
    """
    counts = {"linked_page": 0, "linked_doc": 0, "linked_market": 0}
    if cache_path is None:
        return {**counts, "total": 0}

    from pptx import Presentation

    cache = json.loads(Path(cache_path).read_text(encoding="utf-8"))
    value_index, market_index = build_link_indexes(cache)
    if not value_index and not market_index:
        return {**counts, "total": 0}

    prs = Presentation(str(pptx_path))
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                tbl = shape.table
                n_rows = len(tbl.rows)
                n_cols = len(tbl.columns)
                for r in range(n_rows):
                    row_label = tbl.cell(r, 0).text  # first cell = line-item label
                    for c in range(n_cols):
                        cell = tbl.cell(r, c)
                        _link_text_frame(cell.text_frame, row_label,
                                         value_index, market_index, counts)
            elif shape.has_text_frame:
                tf = shape.text_frame
                _link_text_frame(tf, tf.text, value_index, market_index, counts)

    prs.save(str(pptx_path))
    counts["total"] = counts["linked_page"] + counts["linked_doc"] + counts["linked_market"]
    return counts


def annotate_pptx_with_sources(pptx_path, *, cache_path=None) -> dict:
    """Append a Sources & Assumptions block to the deck's first-slide speaker
    notes when the cache has a ledger. Returns {"notes_added": int}. Never
    raises on a normal empty/missing input."""
    out = {"notes_added": 0}
    if cache_path is None:
        return out
    cache = json.loads(Path(cache_path).read_text(encoding="utf-8"))
    if not (cache.get("__ledger__", {}) or {}).get("entries"):
        return out
    from src.sources_report import build_sources_report
    from pptx import Presentation
    report = build_sources_report(cache)
    prs = Presentation(str(pptx_path))
    slides = list(prs.slides)
    if not slides:
        return out
    notes_tf = slides[0].notes_slide.notes_text_frame
    existing = notes_tf.text
    notes_tf.text = (existing + "\n\n" + report) if existing else report
    prs.save(str(pptx_path))
    out["notes_added"] = 1
    return out
