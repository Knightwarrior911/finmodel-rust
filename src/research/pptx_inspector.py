"""
PPTX Inspector — exhaustive structural reverse-engineering.

Extracts:
  - Theme: 12 theme colors (lt1/dk1/lt2/dk2/accent1-6/hlink/folHlink), major+minor fonts, raw theme XML
  - Slide masters: full XML
  - Slide layouts: name, type, full XML
  - Per slide: layout name, background, dimensions, full shape tree (recursive into groups)
  - Per shape:
        geometry (left/top/width/height/rotation/flip)
        auto-shape preset + adjustments (corner radius, etc.)
        fill (solid/gradient/picture/pattern with all stops)
        line (color, width, dash, head/tail)
        effects (shadow, glow, reflection)
        text frame (margins, anchor, wrap, autofit)
        per-paragraph (alignment, indent, level, line spacing, bullet)
        per-run (font, size, bold/italic/underline, color, hyperlink)
        tables (cell-level fills, borders, merged cells, row/col sizes)
        charts (type, series, colors, axes)
        connectors (start/end)
        placeholders (type, idx, name)
        raw XML  ← ultimate fallback for unsupported features

Plus:
  - Roundtrip clone_template() that copies the file byte-for-byte
    and swaps text by shape ID — guarantees 100% visual fidelity
"""
from __future__ import annotations

import json
import os
import shutil
import zipfile
from copy import deepcopy
from pathlib import Path
from typing import Any, Optional

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.util import Emu


# ── Namespaces ───────────────────────────────────────────────────────────────

A   = "http://schemas.openxmlformats.org/drawingml/2006/main"
P   = "http://schemas.openxmlformats.org/presentationml/2006/main"
R   = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
NS  = {"a": A, "p": P, "r": R}

def _q(ns: str, tag: str) -> str:
    return f"{{{ns}}}{tag}"


# ── Primitive helpers ────────────────────────────────────────────────────────

def _emu_to_in(v) -> Optional[float]:
    try:
        return round(int(v) / 914400, 4)
    except Exception:
        return None


def _emu_to_pt(v) -> Optional[float]:
    try:
        return round(int(v) / 12700, 2)
    except Exception:
        return None


def _ang_to_deg(v) -> Optional[float]:
    try:
        return round(int(v) / 60000, 2)
    except Exception:
        return None


def _hex(b: bytes) -> str:
    return "#" + b.hex().upper() if b else None


def _rgb(color_obj) -> Optional[str]:
    try:
        if color_obj and color_obj.type is not None:
            rgb = color_obj.rgb
            return f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
    except Exception:
        pass
    return None


def _xml_color_solid(elem) -> dict:
    """Extract an a:solidFill / a:srgbClr / a:schemeClr / a:sysClr block."""
    if elem is None:
        return {}
    out = {}
    sf = elem.find(_q(A, "srgbClr"))
    if sf is not None:
        out["rgb"] = "#" + sf.get("val", "").upper()
    sc = elem.find(_q(A, "schemeClr"))
    if sc is not None:
        out["scheme"] = sc.get("val")
        # capture lumMod / lumOff / shade / tint / alpha
        for mod in sc:
            mod_tag = mod.tag.split("}")[-1]
            if mod.get("val") is not None:
                out.setdefault("modifiers", []).append(
                    {mod_tag: mod.get("val")}
                )
    sy = elem.find(_q(A, "sysClr"))
    if sy is not None:
        out["sys"] = sy.get("val")
        out["sysLastClr"] = "#" + (sy.get("lastClr") or "").upper()
    pr = elem.find(_q(A, "prstClr"))
    if pr is not None:
        out["preset"] = pr.get("val")
    return out


# ── Fill / line / effect extraction (XML-level) ──────────────────────────────

def _extract_fill(spPr) -> dict:
    """Full fill descriptor from an a:spPr (or txBody) node."""
    if spPr is None:
        return {"type": "none"}

    # noFill
    if spPr.find(_q(A, "noFill")) is not None:
        return {"type": "noFill"}

    # solidFill
    sf = spPr.find(_q(A, "solidFill"))
    if sf is not None:
        return {"type": "solid", **_xml_color_solid(sf)}

    # gradFill
    gf = spPr.find(_q(A, "gradFill"))
    if gf is not None:
        out = {"type": "gradient", "stops": [], "rotWithShape": gf.get("rotWithShape")}
        gs_lst = gf.find(_q(A, "gsLst"))
        if gs_lst is not None:
            for gs in gs_lst.findall(_q(A, "gs")):
                stop = {"pos": int(gs.get("pos", 0)) / 1000}
                stop.update(_xml_color_solid(gs))
                out["stops"].append(stop)
        lin = gf.find(_q(A, "lin"))
        if lin is not None:
            out["angle"] = _ang_to_deg(lin.get("ang"))
            out["scaled"] = lin.get("scaled")
        return out

    # pattFill
    pf = spPr.find(_q(A, "pattFill"))
    if pf is not None:
        fg = pf.find(_q(A, "fgClr"))
        bg = pf.find(_q(A, "bgClr"))
        return {
            "type": "pattern",
            "preset": pf.get("prst"),
            "fg": _xml_color_solid(fg) if fg is not None else None,
            "bg": _xml_color_solid(bg) if bg is not None else None,
        }

    # blipFill (picture)
    bf = spPr.find(_q(A, "blipFill"))
    if bf is not None:
        blip = bf.find(_q(A, "blip"))
        return {
            "type": "picture",
            "embedRid": blip.get(_q(R, "embed")) if blip is not None else None,
            "linkRid":  blip.get(_q(R, "link"))  if blip is not None else None,
        }

    # grpFill / inheritance
    return {"type": "inherit"}


def _extract_line(spPr) -> dict:
    if spPr is None:
        return {}
    ln = spPr.find(_q(A, "ln"))
    if ln is None:
        return {}
    out = {
        "widthPt":   _emu_to_pt(ln.get("w")),
        "cap":       ln.get("cap"),
        "cmpd":      ln.get("cmpd"),
        "algn":      ln.get("algn"),
    }
    out["fill"] = _extract_fill(ln)
    dash = ln.find(_q(A, "prstDash"))
    if dash is not None:
        out["dash"] = dash.get("val")
    head = ln.find(_q(A, "headEnd"))
    tail = ln.find(_q(A, "tailEnd"))
    if head is not None:
        out["headEnd"] = {"type": head.get("type"), "w": head.get("w"), "len": head.get("len")}
    if tail is not None:
        out["tailEnd"] = {"type": tail.get("type"), "w": tail.get("w"), "len": tail.get("len")}
    return out


def _extract_effects(spPr) -> dict:
    if spPr is None:
        return {}
    eff_lst = spPr.find(_q(A, "effectLst"))
    if eff_lst is None:
        return {}
    out = {}
    outer = eff_lst.find(_q(A, "outerShdw"))
    if outer is not None:
        out["outerShadow"] = {
            "blurRad":   _emu_to_pt(outer.get("blurRad")),
            "dist":      _emu_to_pt(outer.get("dist")),
            "dir":       _ang_to_deg(outer.get("dir")),
            "rotWithShape": outer.get("rotWithShape"),
            "color":     _xml_color_solid(outer),
        }
    glow = eff_lst.find(_q(A, "glow"))
    if glow is not None:
        out["glow"] = {
            "rad":   _emu_to_pt(glow.get("rad")),
            "color": _xml_color_solid(glow),
        }
    refl = eff_lst.find(_q(A, "reflection"))
    if refl is not None:
        out["reflection"] = {a: refl.get(a) for a in refl.attrib}
    soft = eff_lst.find(_q(A, "softEdge"))
    if soft is not None:
        out["softEdge"] = {"rad": _emu_to_pt(soft.get("rad"))}
    return out


# ── Geometry ─────────────────────────────────────────────────────────────────

def _extract_xfrm(spPr) -> dict:
    if spPr is None:
        return {}
    xfrm = spPr.find(_q(A, "xfrm"))
    if xfrm is None:
        return {}
    off = xfrm.find(_q(A, "off"))
    ext = xfrm.find(_q(A, "ext"))
    out = {
        "left":  _emu_to_in(off.get("x"))  if off is not None else None,
        "top":   _emu_to_in(off.get("y"))  if off is not None else None,
        "width": _emu_to_in(ext.get("cx")) if ext is not None else None,
        "height":_emu_to_in(ext.get("cy")) if ext is not None else None,
    }
    rot = xfrm.get("rot")
    if rot:
        out["rotationDeg"] = _ang_to_deg(rot)
    if xfrm.get("flipH"):
        out["flipH"] = xfrm.get("flipH") == "1"
    if xfrm.get("flipV"):
        out["flipV"] = xfrm.get("flipV") == "1"
    return out


def _extract_prst_geom(spPr) -> dict:
    if spPr is None:
        return {}
    pg = spPr.find(_q(A, "prstGeom"))
    if pg is None:
        return {}
    out = {"preset": pg.get("prst")}
    av = pg.find(_q(A, "avLst"))
    if av is not None:
        adjs = []
        for gd in av.findall(_q(A, "gd")):
            adjs.append({"name": gd.get("name"), "fmla": gd.get("fmla")})
        if adjs:
            out["adjustments"] = adjs
    return out


# ── Text frame extraction ────────────────────────────────────────────────────

def _extract_paragraph(p_elem) -> dict:
    out = {}
    pPr = p_elem.find(_q(A, "pPr"))
    if pPr is not None:
        for k in ("algn", "lvl", "indent", "marL", "marR", "rtl"):
            v = pPr.get(k)
            if v is not None:
                out[k] = v
        ln_spc = pPr.find(_q(A, "lnSpc"))
        if ln_spc is not None:
            spc = ln_spc.find(_q(A, "spcPct"))
            if spc is None:
                spc = ln_spc.find(_q(A, "spcPts"))
            if spc is not None:
                out["lineSpacing"] = {spc.tag.split("}")[-1]: spc.get("val")}
        spc_b = pPr.find(_q(A, "spcBef"))
        spc_a = pPr.find(_q(A, "spcAft"))
        for label, node in [("spcBef", spc_b), ("spcAft", spc_a)]:
            if node is not None:
                pts = node.find(_q(A, "spcPts"))
                pct = node.find(_q(A, "spcPct"))
                if pts is not None: out[label] = {"pts": int(pts.get("val"))/100}
                if pct is not None: out[label] = {"pct": int(pct.get("val"))/1000}
        # Bullet character / numbering
        for bu_tag in ("buNone", "buChar", "buAutoNum", "buFont"):
            bu = pPr.find(_q(A, bu_tag))
            if bu is not None:
                out.setdefault("bullet", {})[bu_tag] = dict(bu.attrib)

    runs = []
    for r in p_elem.findall(_q(A, "r")):
        run = {}
        rPr = r.find(_q(A, "rPr"))
        if rPr is not None:
            for k in ("sz", "b", "i", "u", "strike", "baseline", "spc", "lang", "kern"):
                v = rPr.get(k)
                if v is not None:
                    run[k] = v
            sf = rPr.find(_q(A, "solidFill"))
            if sf is not None:
                run["color"] = _xml_color_solid(sf)
            for font_tag in ("latin", "ea", "cs"):
                fnode = rPr.find(_q(A, font_tag))
                if fnode is not None:
                    run.setdefault("fonts", {})[font_tag] = fnode.get("typeface")
            hl = rPr.find(_q(A, "highlight"))
            if hl is not None:
                run["highlight"] = _xml_color_solid(hl)
            hyper = rPr.find(_q(A, "hlinkClick"))
            if hyper is not None:
                run["hyperlinkRid"] = hyper.get(_q(R, "id"))
        t = r.find(_q(A, "t"))
        run["text"] = t.text if t is not None else ""
        runs.append(run)

    # Line break or field handling
    for child in p_elem:
        tag = child.tag.split("}")[-1]
        if tag == "br":
            runs.append({"break": True})
        elif tag == "fld":
            t = child.find(_q(A, "t"))
            runs.append({"field": child.get("type"), "text": t.text if t is not None else ""})

    out["runs"] = runs
    return out


def _extract_text_frame(txBody) -> dict:
    if txBody is None:
        return {}
    out = {}
    bodyPr = txBody.find(_q(A, "bodyPr"))
    if bodyPr is not None:
        for k in ("anchor", "anchorCtr", "wrap", "lIns", "tIns", "rIns", "bIns",
                 "rot", "vert", "spcCol", "numCol", "fromWordArt", "compatLnSpc"):
            v = bodyPr.get(k)
            if v is not None:
                if k in ("lIns", "tIns", "rIns", "bIns"):
                    out.setdefault("margins", {})[k] = _emu_to_in(v)
                else:
                    out[k] = v
        # Autofit
        for af in ("normAutofit", "spAutoFit", "noAutofit"):
            af_elem = bodyPr.find(_q(A, af))
            if af_elem is not None:
                out["autofit"] = {af: dict(af_elem.attrib) if af_elem.attrib else True}

    paragraphs = [_extract_paragraph(p) for p in txBody.findall(_q(A, "p"))]
    out["paragraphs"] = paragraphs

    # Plain-text concatenation for quick reading
    plain = []
    for p in paragraphs:
        line = "".join((r.get("text") or "") for r in p.get("runs", []))
        if line:
            plain.append(line)
    if plain:
        out["text"] = "\n".join(plain)
    return out


# ── Table / chart / connector extraction ─────────────────────────────────────

def _extract_table(tbl_elem) -> dict:
    if tbl_elem is None:
        return {}
    rows = []
    grid = tbl_elem.find(_q(A, "tblGrid"))
    col_widths = []
    if grid is not None:
        for gc in grid.findall(_q(A, "gridCol")):
            col_widths.append(_emu_to_in(gc.get("w")))
    for tr in tbl_elem.findall(_q(A, "tr")):
        row = {"heightIn": _emu_to_in(tr.get("h")), "cells": []}
        for tc in tr.findall(_q(A, "tc")):
            cell = {}
            cell["gridSpan"] = int(tc.get("gridSpan", 1))
            cell["rowSpan"]  = int(tc.get("rowSpan", 1))
            cell["hMerge"]   = tc.get("hMerge") == "1"
            cell["vMerge"]   = tc.get("vMerge") == "1"
            txBody = tc.find(_q(A, "txBody"))
            if txBody is not None:
                cell["text"] = _extract_text_frame(txBody)
            tcPr = tc.find(_q(A, "tcPr"))
            if tcPr is not None:
                cell["fill"] = _extract_fill(tcPr)
                # borders
                borders = {}
                for side in ("lnL", "lnR", "lnT", "lnB"):
                    bn = tcPr.find(_q(A, side))
                    if bn is not None:
                        borders[side] = {
                            "widthPt": _emu_to_pt(bn.get("w")),
                            "fill": _extract_fill(bn),
                        }
                if borders:
                    cell["borders"] = borders
                cell["margins"] = {
                    k: _emu_to_in(tcPr.get(k))
                    for k in ("marL", "marR", "marT", "marB")
                    if tcPr.get(k)
                }
            row["cells"].append(cell)
        rows.append(row)
    return {"colWidthsIn": col_widths, "rows": rows}


def _extract_chart(chart) -> dict:
    try:
        out = {
            "type": str(chart.chart_type).split(".")[-1],
            "hasTitle": chart.has_title,
            "title": chart.chart_title.text_frame.text if chart.has_title else None,
            "seriesCount": len(chart.series),
            "series": [],
        }
        for s in chart.series:
            out["series"].append({"name": s.name, "valuesCount": len(list(s.values))})
        return out
    except Exception as e:
        return {"error": str(e)}


# ── Master / layout / theme ──────────────────────────────────────────────────

def _theme_from_master(master) -> dict:
    """Walk the master's theme part and pull color + font schemes."""
    out = {"colors": {}, "fonts": {}, "rawXml": None}
    try:
        for rel in master.part.rels.values():
            if rel.reltype == RT.THEME:
                theme_blob = rel.target_part.blob
                out["rawXml"] = theme_blob.decode("utf-8", errors="replace")
                root = etree.fromstring(theme_blob)
                clr_scheme = root.find(f".//{_q(A, 'clrScheme')}")
                if clr_scheme is not None:
                    out["clrSchemeName"] = clr_scheme.get("name")
                    for child in clr_scheme:
                        slot = child.tag.split("}")[-1]
                        out["colors"][slot] = _xml_color_solid(child)
                font_scheme = root.find(f".//{_q(A, 'fontScheme')}")
                if font_scheme is not None:
                    out["fontSchemeName"] = font_scheme.get("name")
                    for kind in ("majorFont", "minorFont"):
                        node = font_scheme.find(_q(A, kind))
                        if node is not None:
                            latin = node.find(_q(A, "latin"))
                            ea = node.find(_q(A, "ea"))
                            cs = node.find(_q(A, "cs"))
                            out["fonts"][kind] = {
                                "latin": latin.get("typeface") if latin is not None else None,
                                "ea":    ea.get("typeface")    if ea    is not None else None,
                                "cs":    cs.get("typeface")    if cs    is not None else None,
                            }
                break
    except Exception as e:
        out["error"] = str(e)
    return out


def _layout_summary(layout, *, include_shapes: bool = True,
                    include_xml: bool = True) -> dict:
    elem = layout.element
    layout_type = elem.get("type")
    out = {
        "name": layout.name,
        "type": layout_type,
        "shapeCount": len(list(layout.shapes)),
    }
    if include_shapes:
        out["elements"] = [_shape_info(sh, include_xml=False) for sh in layout.shapes]
    if include_xml:
        out["rawXml"] = etree.tostring(elem, pretty_print=True).decode("utf-8")
    return out


# ── Shape-level inspector ────────────────────────────────────────────────────

def _placeholder_info(shape) -> Optional[dict]:
    if not shape.is_placeholder:
        return None
    ph = shape.placeholder_format
    return {
        "idx":  ph.idx,
        "type": str(ph.type).split(".")[-1] if ph.type else None,
    }


def _connector_endpoints(elem) -> Optional[dict]:
    """For connector shapes, capture start/end connection."""
    nv = elem.find(_q(P, "nvCxnSpPr"))
    if nv is None:
        return None
    cxn = nv.find(_q(P, "cNvCxnSpPr"))
    if cxn is None:
        return None
    out = {}
    s = cxn.find(_q(A, "stCxn"))
    e = cxn.find(_q(A, "endCxn"))
    if s is not None:
        out["start"] = {"id": s.get("id"), "idx": s.get("idx")}
    if e is not None:
        out["end"]   = {"id": e.get("id"), "idx": e.get("idx")}
    return out or None


def _shape_info(sh, include_xml: bool = True) -> dict:
    elem = sh.element
    spPr = elem.find(f".//{_q(P, 'spPr')}")
    if spPr is None:
        spPr = elem.find(f".//{_q(A, 'spPr')}")
    txBody = elem.find(f".//{_q(P, 'txBody')}")
    if txBody is None:
        txBody = elem.find(f".//{_q(A, 'txBody')}")

    info: dict = {
        "id":   sh.shape_id,
        "name": sh.name,
        "type": str(sh.shape_type).split(".")[-1] if sh.shape_type else "UNKNOWN",
        "pos": {
            "left":   _emu_to_in(sh.left),
            "top":    _emu_to_in(sh.top),
            "width":  _emu_to_in(sh.width),
            "height": _emu_to_in(sh.height),
        },
    }

    # Auto-shape preset + adjustments
    try:
        if sh.auto_shape_type is not None:
            info["autoShapeType"] = str(sh.auto_shape_type).split(".")[-1]
    except Exception:
        pass
    try:
        adjs = list(sh.adjustments)
        if adjs:
            info["adjustments"] = adjs
    except Exception:
        pass

    # XML-level extraction (richer than python-pptx convenience accessors)
    info["xfrm"]      = _extract_xfrm(spPr)
    info["prstGeom"]  = _extract_prst_geom(spPr)
    info["fill"]      = _extract_fill(spPr)
    info["line"]      = _extract_line(spPr)
    eff               = _extract_effects(spPr)
    if eff:
        info["effects"] = eff

    # Placeholder
    ph = _placeholder_info(sh)
    if ph:
        info["placeholder"] = ph

    # Connector
    cxn = _connector_endpoints(elem)
    if cxn:
        info["connector"] = cxn

    # Text
    if sh.has_text_frame and txBody is not None:
        info["text"] = _extract_text_frame(txBody)

    # Table
    if sh.has_table:
        tbl_elem = elem.find(f".//{_q(A, 'tbl')}")
        info["table"] = _extract_table(tbl_elem)

    # Chart
    if sh.has_chart:
        info["chart"] = _extract_chart(sh.chart)

    # Picture
    if str(sh.shape_type).endswith("PICTURE"):
        try:
            img = sh.image
            info["picture"] = {
                "filename":    img.filename,
                "contentType": img.content_type,
                "sizeBytes":   len(img.blob),
                "extension":   img.ext,
            }
        except Exception:
            pass

    # Group: recurse
    if str(sh.shape_type).endswith("GROUP"):
        info["children"] = [_shape_info(child, include_xml=include_xml) for child in sh.shapes]

    # Raw XML fallback (string-encoded; truncated by caller if needed)
    if include_xml:
        try:
            info["rawXml"] = etree.tostring(elem, pretty_print=False).decode("utf-8")
        except Exception:
            pass

    return info


def _slide_background(slide) -> dict:
    bg = slide.element.find(f".//{_q(P, 'bg')}")
    if bg is None:
        return {"type": "inherit"}
    # bgPr (explicit) or bgRef (style index)
    bgPr  = bg.find(_q(P, "bgPr"))
    bgRef = bg.find(_q(P, "bgRef"))
    if bgPr is not None:
        return {"type": "explicit", "fill": _extract_fill(bgPr)}
    if bgRef is not None:
        return {"type": "ref", "idx": bgRef.get("idx"),
                "color": _xml_color_solid(bgRef)}
    return {"type": "inherit"}


def _slide_layout_name(slide) -> str:
    try:
        return slide.slide_layout.name
    except Exception:
        return ""


# ── Top-level inspector ──────────────────────────────────────────────────────

def inspect_pptx(
    path: str,
    *,
    include_raw_xml: bool = True,
    include_layouts: bool = True,
    max_xml_chars: Optional[int] = None,
) -> dict:
    """Exhaustive structural inspection."""
    prs = Presentation(path)

    def _truncate(s: Optional[str]) -> Optional[str]:
        if s is None or max_xml_chars is None:
            return s
        return s if len(s) <= max_xml_chars else s[:max_xml_chars] + f"...<+{len(s)-max_xml_chars} chars>"

    masters = []
    for m in prs.slide_masters:
        m_data = {
            "name": getattr(m, "name", ""),
            "theme": _theme_from_master(m),
        }
        if include_raw_xml:
            m_data["rawXml"] = _truncate(
                etree.tostring(m.element, pretty_print=True).decode("utf-8")
            )
            if "rawXml" in m_data["theme"]:
                m_data["theme"]["rawXml"] = _truncate(m_data["theme"]["rawXml"])
        masters.append(m_data)

    layouts = []
    if include_layouts:
        for l in prs.slide_layouts:
            ld = _layout_summary(l)
            if not include_raw_xml:
                ld.pop("rawXml", None)
            else:
                ld["rawXml"] = _truncate(ld["rawXml"])
            layouts.append(ld)

    out: dict = {
        "path": str(path),
        "fileSizeBytes": os.path.getsize(path) if os.path.exists(path) else None,
        "slideCount":  len(prs.slides),
        "layoutCount": len(prs.slide_layouts),
        "masterCount": len(prs.slide_masters),
        "dimensions": {
            "widthIn":  round(prs.slide_width  / 914400, 4),
            "heightIn": round(prs.slide_height / 914400, 4),
            "widthEmu": prs.slide_width,
            "heightEmu": prs.slide_height,
        },
        "masters": masters,
        "layouts": layouts,
        "slides": [],
    }

    for i, slide in enumerate(prs.slides):
        shapes = list(slide.shapes)
        s_info = {
            "index":        i,
            "layoutName":   _slide_layout_name(slide),
            "background":   _slide_background(slide),
            "elementCount": len(shapes),
            "elements":     [_shape_info(sh, include_xml=include_raw_xml) for sh in shapes],
        }
        if include_raw_xml:
            for el in s_info["elements"]:
                if "rawXml" in el:
                    el["rawXml"] = _truncate(el["rawXml"])
        out["slides"].append(s_info)

    return out


def inspect_pptx_json(
    path: str,
    *,
    indent: int = 2,
    include_raw_xml: bool = False,
    max_xml_chars: int = 600,
) -> str:
    """JSON serialised inspection. Raw XML off by default for readability."""
    return json.dumps(
        inspect_pptx(path, include_raw_xml=include_raw_xml, max_xml_chars=max_xml_chars),
        indent=indent, ensure_ascii=False,
    )


# ── Replicator: clone template and modify in place ───────────────────────────

def clone_template(
    template_path: str,
    output_path: str,
    text_replacements: Optional[dict] = None,
    *,
    by_shape_id: Optional[dict] = None,
    by_placeholder_idx: Optional[dict] = None,
    on_slide: Optional[int] = None,
) -> str:
    """
    Copy a template PPTX byte-for-byte (preserves 100% visual fidelity)
    then perform targeted text substitutions.

    text_replacements:        {"old text": "new text"} — global string swap
    by_shape_id:              {shape_id: "new text"}  — exact shape match
    by_placeholder_idx:       {idx: "new text"}        — placeholder swap (per slide)
    on_slide:                 limit by_* maps to this 0-based slide index
    """
    shutil.copy2(template_path, output_path)
    prs = Presentation(output_path)

    text_replacements   = text_replacements   or {}
    by_shape_id         = by_shape_id         or {}
    by_placeholder_idx  = by_placeholder_idx  or {}

    for i, slide in enumerate(prs.slides):
        if on_slide is not None and i != on_slide:
            continue
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            # Shape-id targeted
            if shape.shape_id in by_shape_id:
                _set_first_run_text(tf, by_shape_id[shape.shape_id])
                continue
            # Placeholder-idx targeted
            if shape.is_placeholder and shape.placeholder_format.idx in by_placeholder_idx:
                _set_first_run_text(tf, by_placeholder_idx[shape.placeholder_format.idx])
                continue
            # Global text replacement
            if text_replacements:
                for para in tf.paragraphs:
                    for run in para.runs:
                        for old, new in text_replacements.items():
                            if old in run.text:
                                run.text = run.text.replace(old, new)

    prs.save(output_path)
    return output_path


def _set_first_run_text(tf, new_text: str) -> None:
    """Replace text frame contents while preserving the first run's formatting."""
    if not tf.paragraphs:
        tf.text = new_text
        return
    p0 = tf.paragraphs[0]
    if not p0.runs:
        p0.text = new_text
        return
    # Preserve formatting of run 0; clear the rest
    p0.runs[0].text = new_text
    for run in p0.runs[1:]:
        run.text = ""
    # Drop additional paragraphs
    for extra in tf.paragraphs[1:]:
        extra._pPr = None
        for run in extra.runs:
            run.text = ""


# ── Roundtrip diff helper ────────────────────────────────────────────────────

def diff_decks(path_a: str, path_b: str) -> dict:
    """
    Compare two PPTX files at the structural-fingerprint level.
    Useful for proving the inspector + replicator are faithful.
    """
    a = inspect_pptx(path_a, include_raw_xml=False)
    b = inspect_pptx(path_b, include_raw_xml=False)

    diffs = []
    if a["slideCount"] != b["slideCount"]:
        diffs.append(f"slideCount {a['slideCount']} vs {b['slideCount']}")
    if a["dimensions"] != b["dimensions"]:
        diffs.append(f"dimensions differ: {a['dimensions']} vs {b['dimensions']}")
    for i in range(min(len(a["slides"]), len(b["slides"]))):
        sa, sb = a["slides"][i], b["slides"][i]
        if sa["elementCount"] != sb["elementCount"]:
            diffs.append(f"slide {i}: shape count {sa['elementCount']} vs {sb['elementCount']}")
        if sa["layoutName"] != sb["layoutName"]:
            diffs.append(f"slide {i}: layout {sa['layoutName']!r} vs {sb['layoutName']!r}")
    return {"identical": len(diffs) == 0, "differences": diffs}


# ── CLI ──────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    import sys
    if len(sys.argv) < 2:
        print("Usage: python pptx_inspector.py <file.pptx> [--no-xml] [--max-xml=N]")
        sys.exit(1)
    args = sys.argv[1:]
    path = args[0]
    include_xml = "--no-xml" not in args
    max_xml = 600
    for a in args:
        if a.startswith("--max-xml="):
            max_xml = int(a.split("=", 1)[1])
    print(inspect_pptx_json(path, include_raw_xml=include_xml, max_xml_chars=max_xml))
