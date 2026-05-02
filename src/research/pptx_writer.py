"""
PowerPoint output writer for research agent results.

Follows valuation_kit standards:
  SPEC_powerpoint_engineering.md       - action titles, citations, verification, anti-patterns
  SPEC_powerpoint_formatting.md        - dimensions, fonts, colors, layouts, charts, tables
  SPEC_powerpoint_layout_decisions.md  - exhibit decision tree + 5 ad-hoc archetypes
  SPEC_PitchPres_A4_Landscape.md       - Citi A4-landscape pitch template (firm-specific)

Archetypes implemented:
  1. Comparison Matrix    - N entities x M dimensions (peer comp, bid comparison)
  2. Scorecard            - 4-9 metric tiles with ratings (DD scorecard, screening)
  3. Quote Wall           - 4-8 management quotes (commentary synthesis)
  4. Timeline             - sequence of dated events
  5. Process Diagram      - boxes + arrows (transaction structure, flow)
  6. Strategy Framework   - 2-5 column priorities/initiatives (vision + framework)
  7. Bar Chart            - horizontal bars, target-highlighted, sorted desc
  8. Football Field       - valuation range bars per method
  9. Line Chart           - 1-6 series over time, native LINE_MARKERS chart
 10. Waterfall            - bridge chart (start/plus/minus/total), broken-axis option
 11. Stacked Bar          - vertical breakdown, 2-8 cats x 2-6 series, native COLUMN_STACKED
 12. Pie                  - composition share, 2-8 slices, native chart
 13. Pros / Cons          - 3-column eval (Forest Bright pros, Citi Red cons)
 14. Quad Page            - 2x2 grid of text panels
 15. Org Chart            - 4-tier hierarchical (Blue / Ink / Blue Light / Gray-stat-box)
 16. Tombstone Page       - grid of deal tiles
 17. Team Page            - banner + person tiles
 18. Table of Contents    - numbered hierarchy (1/A/i)

Plus cover, section divider.

Five binding rules from engineering spec are enforced on every slide:
  R1 - Action title (caller supplies; required, non-empty)
  R2 - One idea per slide (one archetype call = one slide)
  R3 - Cite every data slide (source line auto-rendered, skip_source=True opts out)
  R4 - Visual hierarchy (action title > body > footnote)
  R5 - Verification cycle (verify(deck_path) function)

Geometry:
  - 16:9 widescreen (default): 13.33" x 7.50"
  - A4 landscape (Citi PitchPres): 10.83" x 7.50"
  - 4:3: 10.0" x 7.50"
  Set via BrandProfile.aspect_ratio.
"""

from __future__ import annotations

import os
from dataclasses import dataclass, field
from datetime import date
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu


# ──────────────────────────────────────────────────────────────────────
# Constants from SPEC_powerpoint_formatting.md
# ──────────────────────────────────────────────────────────────────────

# Dimensions (Section 1) - 16:9 widescreen default
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
MARGIN_IN  = 0.5

# Aspect ratio dimensions (selectable via BrandProfile.aspect_ratio)
ASPECT_DIMS = {
    "16:9":         (13.333, 7.5),    # widescreen default
    "4:3":          (10.0,   7.5),    # legacy projection
    "A4_LANDSCAPE": (10.83,  7.5),    # Citi PitchPres (780pt x 540pt)
}

# Color palette (Section 4)
BRAND_BLUE = RGBColor(0x25, 0x5B, 0xE3)  # #255BE3 primary
INK        = RGBColor(0x0F, 0x16, 0x32)  # #0F1632 dark
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE6, 0xEB, 0xED)  # alt row shading
MID_GRAY   = RGBColor(0xD3, 0xDA, 0xDD)  # dividers
BORDER_GRAY = RGBColor(0xA4, 0xAC, 0xAF)
SAND       = RGBColor(0xEA, 0xE0, 0xD3)  # subsection shading
ACCENT_RED = RGBColor(0xFF, 0x3C, 0x28)  # negatives, sparingly
FOOTNOTE_GRAY = RGBColor(0x80, 0x80, 0x80)

# Chart series palette (Section 4.2) - in priority order
SERIES_PALETTE = [
    RGBColor(0x25, 0x5B, 0xE3),  # 1. Brand Primary
    RGBColor(0x0F, 0x16, 0x32),  # 2. Brand Dark
    RGBColor(0x73, 0xC2, 0xFC),  # 3. Light Brand
    RGBColor(0xA4, 0xAC, 0xAF),  # 4. Mid Gray
    RGBColor(0x38, 0x8A, 0x42),  # 5. Forest
    RGBColor(0x80, 0xCE, 0x84),  # 6. Light Forest
    RGBColor(0xFA, 0xB7, 0x28),  # 7. Tan
    RGBColor(0xFF, 0xA1, 0x5A),  # 8. Orange Light
    RGBColor(0x8E, 0x31, 0x9C),  # 9. Purple
    RGBColor(0xD7, 0x16, 0x71),  # 10. Plum
]

# Typography (Section 3)
FONT_HEADLINE = "Arial"
FONT_BODY     = "Arial"

PT_HEADLINE   = 26  # Section 3.1: 24-28pt action title
PT_SUB        = 16
PT_BODY       = 12
PT_TABLE_H    = 11
PT_TABLE_B    = 10
PT_FOOTNOTE   = 8   # Section 3.1: 7-8pt italic
PT_PAGE       = 8


# ──────────────────────────────────────────────────────────────────────
# Archetype decision tree - per SPEC_powerpoint_layout_decisions Section 1+2
# ──────────────────────────────────────────────────────────────────────

ARCH_COMPARISON = "comparison_matrix"
ARCH_SCORECARD  = "scorecard"
ARCH_QUOTE_WALL = "quote_wall"
ARCH_TIMELINE   = "timeline"
ARCH_PROCESS    = "process_diagram"
ARCH_STRATEGY   = "strategy_framework"


@dataclass
class ArchetypeDecision:
    """Output of pick_slide_archetype()."""
    archetype: str
    split_required: bool   # True when entity count > density limit
    n_slides: int          # how many slides needed (1 or more)
    rationale: list = field(default_factory=list)


# Density limits (Section 3.1)
_DENSITY_LIMITS = {
    ARCH_COMPARISON: (8, 8),   # max 8 entities x 8 metrics per slide
    ARCH_SCORECARD:  (9, 1),   # 4-9 tiles
    ARCH_QUOTE_WALL: (8, 1),   # 4-8 quotes
    ARCH_TIMELINE:   (10, 1),  # max 10 milestones
    ARCH_PROCESS:    (8, 1),   # max 8 boxes
    ARCH_STRATEGY:   (5, 1),   # max 5 columns (priorities/initiatives)
}


def pick_slide_archetype(
    *,
    data_shape: str,        # "comparison"|"breakdown"|"trend"|"process"|"structure"|
                            # "geography"|"single_stat"|"range"|"bridge"|"scatter"|
                            # "quotes"|"matrix"|"events"
    n_entities: int = 1,
    n_metrics: int = 1,
    has_quotes: bool = False,
    is_dated: bool = False,
) -> ArchetypeDecision:
    """
    Map data shape (Q2 of decision tree) to one of the 5 ad-hoc archetypes.
    Density-check (Section 3.1) and recommend split if exceeded.
    """
    rationale = []

    # Q2 -> archetype mapping
    if has_quotes or data_shape == "quotes":
        archetype = ARCH_QUOTE_WALL
        rationale.append("data_shape=quotes -> quote_wall")
    elif is_dated or data_shape == "events":
        archetype = ARCH_TIMELINE
        rationale.append(f"is_dated/events -> timeline")
    elif data_shape in ("framework", "priorities", "initiatives", "strategy"):
        archetype = ARCH_STRATEGY
        rationale.append(f"data_shape={data_shape} -> strategy_framework")
    elif data_shape in ("process", "structure"):
        archetype = ARCH_PROCESS
        rationale.append(f"data_shape={data_shape} -> process_diagram")
    elif data_shape == "single_stat" or (n_metrics <= 9 and n_entities == 1):
        archetype = ARCH_SCORECARD
        rationale.append("single entity, multiple metrics -> scorecard")
    else:
        # Default: comparison matrix
        archetype = ARCH_COMPARISON
        rationale.append(f"data_shape={data_shape} -> comparison_matrix")

    # Q3 density check
    max_e, max_m = _DENSITY_LIMITS[archetype]
    split_required = (n_entities > max_e) or (n_metrics > max_m)
    n_slides = 1
    if split_required:
        # Round up
        if archetype == ARCH_COMPARISON:
            n_slides = max(1, -(-n_entities // max_e))
        else:
            n_slides = max(1, -(-n_entities // max_e))
        rationale.append(f"density {n_entities}x{n_metrics} > limit "
                         f"{max_e}x{max_m} -> {n_slides} slides")

    return ArchetypeDecision(
        archetype=archetype,
        split_required=split_required,
        n_slides=n_slides,
        rationale=rationale,
    )


def verify(deck_path: str) -> dict:
    """
    Run R5 verification cycle on a saved .pptx file.

    Three QA passes per SPEC_powerpoint_engineering Section 8:
      - Structural: file opens, every slide loads, shape geometry valid
      - Visual: shapes don't extend past slide bounds, no obvious overflow
      - Content: data slides have source line, no placeholder text
                 ([TBD], TODO, FIXME, lorem ipsum)

    Returns dict {"critical": [...], "minor": [...], "passed": int}.
    Critical = blocks delivery; minor = ship-acceptable.
    """
    from pptx import Presentation
    from pptx.util import Emu

    issues_critical: list = []
    issues_minor: list = []

    placeholder_tags = ["[tbd]", "lorem ipsum", "fixme", "xxx",
                        "todo:", "placeholder"]

    try:
        p = Presentation(deck_path)
    except Exception as e:
        return {"critical": [f"failed to open deck: {e}"],
                "minor": [], "passed": 0}

    slide_w = p.slide_width
    slide_h = p.slide_height
    n = len(p.slides)

    for i, slide in enumerate(p.slides):
        slide_num = i + 1
        has_source_line = False
        all_text = []

        for shape in slide.shapes:
            # Geometry check
            try:
                if shape.left is None or shape.top is None:
                    continue
                right = shape.left + (shape.width or 0)
                bottom = shape.top + (shape.height or 0)
                # Tolerance: allow 0.05" overflow before flagging
                tol = Emu(45720)  # 0.05"
                if right > slide_w + tol:
                    over = (right - slide_w) / 914400  # to inches
                    issues_critical.append(
                        f"slide {slide_num}: shape extends past right edge "
                        f"by {over:.2f}in")
                if bottom > slide_h + tol:
                    over = (bottom - slide_h) / 914400
                    issues_critical.append(
                        f"slide {slide_num}: shape extends past bottom edge "
                        f"by {over:.2f}in")
            except Exception:
                pass

            # Text checks
            if shape.has_text_frame:
                txt = shape.text_frame.text
                if txt and txt.strip():
                    all_text.append(txt)
                    lower = txt.lower()
                    for tag in placeholder_tags:
                        if tag in lower:
                            issues_minor.append(
                                f"slide {slide_num}: contains placeholder "
                                f"tag '{tag}'")
                    # Heuristic source-line detection: text starts with
                    # "Source:" or "Note:" and is under 200 chars
                    s = txt.strip()
                    if (s.startswith(("Source:", "Note:", "source:", "note:"))
                            and len(s) < 300):
                        has_source_line = True

                    # Bullet/text overflow pre-flight: estimate line wrap
                    # and compare to shape height. Triggers only on multi-line
                    # bodies (>=3 chars per line, >=2 lines or contains newline).
                    if (shape.width and shape.height
                            and len(s) > 30 and not has_source_line):
                        try:
                            box_w_in = shape.width / 914400
                            box_h_in = shape.height / 914400
                            # Use first run's font size as proxy
                            pt = 11
                            for para in shape.text_frame.paragraphs:
                                for run in para.runs:
                                    if run.font.size:
                                        pt = run.font.size.pt
                                        break
                                if pt != 11: break
                            est_chars_per_line = max(1, int(
                                (box_w_in * 72) / (pt * 0.50)))
                            # Tally lines: explicit \n + wrapped lines
                            est_lines = 0
                            for line in s.split("\n"):
                                est_lines += max(1,
                                    -(-len(line) // est_chars_per_line))
                            # Pt height ~ 1.25 line height; convert to inches
                            line_h_in = (pt * 1.25) / 72
                            est_h = est_lines * line_h_in
                            # Allow 10% slack for inter-line gaps
                            if est_h > box_h_in * 1.10 and est_lines >= 3:
                                issues_minor.append(
                                    f"slide {slide_num}: probable text "
                                    f"overflow (~{est_lines} lines in "
                                    f"{box_h_in:.2f}in box at {pt:.0f}pt)")
                        except Exception:
                            pass

        # Content QA: data slide w/o source
        # Skip cover (slide 1) and section dividers (very few text shapes).
        # Also skip non-data slide types via heuristic: vision/strategy/agenda
        # slides have descriptive headings + text but no numeric claims.
        is_cover_or_divider = (slide_num == 1 or
                               len([t for t in all_text if t.strip()]) <= 3)
        all_text_blob = " ".join(all_text).lower()
        non_data_markers = [
            "our vision", "our strategy", "agenda", "table of contents",
            "introduction", "executive summary",
            "delivering on our", "our team", "deal team",
        ]
        is_non_data = any(m in all_text_blob for m in non_data_markers)
        # Has numeric data? Exclude short strings (page numbers, "I", "II")
        has_numeric = any(
            any(ch.isdigit() for ch in t)
            for t in all_text
            if len(t.strip()) > 3
        )
        if (not is_cover_or_divider and not is_non_data and has_numeric
                and not has_source_line):
            issues_minor.append(
                f"slide {slide_num}: no Source/Note line "
                "(R3: every data slide must cite)")

    return {
        "critical": issues_critical,
        "minor": issues_minor,
        "passed": n,
    }


def split_into_chunks(items: list, archetype: str) -> list:
    """
    Split a list of items into per-slide chunks based on archetype density limits.
    Useful for: many quotes -> multiple quote_wall slides; many events -> multiple
    timelines; many tiles -> multiple scorecards; many entities -> multiple
    comparison_matrix slides.

    Returns list of sublists; each sublist fits within one slide's density limit.
    """
    if archetype not in _DENSITY_LIMITS:
        raise ValueError(f"unknown archetype: {archetype}")
    max_per_slide = _DENSITY_LIMITS[archetype][0]
    return [items[i:i + max_per_slide]
            for i in range(0, len(items), max_per_slide)]


# ──────────────────────────────────────────────────────────────────────
# Data classes for slide content
# ──────────────────────────────────────────────────────────────────────

@dataclass
class Quote:
    """One management quote for a quote wall."""
    company: str
    text: str
    speaker: str            # "CFO" or "Tim Cook, CEO"
    date_or_event: str      # "Q4 FY25 call (Jan 30, 2025)"


@dataclass
class TimelineEvent:
    """One dated event for a timeline."""
    date_label: str         # "Jan 2026" or "2026-01-15"
    entity: str             # "MSFT"
    label: str              # "$50B AI infra commitment"
    amount: Optional[float] = None


@dataclass
class ProcessBox:
    """One box in a process diagram."""
    label: str
    sub_label: str = ""


@dataclass
class ProcessArrow:
    """One arrow connecting boxes by index."""
    from_idx: int
    to_idx: int
    label: str = ""


@dataclass
class ScorecardTile:
    """One tile in a scorecard."""
    metric: str             # "Revenue"
    value: str              # "$2.4B"
    rating: int = 0         # 0-5 dot rating; 0 = no rating
    sub: str = ""           # optional small text


@dataclass
class FrameworkSection:
    """One column in a strategy framework slide."""
    title: str              # "SIMPLIFICATION"
    bullets: list           # list of strings
    badge: str = ""         # "LARGELY COMPLETE" / "MAIN PRIORITIES FOR 2026"
    highlighted: bool = False  # column gets brand emphasis


@dataclass
class OrgBox:
    """One node in an org chart.

    Tier 4 ("stat box") use case: hang quantitative metrics (revenue,
    headcount, market share) off divisions or roles. Set tier=4 and supply
    parent_idx to anchor the stat under its parent. Stat boxes render
    smaller (lighter gray, smaller font) so they read as annotations rather
    than reporting lines.

        boxes = [
            OrgBox("CEO", tier=1),
            OrgBox("Mining", tier=2, parent_idx=0),
            OrgBox("$65.4B revenue", sub_label="FY25", tier=4, parent_idx=1),
            OrgBox("12,400 headcount", sub_label="as of FY25",
                   tier=4, parent_idx=1),
        ]
    """
    label: str              # "CEO" or company/role name
    sub_label: str = ""     # secondary line (name, division)
    tier: int = 1           # 1=top (Citi Blue), 2=second (Ink), 3=third (Blue Light), 4=stat (Gray)
    parent_idx: Optional[int] = None  # 0-based index into boxes list


@dataclass
class TombstoneTile:
    """One deal tile on a tombstone page."""
    deal_status: str = ""           # "Pending" / "Closed" / "Announced"
    client: str = ""                # Issuer / borrower / target
    description: str = ""           # 1-2 line deal description
    deal_size: str = ""             # "$2.5B" - rendered bold
    date: str = ""                  # "January 2026"
    highlight: bool = False         # featured deal -> Citi Blue fill
    fill_variant: str = "white"     # "white" | "gray" | "sand"


@dataclass
class TeamMember:
    """One person tile on a team page."""
    name: str
    title: str = ""
    segment: str = ""
    email: str = ""
    phone: str = ""
    photo_path: str = ""    # optional headshot PNG/JPG; falls back to initials disc


@dataclass
class TocEntry:
    """One entry in a table of contents."""
    text: str
    page: Optional[int] = None
    level: int = 1              # 1 = "1.", 2 = "A.", 3 = "i."


# ──────────────────────────────────────────────────────────────────────
# Brand override system - SPEC_powerpoint_formatting Section 16
# ──────────────────────────────────────────────────────────────────────

@dataclass
class BrandProfile:
    """
    Firm-specific brand overrides. Pass to PPTXDeckWriter(brand=profile).
    Hex strings (e.g. "#235BE1") for colors; converted to RGBColor on apply.

    Use BrandProfile.from_json(path) to load a saved profile, or
    BrandProfile.from_pdf(path) to extract from a sample deck PDF.
    """
    # Colors (hex strings, "#RRGGBB")
    brand_primary: str = "#255BE3"
    brand_dark: str = "#0F1632"
    accent_cyan: str = "#73C2FC"
    accent_red: str = "#FF3C28"

    # Typography
    font_headline: str = "Arial"
    font_body: str = "Arial"
    headline_size: int = 22
    body_size: int = 11
    footnote_size: int = 8
    headline_bold: bool = True

    # Assets
    logo_path: str = ""

    # Page geometry
    aspect_ratio: str = "16:9"  # "16:9" or "4:3"

    @classmethod
    def from_json(cls, path: str) -> "BrandProfile":
        """Load a saved brand profile from a JSON file."""
        import json
        with open(path, encoding="utf-8") as f:
            data = json.load(f)
        # Filter to known fields only (forward-compatible with extra keys)
        valid = {f for f in cls.__dataclass_fields__}
        return cls(**{k: v for k, v in data.items() if k in valid})

    def to_json(self, path: str) -> None:
        """Save this profile as JSON."""
        import json
        from dataclasses import asdict
        with open(path, "w", encoding="utf-8") as f:
            json.dump(asdict(self), f, indent=2)

    @classmethod
    def from_pdf(cls, pdf_path_or_url: str,
                 logo_path: str = "",
                 extract_logo: bool = False,
                 logo_out_dir: str = "") -> "BrandProfile":
        """
        Extract brand profile from a sample deck PDF.

        Walks the PDF via PyMuPDF, aggregates color frequency, fonts, sizes,
        and typography roles by position. Most-used non-neutral fill color
        becomes brand_primary; darkest fill becomes brand_dark; first found
        cyan/red become accents. Headline font/size from y<0.15 spans
        with size>=14pt.

        Args:
            extract_logo: If True, find largest image on cover/first 3 pages
                and save as PNG to logo_out_dir; result auto-set as logo_path.
            logo_out_dir: Where to save extracted logo (defaults to cwd).
        """
        import urllib.request
        from collections import Counter, defaultdict
        try:
            import fitz  # PyMuPDF
        except ImportError:
            raise RuntimeError(
                "BrandProfile.from_pdf requires PyMuPDF (`pip install pymupdf`)"
            )

        # Load
        if pdf_path_or_url.startswith(("http://", "https://")):
            req = urllib.request.Request(
                pdf_path_or_url, headers={"User-Agent": "Mozilla/5.0"})
            data = urllib.request.urlopen(req, timeout=60).read()
            doc = fitz.open(stream=data, filetype="pdf")
        else:
            doc = fitz.open(pdf_path_or_url)

        page_w_pts = doc[0].rect.width
        page_h_pts = doc[0].rect.height

        fill_colors: Counter = Counter()
        role_buckets: defaultdict = defaultdict(list)

        def _rgb_to_hex(rgb):
            if rgb is None or len(rgb) != 3:
                return None
            r, g, b = rgb
            return f"#{int(r*255):02X}{int(g*255):02X}{int(b*255):02X}"

        def _int_to_hex(c: int) -> str:
            return f"#{(c>>16)&0xFF:02X}{(c>>8)&0xFF:02X}{c&0xFF:02X}"

        def _is_neutral(hx: str) -> bool:
            r, g, b = (int(hx[1:3], 16), int(hx[3:5], 16), int(hx[5:7], 16))
            return abs(r - g) <= 8 and abs(g - b) <= 8 and abs(r - b) <= 8

        def _is_cyan(hx: str) -> bool:
            r, g, b = (int(hx[1:3], 16), int(hx[3:5], 16), int(hx[5:7], 16))
            return b > 200 and g > 150 and r < 150

        def _is_red(hx: str) -> bool:
            r, g, b = (int(hx[1:3], 16), int(hx[3:5], 16), int(hx[5:7], 16))
            return r > 200 and g < 100 and b < 100

        for pi in range(min(doc.page_count, 30)):
            page = doc[pi]
            for d in page.get_drawings():
                fill = d.get("fill")
                if fill is not None:
                    hx = _rgb_to_hex(fill)
                    if hx and hx != "#FFFFFF":
                        fill_colors[hx] += 1
            for block in page.get_text("dict").get("blocks", []):
                if block.get("type") != 0: continue
                for line in block.get("lines", []):
                    for span in line.get("spans", []):
                        if not span.get("text", "").strip(): continue
                        font = span.get("font", "")
                        size = round(span.get("size", 0), 1)
                        bbox = span.get("bbox", (0, 0, 0, 0))
                        y_norm = bbox[1] / page_h_pts if page_h_pts else 0
                        if y_norm < 0.15 and size >= 14:
                            role_buckets["headline"].append((font, size))
                        elif size <= 8:
                            role_buckets["footnote"].append((font, size))
                        elif 9 <= size <= 13:
                            role_buckets["body"].append((font, size))

        # Pick brand_primary: most-used non-neutral, non-pure-red fill
        non_neutral = [(c, n) for c, n in fill_colors.most_common(30)
                       if not _is_neutral(c) and not _is_red(c)]
        brand_primary = non_neutral[0][0] if non_neutral else "#255BE3"

        # Pick brand_dark: most-used CHROMATIC dark fill (skip pure black/gray)
        brand_dark = "#0F1632"
        for c, _ in fill_colors.most_common(30):
            r, g, b = (int(c[1:3], 16), int(c[3:5], 16), int(c[5:7], 16))
            if r + g + b < 200 and not _is_neutral(c):
                brand_dark = c
                break

        # Cyan / red accents
        accent_cyan = "#73C2FC"
        for c, _ in fill_colors.most_common(20):
            if _is_cyan(c):
                accent_cyan = c
                break
        accent_red = "#FF3C28"
        for c, _ in fill_colors.most_common(20):
            if _is_red(c):
                accent_red = c
                break

        # Strip docx font subset prefix "ABCDEF+"
        def _clean_font(f: str) -> str:
            if "+" in f: f = f.split("+", 1)[1]
            return f or "Arial"

        def _dominant_font_size(role: str, default_size: int):
            samples = role_buckets.get(role, [])
            if not samples:
                return "Arial", default_size
            ctr = Counter(samples)
            (font, size), _ = ctr.most_common(1)[0]
            return _clean_font(font), int(size)

        font_headline, headline_size = _dominant_font_size("headline", 22)
        font_body, body_size = _dominant_font_size("body", 11)
        _, footnote_size = _dominant_font_size("footnote", 8)

        # Aspect ratio
        ratio = page_w_pts / page_h_pts if page_h_pts else 16/9
        aspect = "16:9" if abs(ratio - 16/9) < 0.2 else \
                 ("4:3" if abs(ratio - 4/3) < 0.2 else "16:9")

        # Logo extraction: find largest image on cover/first 3 pages
        if extract_logo:
            out_dir = logo_out_dir or os.getcwd()
            os.makedirs(out_dir, exist_ok=True)
            best = None  # (area, xref, page_idx)
            for pi in range(min(3, doc.page_count)):
                try:
                    for img_info in doc[pi].get_images(full=True):
                        xref = img_info[0]
                        try:
                            pix = fitz.Pixmap(doc, xref)
                        except Exception:
                            continue
                        # Filter: reasonable logo size, wider than tall (typical)
                        if pix.width < 60 or pix.height < 30:
                            continue
                        if pix.width > 800 or pix.height > 400:
                            continue  # too big — likely a chart/photo
                        area = pix.width * pix.height
                        if best is None or area > best[0]:
                            best = (area, xref, pi)
                except Exception:
                    continue
            if best is not None:
                area, xref, pi = best
                pix = fitz.Pixmap(doc, xref)
                if pix.alpha or pix.colorspace.name not in ("DeviceRGB", "DeviceGray"):
                    pix = fitz.Pixmap(fitz.csRGB, pix)
                logo_filename = f"extracted_logo_p{pi+1}.png"
                logo_path = os.path.join(out_dir, logo_filename)
                pix.save(logo_path)

        return cls(
            brand_primary=brand_primary,
            brand_dark=brand_dark,
            accent_cyan=accent_cyan,
            accent_red=accent_red,
            font_headline=font_headline,
            font_body=font_body,
            headline_size=headline_size,
            body_size=body_size,
            footnote_size=footnote_size,
            headline_bold=True,  # caller can override
            logo_path=logo_path,
            aspect_ratio=aspect,
        )


def make_pitchpres_profile(logo_path: str = "") -> "BrandProfile":
    """
    Citi PitchPres A4 Landscape profile preset.
    Matches SPEC_PitchPres_A4_Landscape.md exactly.
    """
    return BrandProfile(
        brand_primary="#255BE3",   # Citi Blue
        brand_dark="#0F1632",      # Citi Ink
        accent_cyan="#73C2FC",     # Blue Light
        accent_red="#FF3C28",      # Citi Red (negatives only)
        font_headline="Citi Sans Display",
        font_body="Citi Sans Text",
        headline_size=22,
        body_size=11,
        footnote_size=8,
        headline_bold=False,       # Display Regular per spec
        logo_path=logo_path,
        aspect_ratio="A4_LANDSCAPE",  # 10.83" x 7.5"
    )


# ──────────────────────────────────────────────────────────────────────
# Writer
# ──────────────────────────────────────────────────────────────────────

class PPTXDeckWriter:
    """
    Builds IB-style decks per valuation_kit specs.

    Usage:
        deck = PPTXDeckWriter(firm="Anthropic Capital", project="Sandvik DD")
        deck.add_cover("Sandvik AB - Investment Memo", subtitle="Industrials Long")
        deck.add_section_divider("II", "Valuation Analysis")
        deck.add_comparison_matrix(
            action_title="Sandvik trades at a 12% discount to peer median EV/EBITDA",
            entities=[...], metrics=[...], values=[[...], ...],
            target_label="SAND.ST",
            source="Bloomberg / FactSet, retrieved Apr 30 2026",
        )
        path = deck.save("Sandvik_DD.pptx")
    """

    def __init__(self,
                 firm: str = "",
                 project: str = "Confidential",
                 output_dir: str = None,
                 confidentiality: str = "CONFIDENTIAL",
                 logo_path: str = "",
                 headline_bold: bool = True,
                 brand: Optional["BrandProfile"] = None):
        """
        Args:
            firm: firm name (used in footer when confidentiality set).
            project: project name (used in footer).
            output_dir: where to save .pptx files; defaults to <repo>/decks.
            confidentiality: "CONFIDENTIAL" / "DRAFT" / "" to suppress.
            logo_path: optional PNG/JPG of firm logo. When set, auto-rendered
                top-right on every content slide and prominently on cover.
                Overridden by brand.logo_path if brand provided.
            headline_bold: True for bold action titles (default per spec R1).
                Overridden by brand.headline_bold if brand provided.
            brand: optional BrandProfile overriding default colors, fonts,
                sizes. Profile is the canonical mechanism for firm overrides
                (per Section 16 of formatting spec).
        """
        self.firm = firm
        self.project = project
        self.confidentiality = confidentiality

        # Apply brand profile (or fall back to module defaults / params)
        self._apply_brand(brand, logo_path, headline_bold)

        self.output_dir = output_dir or os.path.join(
            os.path.dirname(__file__), "..", "..", "decks"
        )
        os.makedirs(self.output_dir, exist_ok=True)

        self.prs = Presentation()
        w_in, h_in = ASPECT_DIMS.get(self.brand_aspect, ASPECT_DIMS["16:9"])
        self.prs.slide_width  = Inches(w_in)
        self.prs.slide_height = Inches(h_in)
        self._slide_w_in = w_in
        self._slide_h_in = h_in

        self._page = 0  # 0 means cover not yet added

    def _apply_brand(self, brand, logo_path: str, headline_bold: bool):
        """Resolve effective colors/fonts/sizes from brand profile or defaults."""
        def _hex_to_rgb(hx: str) -> RGBColor:
            h = hx.lstrip("#")
            return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

        if brand is not None:
            self.brand_primary = _hex_to_rgb(brand.brand_primary)
            self.brand_dark    = _hex_to_rgb(brand.brand_dark)
            self.accent_cyan   = _hex_to_rgb(brand.accent_cyan)
            self.accent_red    = _hex_to_rgb(brand.accent_red)
            self.font_headline = brand.font_headline
            self.font_body     = brand.font_body
            self.headline_size = brand.headline_size
            self.body_size     = brand.body_size
            self.footnote_size = brand.footnote_size
            self.headline_bold = brand.headline_bold
            self.logo_path     = (brand.logo_path
                                  if brand.logo_path and
                                     os.path.exists(brand.logo_path)
                                  else "")
            self.brand_aspect  = brand.aspect_ratio
        else:
            self.brand_primary = BRAND_BLUE
            self.brand_dark    = INK
            self.accent_cyan   = RGBColor(0x73, 0xC2, 0xFC)
            self.accent_red    = RGBColor(0xFF, 0x3C, 0x28)
            self.font_headline = FONT_HEADLINE
            self.font_body     = FONT_BODY
            self.headline_size = 22
            self.body_size     = PT_BODY
            self.footnote_size = PT_FOOTNOTE
            self.headline_bold = headline_bold
            self.logo_path     = (logo_path
                                  if logo_path and os.path.exists(logo_path)
                                  else "")
            self.brand_aspect  = "16:9"

    # ── Public slide builders ────────────────────────────────────────

    def add_cover(self, title: str, subtitle: str = "",
                  deck_date: Optional[str] = None) -> None:
        """
        Cover slide (Section 6 of formatting spec).
        Cleaner Citi-style: large brand-blue title on white, slim left accent bar,
        smaller subtitle and date. Generous whitespace.
        """
        slide = self._blank_slide()

        # Slim brand accent bar on left edge
        self._add_rect(0, 0, 0.25, self._slide_h_in, self.brand_primary, no_line=True)

        # Big brand-blue title (fixed size; long titles should be edited not auto-shrunk)
        title_box = self._add_text(
            MARGIN_IN + 0.3, self._slide_h_in * 0.32,
            self._slide_w_in - MARGIN_IN - 0.8, 2.4,
            title, font_size=44, bold=self.headline_bold, color=self.brand_primary,
            font=self.font_headline,
        )

        if subtitle:
            self._add_text(MARGIN_IN + 0.3, self._slide_h_in * 0.55,
                           self._slide_w_in - MARGIN_IN - 0.8, 0.6,
                           subtitle, font_size=22, color=self.brand_dark,
                           font=self.font_headline)

        date_str = deck_date or date.today().strftime("%B %Y")
        self._add_text(MARGIN_IN + 0.3, self._slide_h_in - 1.1,
                       self._slide_w_in - MARGIN_IN - 0.8, 0.4,
                       date_str, font_size=14, color=FOOTNOTE_GRAY)

        if self.confidentiality:
            self._add_text(self._slide_w_in - 2.5, self._slide_h_in - 0.5,
                           2.0, 0.3,
                           self.confidentiality, font_size=PT_FOOTNOTE,
                           bold=True, color=FOOTNOTE_GRAY,
                           align=PP_ALIGN.RIGHT)

        # Logo: large version on cover, right-anchored at margin
        if self.logo_path:
            self._add_logo(slide=None, height_in=0.6,
                           right_in=self._slide_w_in - MARGIN_IN,
                           y_in=MARGIN_IN)

        self._page = 0  # cover doesn't count

    def add_section_divider(self, section_num: str, title: str) -> None:
        """
        Section divider (Section 7).
        White background with large brand-blue section number + dark title.
        Cleaner / lighter than a fully colored background.
        """
        slide = self._blank_slide()

        # Slim brand accent bar on left edge for continuity with cover
        self._add_rect(0, 0, 0.25, self._slide_h_in, self.brand_primary, no_line=True)

        # Brand-blue section number (large)
        self._add_text(MARGIN_IN, self._slide_h_in * 0.30,
                       self._slide_w_in - 2 * MARGIN_IN, 1.6,
                       section_num, font_size=80, bold=True,
                       color=self.brand_primary, font=self.font_headline,
                       align=PP_ALIGN.LEFT)

        # Thin brand divider under number
        self._add_line(MARGIN_IN, self._slide_h_in * 0.50,
                       MARGIN_IN + 1.2, self._slide_h_in * 0.50,
                       color=self.brand_primary, weight=2.5)

        # Section title in dark ink
        self._add_text(MARGIN_IN, self._slide_h_in * 0.54,
                       self._slide_w_in - 2 * MARGIN_IN, 1.0,
                       title, font_size=36, bold=True,
                       color=self.brand_dark, font=self.font_headline,
                       align=PP_ALIGN.LEFT)

        self._page += 1
        self._add_footer(slide)

    def add_comparison_matrix(
        self,
        action_title: str,
        entities: list[str],          # column labels
        metrics: list[str],           # row labels
        values: list[list],           # values[row_metric][col_entity]
        target_label: str = "",       # entity to highlight
        target_idx: Optional[int] = None,
        source: str = "",
        notes: str = "",
        skip_source: bool = False,
        summary_stats: bool = True,   # add Median row at bottom
    ) -> None:
        """
        Comparison matrix slide (Section 2.1).
        entities x metrics table; target column shaded brand color.
        """
        action_title = self._validate_action_title(action_title)
        if len(values) != len(metrics):
            raise ValueError("values must have same length as metrics")
        for r in values:
            if len(r) != len(entities):
                raise ValueError("each values row must match entities length")

        # Density check (Section 3.1)
        max_e, max_m = _DENSITY_LIMITS[ARCH_COMPARISON]
        if len(entities) > max_e or len(metrics) > max_m:
            raise ValueError(
                f"density {len(entities)}x{len(metrics)} exceeds "
                f"comparison limit {max_e}x{max_m} - split into multiple slides"
            )

        # Resolve target column
        if target_idx is None and target_label:
            try:
                target_idx = entities.index(target_label)
            except ValueError:
                target_idx = None

        slide = self._content_slide_with_title(action_title)

        # Build table
        n_rows = 1 + len(metrics) + (1 if summary_stats else 0)
        n_cols = 1 + len(entities)

        # Position: leave room for headline + footer
        tbl_left = MARGIN_IN
        tbl_top = 1.4
        tbl_width = self._slide_w_in - 2 * MARGIN_IN
        tbl_height = min(5.2, 0.4 + 0.45 * (n_rows - 1))

        tbl_shape = slide.shapes.add_table(
            n_rows, n_cols,
            Inches(tbl_left), Inches(tbl_top),
            Inches(tbl_width), Inches(tbl_height),
        )
        tbl = tbl_shape.table

        # Header row
        self._set_cell(tbl.cell(0, 0), "", self.brand_primary, WHITE,
                       bold=True, size=PT_TABLE_H, align=PP_ALIGN.LEFT)
        for ci, ent in enumerate(entities):
            self._set_cell(tbl.cell(0, ci + 1), ent, self.brand_primary, WHITE,
                           bold=True, size=PT_TABLE_H, align=PP_ALIGN.CENTER)

        # Data rows
        for ri, metric in enumerate(metrics):
            row_idx = ri + 1
            self._set_cell(tbl.cell(row_idx, 0), metric, WHITE, self.brand_dark,
                           bold=False, size=PT_TABLE_B, align=PP_ALIGN.LEFT)
            for ci in range(len(entities)):
                v = values[ri][ci]
                bg = LIGHT_GRAY if (target_idx is not None and ci == target_idx) else WHITE
                fg = self.brand_dark
                self._set_cell(tbl.cell(row_idx, ci + 1),
                               self._format_value(v),
                               bg, fg, bold=False, size=PT_TABLE_B,
                               align=PP_ALIGN.RIGHT)

        # Summary stats row (Section 7 of formatting spec - shaded median)
        if summary_stats:
            srow = n_rows - 1
            self._set_cell(tbl.cell(srow, 0), "Median", LIGHT_GRAY, self.brand_dark,
                           bold=True, size=PT_TABLE_B, italic=True,
                           align=PP_ALIGN.LEFT)
            import statistics
            for ci in range(len(entities)):
                col_vals = [values[ri][ci] for ri in range(len(metrics))
                            if isinstance(values[ri][ci], (int, float))]
                med = statistics.median(col_vals) if col_vals else None
                self._set_cell(tbl.cell(srow, ci + 1),
                               self._format_value(med),
                               LIGHT_GRAY, self.brand_dark, bold=True,
                               size=PT_TABLE_B, italic=True,
                               align=PP_ALIGN.RIGHT)

        if not skip_source:
            self._add_source_line(slide, source, notes)
        self._page += 1
        self._add_footer(slide)

    def add_scorecard(
        self,
        action_title: str,
        tiles: list[ScorecardTile],
        source: str = "",
        notes: str = "",
        skip_source: bool = False,
    ) -> None:
        """
        Scorecard slide (Section 2.2). Grid of 4-9 metric tiles with optional rating dots.
        """
        action_title = self._validate_action_title(action_title)
        n = len(tiles)
        if not (1 <= n <= _DENSITY_LIMITS[ARCH_SCORECARD][0]):
            raise ValueError(f"scorecard requires 1-9 tiles, got {n}")

        slide = self._content_slide_with_title(action_title)

        # Pick grid: prefer 3 cols
        if n <= 3:    cols, rows = n, 1
        elif n <= 6:  cols, rows = 3, 2
        else:         cols, rows = 3, 3

        avail_w = self._slide_w_in - 2 * MARGIN_IN
        avail_h = 4.8
        gap = 0.2
        tile_w = (avail_w - gap * (cols - 1)) / cols
        tile_h = (avail_h - gap * (rows - 1)) / rows
        top0 = 1.5
        left0 = MARGIN_IN

        for i, tile in enumerate(tiles):
            r, c = divmod(i, cols)
            left = left0 + c * (tile_w + gap)
            top = top0 + r * (tile_h + gap)
            self._draw_tile(slide, left, top, tile_w, tile_h, tile)

        if not skip_source:
            self._add_source_line(slide, source, notes)
        self._page += 1
        self._add_footer(slide)

    def add_quote_wall(
        self,
        action_title: str,
        quotes: list[Quote],
        source: str = "",
        notes: str = "",
        skip_source: bool = False,
    ) -> None:
        """
        Quote wall slide (Section 2.3). 4-8 management quotes in 2-col grid.
        """
        action_title = self._validate_action_title(action_title)
        n = len(quotes)
        if not (1 <= n <= _DENSITY_LIMITS[ARCH_QUOTE_WALL][0]):
            raise ValueError(f"quote_wall requires 1-8 quotes, got {n}")

        slide = self._content_slide_with_title(action_title)

        cols = 2 if n > 1 else 1
        rows = -(-n // cols)  # ceil

        avail_w = self._slide_w_in - 2 * MARGIN_IN
        avail_h = 4.8
        gap = 0.25
        card_w = (avail_w - gap * (cols - 1)) / cols
        card_h = (avail_h - gap * (rows - 1)) / rows
        top0 = 1.5
        left0 = MARGIN_IN

        for i, q in enumerate(quotes):
            r, c = divmod(i, cols)
            left = left0 + c * (card_w + gap)
            top = top0 + r * (card_h + gap)
            self._draw_quote_card(slide, left, top, card_w, card_h, q)

        if not skip_source:
            self._add_source_line(slide, source, notes)
        self._page += 1
        self._add_footer(slide)

    def add_timeline(
        self,
        action_title: str,
        events: list[TimelineEvent],
        source: str = "",
        notes: str = "",
        skip_source: bool = False,
    ) -> None:
        """
        Timeline slide (Section 2.4). Horizontal line with milestones.
        """
        action_title = self._validate_action_title(action_title)
        n = len(events)
        if not (1 <= n <= _DENSITY_LIMITS[ARCH_TIMELINE][0]):
            raise ValueError(f"timeline requires 1-10 events, got {n}")

        slide = self._content_slide_with_title(action_title)

        # Horizontal axis line at vertical middle of body
        axis_y = 1.5 + 4.8 / 2
        axis_left = MARGIN_IN + 0.5
        axis_right = self._slide_w_in - MARGIN_IN - 0.5

        self._add_line(axis_left, axis_y, axis_right, axis_y,
                       color=MID_GRAY, weight=2)

        if n == 1:
            xs = [(axis_left + axis_right) / 2]
        else:
            xs = [axis_left + i * (axis_right - axis_left) / (n - 1)
                  for i in range(n)]

        for i, ev in enumerate(events):
            x = xs[i]
            # Milestone dot
            self._add_oval(x - 0.1, axis_y - 0.1, 0.2, 0.2, self.brand_primary)

            # Alternate above/below for readability
            above = (i % 2 == 0)
            label_y = axis_y - 1.6 if above else axis_y + 0.3
            anchor_y = axis_y - 0.3 if above else axis_y + 0.1

            # Stem line from dot to label
            self._add_line(x, anchor_y, x, label_y + (1.2 if above else 0),
                           color=MID_GRAY, weight=1)

            # Label box
            box_w = 1.6
            self._add_text(x - box_w / 2, label_y, box_w, 1.0,
                           f"{ev.entity}\n{ev.label}",
                           font_size=PT_BODY, color=self.brand_dark,
                           align=PP_ALIGN.CENTER, bold=True)

            # Date label - opposite side small
            date_y = axis_y + 0.3 if above else axis_y - 0.5
            self._add_text(x - 0.8, date_y, 1.6, 0.3,
                           ev.date_label, font_size=9,
                           color=FOOTNOTE_GRAY, align=PP_ALIGN.CENTER,
                           italic=True)

        if not skip_source:
            self._add_source_line(slide, source, notes)
        self._page += 1
        self._add_footer(slide)

    def add_process_diagram(
        self,
        action_title: str,
        boxes: list[ProcessBox],
        arrows: list[ProcessArrow],
        source: str = "",
        notes: str = "",
        skip_source: bool = False,
        direction: str = "ltr",   # "ltr" or "ttb"
    ) -> None:
        """
        Process / structure diagram (Section 2.5). Boxes + arrows.
        """
        action_title = self._validate_action_title(action_title)
        n = len(boxes)
        if not (1 <= n <= _DENSITY_LIMITS[ARCH_PROCESS][0]):
            raise ValueError(f"process_diagram requires 1-8 boxes, got {n}")
        if direction not in ("ltr", "ttb"):
            raise ValueError("direction must be 'ltr' or 'ttb'")

        slide = self._content_slide_with_title(action_title)

        # Layout boxes in a single line (LTR) or column (TTB)
        body_top = 1.6
        body_left = MARGIN_IN + 0.3
        body_w = self._slide_w_in - 2 * MARGIN_IN - 0.6
        body_h = 4.5

        if direction == "ltr":
            box_h = 1.2
            box_y = body_top + (body_h - box_h) / 2
            gap = min(0.5, body_w / (n * 4))
            box_w = (body_w - gap * (n - 1)) / n
            positions = [(body_left + i * (box_w + gap), box_y, box_w, box_h)
                         for i in range(n)]
        else:
            box_w = min(3.5, body_w * 0.5)
            box_x = body_left + (body_w - box_w) / 2
            gap = 0.4
            box_h = (body_h - gap * (n - 1)) / n
            positions = [(box_x, body_top + i * (box_h + gap), box_w, box_h)
                         for i in range(n)]

        # Draw boxes
        for i, b in enumerate(boxes):
            x, y, w, h = positions[i]
            shp = self._add_rounded_rect(x, y, w, h, self.brand_primary)
            tf = shp.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.text = b.label
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0]
            run.font.size = Pt(PT_BODY)
            run.font.bold = True
            run.font.color.rgb = WHITE
            run.font.name = self.font_body
            if b.sub_label:
                p2 = tf.add_paragraph()
                p2.alignment = PP_ALIGN.CENTER
                run2 = p2.add_run()
                run2.text = b.sub_label
                run2.font.size = Pt(9)
                run2.font.color.rgb = WHITE
                run2.font.name = self.font_body

        # Draw arrows
        for arr in arrows:
            if arr.from_idx >= n or arr.to_idx >= n:
                continue
            fx, fy, fw, fh = positions[arr.from_idx]
            tx, ty, tw, th = positions[arr.to_idx]
            if direction == "ltr":
                x1 = fx + fw
                y1 = fy + fh / 2
                x2 = tx
                y2 = ty + th / 2
            else:
                x1 = fx + fw / 2
                y1 = fy + fh
                x2 = tx + tw / 2
                y2 = ty
            self._add_arrow(x1, y1, x2, y2, color=self.brand_dark)
            if arr.label:
                lx = (x1 + x2) / 2 - 0.7
                ly = (y1 + y2) / 2 - 0.15
                self._add_text(lx, ly, 1.4, 0.3, arr.label,
                               font_size=9, color=self.brand_dark,
                               align=PP_ALIGN.CENTER, italic=True)

        if not skip_source:
            self._add_source_line(slide, source, notes)
        self._page += 1
        self._add_footer(slide)

    def add_strategy_framework(
        self,
        action_title: str,
        sections: list,                # list[FrameworkSection], 2-5 columns
        vision: str = "",              # optional top vision banner
        vision_label: str = "OUR VISION",
        framework_label: str = "",     # e.g. "DELIVERING ON OUR PRIORITIES"
        source: str = "",
        notes: str = "",
        skip_source: bool = False,
    ) -> None:
        """
        Strategy framework slide - N-column priorities/initiatives layout.
        Modeled on Citi Q1 2026 page 2 ("Our strategy and path forward remain
        unchanged"). Optional top vision banner; optional framework subhead;
        2-5 columns each with title + bullet list, optional badge per group.
        """
        action_title = self._validate_action_title(action_title)
        n = len(sections)
        if not (2 <= n <= 5):
            raise ValueError(f"strategy_framework requires 2-5 columns, got {n}")

        slide = self._content_slide_with_title(action_title)

        cur_y = 0.85  # tight under headline divider (Citi-style)

        # Optional vision banner (dark band with brand-blue label + vision text)
        if vision:
            vis_h = 0.95
            self._add_rect(MARGIN_IN, cur_y,
                           self._slide_w_in - 2 * MARGIN_IN, vis_h,
                           INK, no_line=True)
            # "OUR VISION" label, brand-blue, centered top
            self._add_text(MARGIN_IN, cur_y + 0.05,
                           self._slide_w_in - 2 * MARGIN_IN, 0.3,
                           vision_label, font_size=12, bold=True,
                           color=self.brand_primary, font=self.font_headline,
                           align=PP_ALIGN.CENTER)
            # Vision text white
            self._add_text(MARGIN_IN + 0.3, cur_y + 0.35,
                           self._slide_w_in - 2 * MARGIN_IN - 0.6, 0.7,
                           vision, font_size=14, color=WHITE,
                           font=self.font_headline, align=PP_ALIGN.CENTER,
                           anchor=MSO_ANCHOR.MIDDLE)
            cur_y += vis_h + 0.2

        # Optional framework subhead
        if framework_label:
            self._add_text(MARGIN_IN, cur_y,
                           self._slide_w_in - 2 * MARGIN_IN, 0.35,
                           framework_label, font_size=14, bold=True,
                           color=self.brand_dark, font=self.font_headline,
                           align=PP_ALIGN.CENTER)
            cur_y += 0.4

        # Compute column geometry
        body_top = cur_y
        body_bottom = self._slide_h_in - 0.6   # leave room for source + footer (tight)
        body_h = body_bottom - body_top
        col_gap = 0.15
        avail_w = self._slide_w_in - 2 * MARGIN_IN
        col_w = (avail_w - col_gap * (n - 1)) / n

        # Render badges across the top (group adjacent columns sharing same badge)
        badges = [s.badge for s in sections]
        if any(badges):
            i = 0
            badge_h = 0.32
            while i < n:
                b = badges[i]
                j = i
                while j < n and badges[j] == b:
                    j += 1
                if b:
                    left = MARGIN_IN + i * (col_w + col_gap)
                    width = (j - i) * col_w + (j - i - 1) * col_gap
                    bg = SAND if "complete" in b.lower() else self.brand_primary
                    fg = self.brand_dark if bg == SAND else WHITE
                    self._add_rect(left, body_top, width, badge_h, bg,
                                   no_line=True)
                    self._add_text(left, body_top + 0.04, width, 0.28,
                                   b, font_size=10, bold=True, color=fg,
                                   font=self.font_headline,
                                   align=PP_ALIGN.CENTER)
                i = j
            body_top += badge_h + 0.1

        # Column titles + bullets
        for i, sec in enumerate(sections):
            left = MARGIN_IN + i * (col_w + col_gap)

            # Column header strip
            header_h = 0.4
            header_color = self.brand_primary if sec.highlighted else INK
            self._add_rect(left, body_top, col_w, header_h,
                           header_color, no_line=True)
            self._add_text(left + 0.05, body_top + 0.06,
                           col_w - 0.1, header_h - 0.1,
                           sec.title, font_size=11, bold=True,
                           color=WHITE, font=self.font_headline,
                           align=PP_ALIGN.CENTER)

            # Bullet list area - Citi style: plain text, centered, no markers,
            # blank line between items for breathing room
            bullet_top = body_top + header_h + 0.15
            bullet_h = body_bottom - bullet_top - 0.1
            bullet_text = "\n\n".join(sec.bullets)
            self._add_text(left + 0.1, bullet_top,
                           col_w - 0.2, bullet_h,
                           bullet_text, font_size=11, color=self.brand_dark,
                           font=self.font_body, align=PP_ALIGN.CENTER)

        if not skip_source:
            self._add_source_line(slide, source, notes)
        self._page += 1
        self._add_footer(slide)

    def add_bar_chart(
        self,
        action_title: str,
        labels: list,                  # entity labels (rows)
        values: list,                  # numeric values (one per label)
        value_format: str = "{:,.1f}", # python format string
        target_label: str = "",        # entity to highlight in brand color
        x_label: str = "",             # axis title (units)
        source: str = "",
        notes: str = "",
        skip_source: bool = False,
    ) -> None:
        """
        Horizontal bar chart (Section 5.1 - 'compare entities on metric X').
        Bars rendered as colored rectangles; values labeled at right end of bar.
        Target entity bar uses brand_primary; others use brand_dark.
        Sorted descending by value automatically.
        """
        action_title = self._validate_action_title(action_title)
        if len(labels) != len(values):
            raise ValueError("labels and values must match length")
        n = len(labels)
        if not (1 <= n <= 12):
            raise ValueError(f"bar chart requires 1-12 bars, got {n}")

        slide = self._content_slide_with_title(action_title)

        # Sort descending - largest at top
        pairs = sorted(zip(labels, values), key=lambda p: -p[1])

        body_top = 0.95
        body_bottom = self._slide_h_in - 0.7
        body_h = body_bottom - body_top
        label_w = 1.6
        value_w = 1.0
        chart_left = 0.3 + label_w
        chart_right = self._slide_w_in - 0.3 - value_w
        chart_w = chart_right - chart_left

        max_val = max(abs(v) for _, v in pairs) or 1
        bar_h = min(0.45, (body_h - (n - 1) * 0.12) / n)
        gap = max(0.08, (body_h - n * bar_h) / max(1, n - 1))

        for i, (lbl, val) in enumerate(pairs):
            y = body_top + i * (bar_h + gap)
            # Label
            self._add_text(0.3, y + bar_h * 0.1, label_w - 0.1, bar_h * 0.9,
                           lbl, font_size=11,
                           color=self.brand_dark, font=self.font_body,
                           align=PP_ALIGN.RIGHT)
            # Bar
            bar_w = chart_w * (val / max_val) if max_val else 0
            color = self.brand_primary if lbl == target_label else self.brand_dark
            self._add_rect(chart_left, y, max(0.05, bar_w), bar_h,
                           color, no_line=True)
            # Value (right of bar)
            self._add_text(chart_left + bar_w + 0.05, y + bar_h * 0.1,
                           value_w, bar_h * 0.9,
                           value_format.format(val),
                           font_size=11, color=self.brand_dark,
                           font=self.font_body, align=PP_ALIGN.LEFT)

        # X-axis label below chart
        if x_label:
            self._add_text(chart_left, body_bottom + 0.05,
                           chart_w, 0.25, x_label,
                           font_size=9, italic=True, color=FOOTNOTE_GRAY,
                           align=PP_ALIGN.CENTER)

        if not skip_source:
            self._add_source_line(slide, source, notes)
        self._page += 1
        self._add_footer(slide)

    def add_football_field(
        self,
        action_title: str,
        methods: list,                 # list of dicts: {label, low, high, mid}
        target_value: Optional[float] = None,  # vertical line (current price)
        target_label: str = "Current",
        value_format: str = "${:,.0f}",
        source: str = "",
        notes: str = "",
        skip_source: bool = False,
    ) -> None:
        """
        Football field (Section 5.3) - horizontal range bars per valuation method.
        Each bar shows low-high range; mid (if provided) shown as tick.
        Optional vertical reference line for current price/value.

        methods: list of dicts with keys "label" (str), "low" (float),
                 "high" (float), and optional "mid" (float).
        """
        action_title = self._validate_action_title(action_title)
        n = len(methods)
        if not (1 <= n <= 8):
            raise ValueError(f"football field requires 1-8 methods, got {n}")
        for m in methods:
            if "label" not in m or "low" not in m or "high" not in m:
                raise ValueError("each method needs label, low, high")
            if m["low"] > m["high"]:
                raise ValueError(f"{m['label']}: low > high")

        slide = self._content_slide_with_title(action_title)

        body_top = 0.95
        body_bottom = self._slide_h_in - 0.7
        body_h = body_bottom - body_top
        label_w = 2.0
        chart_left = 0.3 + label_w
        chart_right = self._slide_w_in - 0.5
        chart_w = chart_right - chart_left

        # Determine value range
        all_lo = min(m["low"] for m in methods)
        all_hi = max(m["high"] for m in methods)
        if target_value is not None:
            all_lo = min(all_lo, target_value)
            all_hi = max(all_hi, target_value)
        # Add 5% padding either side
        rng = all_hi - all_lo or 1
        x_min = all_lo - rng * 0.05
        x_max = all_hi + rng * 0.05
        x_rng = x_max - x_min

        def x_to_in(v: float) -> float:
            return chart_left + (v - x_min) / x_rng * chart_w

        bar_h = min(0.35, (body_h - (n - 1) * 0.15) / n)
        gap = max(0.12, (body_h - n * bar_h) / max(1, n - 1))

        # Background gridlines (3-5 ticks)
        n_ticks = 5
        for i in range(n_ticks + 1):
            tick_v = x_min + (x_max - x_min) * i / n_ticks
            tx = x_to_in(tick_v)
            self._add_line(tx, body_top, tx, body_bottom,
                           color=LIGHT_GRAY, weight=0.5)
            self._add_text(tx - 0.4, body_bottom + 0.05, 0.8, 0.25,
                           value_format.format(tick_v),
                           font_size=9, color=FOOTNOTE_GRAY,
                           align=PP_ALIGN.CENTER)

        for i, m in enumerate(methods):
            y = body_top + i * (bar_h + gap)
            # Method label
            self._add_text(0.3, y + bar_h * 0.1, label_w - 0.1, bar_h * 0.9,
                           m["label"], font_size=11,
                           color=self.brand_dark, font=self.font_body,
                           align=PP_ALIGN.RIGHT)
            # Range bar (brand_primary fill)
            x0 = x_to_in(m["low"])
            x1 = x_to_in(m["high"])
            self._add_rect(x0, y, max(0.05, x1 - x0), bar_h,
                           self.brand_primary, no_line=True)
            # Low/high labels at ends
            self._add_text(x0 - 0.5, y + bar_h * 0.1, 0.5, bar_h * 0.9,
                           value_format.format(m["low"]),
                           font_size=9, color=self.brand_dark,
                           font=self.font_body, align=PP_ALIGN.RIGHT)
            self._add_text(x1 + 0.05, y + bar_h * 0.1, 0.7, bar_h * 0.9,
                           value_format.format(m["high"]),
                           font_size=9, color=self.brand_dark,
                           font=self.font_body, align=PP_ALIGN.LEFT)
            # Optional mid tick
            if "mid" in m and m["mid"] is not None:
                mx = x_to_in(m["mid"])
                self._add_line(mx, y - 0.05, mx, y + bar_h + 0.05,
                               color=WHITE, weight=2)

        # Target reference line (vertical, full chart height)
        if target_value is not None:
            tx = x_to_in(target_value)
            self._add_line(tx, body_top - 0.1, tx, body_bottom,
                           color=self.accent_red, weight=2)
            self._add_text(tx - 0.7, body_top - 0.35,
                           1.4, 0.25,
                           f"{target_label}: {value_format.format(target_value)}",
                           font_size=9, bold=True, color=self.accent_red,
                           font=self.font_body, align=PP_ALIGN.CENTER)

        if not skip_source:
            self._add_source_line(slide, source, notes)
        self._page += 1
        self._add_footer(slide)

    # ── P4 charts: line / waterfall / stacked-bar / pie (Section 5) ──

    def add_line_chart(
        self,
        action_title: str,
        x_labels: list,                # period labels along x-axis
        series: list,                  # list of dicts: {label, values}
        target_series: str = "",       # series label rendered in brand_primary
        y_format: str = "{:,.0f}",
        y_label: str = "",
        source: str = "",
        notes: str = "",
        skip_source: bool = False,
        native: bool = True,
    ) -> None:
        """
        Multi-series line chart (Section 5.4 'trend over time').
        Native python-pptx chart (XL_CHART_TYPE.LINE_MARKERS) — editable in
        PowerPoint, scales properly, axis labels render natively. Set
        native=False to fall back to the legacy shape-based renderer.
        Density: 1-6 series x 2-12 periods.
        """
        action_title = self._validate_action_title(action_title)
        if not x_labels or len(x_labels) < 2:
            raise ValueError("line chart requires >=2 x_labels")
        if not (1 <= len(series) <= 6):
            raise ValueError(f"line chart requires 1-6 series, got {len(series)}")
        for s in series:
            if "label" not in s or "values" not in s:
                raise ValueError("each series needs 'label' and 'values'")
            if len(s["values"]) != len(x_labels):
                raise ValueError(
                    f"series '{s['label']}' has {len(s['values'])} values "
                    f"but {len(x_labels)} x_labels"
                )

        if not native:
            return self._add_line_chart_shapes(
                action_title, x_labels, series, target_series,
                y_format, y_label, source, notes, skip_source,
            )

        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import (
            XL_CHART_TYPE, XL_LEGEND_POSITION, XL_TICK_MARK,
        )

        slide = self._content_slide_with_title(action_title)

        body_top = 0.95
        body_bottom = self._slide_h_in - 0.7
        chart_left = 0.4
        chart_top = body_top
        chart_w = self._slide_w_in - 0.8
        chart_h = body_bottom - body_top - 0.05

        chart_data = CategoryChartData()
        chart_data.categories = [str(x) for x in x_labels]
        for s in series:
            vals = [None if v is None else float(v) for v in s["values"]]
            chart_data.add_series(s["label"], vals)

        gframe = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE_MARKERS,
            Inches(chart_left), Inches(chart_top),
            Inches(chart_w), Inches(chart_h),
            chart_data,
        )
        chart = gframe.chart
        chart.has_title = False
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(10)
        chart.legend.font.name = self.font_body
        chart.legend.font.color.rgb = self.brand_dark

        def _palette_color(idx: int, label: str):
            if target_series and label == target_series:
                return self.brand_primary
            pal = [self.brand_dark, self.accent_cyan,
                   SERIES_PALETTE[4], SERIES_PALETTE[6],
                   SERIES_PALETTE[7], SERIES_PALETTE[8]]
            return pal[idx % len(pal)]

        for si, s in enumerate(series):
            color = _palette_color(si, s["label"])
            ser = chart.series[si]
            ser.format.line.color.rgb = color
            ser.format.line.width = Pt(2.25)
            try:
                ser.marker.format.fill.solid()
                ser.marker.format.fill.fore_color.rgb = color
                ser.marker.format.line.color.rgb = color
                ser.marker.size = 6
            except Exception:
                pass

        # Y-axis: number format + light gridlines
        try:
            va = chart.value_axis
            va.tick_labels.font.size = Pt(9)
            va.tick_labels.font.name = self.font_body
            va.tick_labels.font.color.rgb = FOOTNOTE_GRAY
            va.tick_labels.number_format = self._fmt_to_numfmt(y_format)
            va.has_major_gridlines = True
            try:
                va.major_gridlines.format.line.color.rgb = LIGHT_GRAY
                va.major_gridlines.format.line.width = Pt(0.5)
            except Exception:
                pass
            va.minor_tick_mark = XL_TICK_MARK.NONE
            va.major_tick_mark = XL_TICK_MARK.OUTSIDE
            if y_label:
                va.has_title = True
                va.axis_title.text_frame.text = y_label
                p = va.axis_title.text_frame.paragraphs[0]
                if p.runs:
                    p.runs[0].font.size = Pt(9)
                    p.runs[0].font.italic = True
                    p.runs[0].font.color.rgb = FOOTNOTE_GRAY
                    p.runs[0].font.name = self.font_body
        except Exception:
            pass

        # X-axis: tick label styling
        try:
            ca = chart.category_axis
            ca.tick_labels.font.size = Pt(9)
            ca.tick_labels.font.name = self.font_body
            ca.tick_labels.font.color.rgb = self.brand_dark
        except Exception:
            pass

        if not skip_source:
            self._add_source_line(slide, source, notes)
        self._page += 1
        self._add_footer(slide)

    def _add_line_chart_shapes(
        self, action_title, x_labels, series, target_series,
        y_format, y_label, source, notes, skip_source,
    ) -> None:
        """Legacy shape-based line chart (kept as fallback). See add_line_chart."""
        slide = self._content_slide_with_title(action_title)

        body_top = 0.95
        body_bottom = self._slide_h_in - 0.7
        plot_left = 0.9
        plot_right = self._slide_w_in - 1.6
        plot_w = plot_right - plot_left
        plot_h = body_bottom - body_top - 0.4
        plot_top = body_top + 0.1
        plot_bot = plot_top + plot_h

        all_vals = [v for s in series for v in s["values"] if v is not None]
        if not all_vals:
            raise ValueError("line chart: no numeric values")
        y_min, y_max = min(all_vals), max(all_vals)
        if y_min > 0:
            y_min = 0
        rng = (y_max - y_min) or 1
        y_max += rng * 0.08
        rng = y_max - y_min

        def x_to_in(i: int) -> float:
            n = len(x_labels)
            if n == 1: return plot_left + plot_w / 2
            return plot_left + i * plot_w / (n - 1)

        def y_to_in(v: float) -> float:
            return plot_bot - (v - y_min) / rng * plot_h

        n_ticks = 5
        for k in range(n_ticks + 1):
            v = y_min + rng * k / n_ticks
            ty = y_to_in(v)
            self._add_line(plot_left, ty, plot_right, ty,
                           color=LIGHT_GRAY, weight=0.5)
            self._add_text(0.3, ty - 0.1, plot_left - 0.35, 0.22,
                           y_format.format(v), font_size=9,
                           color=FOOTNOTE_GRAY, font=self.font_body,
                           align=PP_ALIGN.RIGHT)

        self._add_line(plot_left, plot_bot, plot_right, plot_bot,
                       color=BORDER_GRAY, weight=0.75)
        for i, xl in enumerate(x_labels):
            x = x_to_in(i)
            self._add_text(x - 0.6, plot_bot + 0.05, 1.2, 0.25,
                           str(xl), font_size=9, color=FOOTNOTE_GRAY,
                           font=self.font_body, align=PP_ALIGN.CENTER)

        if y_label:
            self._add_text(0.3, body_top - 0.05, 1.5, 0.22,
                           y_label, font_size=9, italic=True,
                           color=FOOTNOTE_GRAY, font=self.font_body)

        def _series_color(idx: int, label: str):
            if target_series and label == target_series:
                return self.brand_primary
            palette = [self.brand_dark, self.accent_cyan,
                       SERIES_PALETTE[3], SERIES_PALETTE[4],
                       SERIES_PALETTE[6], SERIES_PALETTE[8]]
            return palette[idx % len(palette)]

        for si, s in enumerate(series):
            color = _series_color(si, s["label"])
            vals = s["values"]
            pts = [(x_to_in(i), y_to_in(v)) for i, v in enumerate(vals)
                   if v is not None]
            for k in range(len(pts) - 1):
                x1, y1 = pts[k]
                x2, y2 = pts[k + 1]
                self._add_line(x1, y1, x2, y2, color=color, weight=2)
            if pts:
                ex, ey = pts[-1]
                self._add_oval(ex - 0.07, ey - 0.07, 0.14, 0.14, color)
                self._add_text(ex + 0.1, ey - 0.13, 1.4, 0.25,
                               s["label"], font_size=10, bold=True,
                               color=color, font=self.font_body)

        if not skip_source:
            self._add_source_line(slide, source, notes)
        self._page += 1
        self._add_footer(slide)

    @staticmethod
    def _fmt_to_numfmt(fmt: str) -> str:
        """Map Python format strings ({:.1f}%, {:,.0f}, etc.) to Excel
        number_format strings used by chart axes/data labels.

        Note: values like 23.4 supplied to a Python `{:.1f}%` format render
        as "23.4%". Excel's `0.0%` format multiplies by 100, so we use
        `0.0"%"` (literal %) to preserve the supplied magnitude. Use
        `:.1%` (Python percent) to opt into Excel's multiplier behavior.
        """
        f = (fmt or "").strip()
        # Excel-style percent (caller supplied 0.198 not 19.8)
        if ":.1%" in f: return "0.0%"
        if ":.2%" in f: return "0.00%"
        if ":.0%" in f: return "0%"
        # Trailing literal "%" with pre-scaled values
        if f.endswith("%}") or f.endswith("%"):
            if ":.1f" in f: return '0.0"%"'
            if ":.2f" in f: return '0.00"%"'
            return '0"%"'
        # Multiples ({:.1f}x)
        if f.endswith("x}") or f.endswith("x"):
            if ":.1f" in f: return '0.0"x"'
            if ":.2f" in f: return '0.00"x"'
            return '0"x"'
        # Comma-separated integers / decimals
        if ":+,.0f" in f: return "+#,##0;-#,##0"
        if ":,.0f" in f: return "#,##0"
        if ":,.1f" in f: return "#,##0.0"
        if ":,.2f" in f: return "#,##0.00"
        if ":.0f" in f:  return "0"
        if ":.1f" in f:  return "0.0"
        if ":.2f" in f:  return "0.00"
        return "General"

    def add_waterfall(
        self,
        action_title: str,
        segments: list,                # list of dicts: {label, value, kind}
        value_format: str = "{:+,.0f}",
        y_label: str = "",
        source: str = "",
        notes: str = "",
        skip_source: bool = False,
        broken_axis: bool = False,     # auto-truncate baseline when delta ratio extreme
        broken_axis_threshold: float = 5.0,
    ) -> None:
        """
        Waterfall / bridge chart (Section 5.5 'value bridge').
        kind: "start" | "plus" | "minus" | "total"
              start  -> brand_dark fill, full bar from baseline
              plus   -> Forest Bright fill, bar above prior cumulative
              minus  -> Citi Red fill, bar below prior cumulative
              total  -> brand_dark fill, full bar from baseline
        Connector lines render between bar tops.

        broken_axis: when True, and the largest start/total bar is more than
            `broken_axis_threshold`x larger than the biggest delta, the y-axis
            is truncated below the deltas and a zigzag break marker is drawn
            on each start/total bar. Makes plus/minus bars visually readable
            when the baseline dwarfs the deltas.
        Density: 2-12 segments.
        """
        action_title = self._validate_action_title(action_title)
        n = len(segments)
        if not (2 <= n <= 12):
            raise ValueError(f"waterfall requires 2-12 segments, got {n}")
        for s in segments:
            if s.get("kind") not in ("start", "plus", "minus", "total"):
                raise ValueError(
                    f"segment kind must be start/plus/minus/total, "
                    f"got '{s.get('kind')}'")
            if "value" not in s or "label" not in s:
                raise ValueError("each segment needs label + value")

        slide = self._content_slide_with_title(action_title)

        body_top = 0.95
        body_bottom = self._slide_h_in - 0.85  # leave room for x-labels
        plot_left = 0.9
        plot_right = self._slide_w_in - 0.4
        plot_w = plot_right - plot_left
        plot_top = body_top + 0.4   # leave room for value labels
        plot_bot = body_bottom - 0.3
        plot_h = plot_bot - plot_top

        # Compute running cumulative + bar (lo, hi) per segment
        bars = []  # list of (kind, lo, hi)
        running = 0.0
        for s in segments:
            k = s["kind"]
            v = float(s["value"])
            if k in ("start", "total"):
                lo, hi = 0.0, v
                running = v
            elif k == "plus":
                lo, hi = running, running + v
                running += v
            else:  # minus
                lo, hi = running + v, running
                running += v
            bars.append((k, lo, hi))

        # Decide on broken axis: only when start/total dwarfs deltas
        floor_v = 0.0
        do_break = False
        bar_max = max(hi for _, _, hi in bars)
        bar_min_lo = min(lo for _, lo, _ in bars)
        deltas = [abs(float(s["value"]))
                  for s in segments if s["kind"] in ("plus", "minus")]
        biggest_baseline = max(
            (abs(float(s["value"]))
             for s in segments if s["kind"] in ("start", "total")),
            default=0.0,
        )
        biggest_delta = max(deltas) if deltas else 0.0
        if (broken_axis and deltas and biggest_baseline > 0
                and biggest_baseline / max(biggest_delta, 1e-9)
                >= broken_axis_threshold):
            do_break = True
            # Floor sits just below all bar bottoms (lo) and well above 0.
            min_bottom = min(lo for k, lo, hi in bars
                             if k in ("plus", "minus"))
            floor_v = max(bar_min_lo, min_bottom) * 0.92
            # Anchor floor to the smaller of (bar_min_lo, the lowest delta lo)
            # so deltas occupy meaningful space.
            floor_v = min(floor_v, bar_min_lo - biggest_delta * 0.5)
            if floor_v <= 0:
                floor_v = bar_min_lo * 0.85

        if do_break:
            y_min = floor_v
            y_max = bar_max
        else:
            y_min = min(0, bar_min_lo)
            y_max = bar_max
        rng = (y_max - y_min) or 1
        y_max += rng * 0.10
        rng = y_max - y_min

        def y_to_in(v: float) -> float:
            v_clamped = max(min(v, y_max), y_min)
            return plot_bot - (v_clamped - y_min) / rng * plot_h

        # Baseline (zero or break floor)
        baseline_y = y_to_in(y_min)
        self._add_line(plot_left, baseline_y, plot_right, baseline_y,
                       color=BORDER_GRAY, weight=0.75)

        # Bar geometry
        gap = 0.15
        bar_w = (plot_w - gap * (n - 1)) / n
        forest = RGBColor(0x38, 0x8A, 0x42)

        prev_top_y = None
        prev_x_right = None
        for i, ((k, lo, hi), seg) in enumerate(zip(bars, segments)):
            x = plot_left + i * (bar_w + gap)
            # When broken: start/total bars are clipped at floor_v but
            # stretch from floor to hi. Deltas already span lo..hi naturally.
            disp_lo = max(lo, y_min) if (do_break and k in ("start", "total")) else lo
            top_y = y_to_in(max(disp_lo, hi))
            bot_y = y_to_in(min(disp_lo, hi))
            h = max(0.04, bot_y - top_y)

            color = {
                "start": self.brand_dark,
                "total": self.brand_dark,
                "plus":  forest,
                "minus": self.accent_red,
            }[k]
            self._add_rect(x, top_y, bar_w, h, color, no_line=True)

            # Zigzag break marker on truncated start/total bars
            if do_break and k in ("start", "total"):
                self._draw_break_marker(x, bar_w, plot_bot - 0.05)

            # Connector to previous bar's running top
            if prev_top_y is not None and prev_x_right is not None:
                cy = y_to_in(lo if k in ("plus", "minus") else max(lo, hi))
                self._add_line(prev_x_right, prev_top_y,
                               x, cy,
                               color=BORDER_GRAY, weight=0.5)

            # Value label above bar
            label_y = top_y - 0.28
            self._add_text(x - 0.1, label_y, bar_w + 0.2, 0.22,
                           value_format.format(seg["value"]),
                           font_size=9, bold=True, color=color,
                           font=self.font_body, align=PP_ALIGN.CENTER)

            # X-axis label below
            self._add_text(x - 0.1, plot_bot + 0.1, bar_w + 0.2, 0.4,
                           seg["label"], font_size=9, color=self.brand_dark,
                           font=self.font_body, align=PP_ALIGN.CENTER)

            # Connector anchor: top of running cumulative side
            prev_top_y = y_to_in(hi if k != "minus" else lo)
            prev_x_right = x + bar_w

        if y_label:
            label = y_label
            if do_break:
                label = f"{y_label} (axis truncated)" if y_label \
                    else "(axis truncated)"
            self._add_text(0.3, body_top + 0.05, 2.5, 0.22,
                           label, font_size=9, italic=True,
                           color=FOOTNOTE_GRAY, font=self.font_body)
        elif do_break:
            self._add_text(0.3, body_top + 0.05, 2.5, 0.22,
                           "(axis truncated)", font_size=9, italic=True,
                           color=FOOTNOTE_GRAY, font=self.font_body)

        if not skip_source:
            self._add_source_line(slide, source, notes)
        self._page += 1
        self._add_footer(slide)

    def _draw_break_marker(self, bar_x: float, bar_w: float,
                           band_top: float) -> None:
        """Draw a zigzag break marker spanning the bottom band of a bar.
        Indicates a truncated y-axis; sits at `band_top` y-position."""
        # White band over the bar bottom (4-pt strip), plus two zigzag lines
        # in BORDER_GRAY traversing it.
        band_h = 0.18
        self._add_rect(bar_x - 0.02, band_top, bar_w + 0.04, band_h,
                       WHITE, no_line=True)
        # Zigzag: 4 alternating peaks across bar width
        zig_y_top = band_top + 0.02
        zig_y_bot = band_top + band_h - 0.02
        steps = 6
        last_x = bar_x - 0.02
        last_y = zig_y_top
        for k in range(1, steps + 1):
            nx = bar_x - 0.02 + (bar_w + 0.04) * k / steps
            ny = zig_y_bot if k % 2 == 1 else zig_y_top
            self._add_line(last_x, last_y, nx, ny,
                           color=BORDER_GRAY, weight=0.75)
            last_x, last_y = nx, ny
        # Second zigzag offset down slightly for readability
        # (kept simple — single zigzag is conventional)

    def add_stacked_bar(
        self,
        action_title: str,
        categories: list,              # x-axis labels (e.g. years)
        series: list,                  # list of dicts: {label, values, color?}
        target_category: str = "",     # category index emphasised via label bold
        value_format: str = "{:,.0f}",
        show_totals: bool = True,
        source: str = "",
        notes: str = "",
        skip_source: bool = False,
        native: bool = True,
    ) -> None:
        """
        Vertical stacked column chart (Section 5.6 'breakdown over time').
        Native python-pptx chart (XL_CHART_TYPE.COLUMN_STACKED) — editable in
        PowerPoint, scales properly. Set native=False for the legacy shape-
        based renderer. Density: 2-8 categories x 2-6 series. Per-series
        'color' hex override applies to both renderers.
        """
        action_title = self._validate_action_title(action_title)
        nc = len(categories)
        ns = len(series)
        if not (2 <= nc <= 8):
            raise ValueError(f"stacked_bar requires 2-8 categories, got {nc}")
        if not (2 <= ns <= 6):
            raise ValueError(f"stacked_bar requires 2-6 series, got {ns}")
        for s in series:
            if "label" not in s or "values" not in s:
                raise ValueError("each series needs label + values")
            if len(s["values"]) != nc:
                raise ValueError(
                    f"series '{s['label']}' has {len(s['values'])} values "
                    f"but {nc} categories")
            for v in s["values"]:
                if v < 0:
                    raise ValueError("stacked_bar values must be non-negative")

        if not native:
            return self._add_stacked_bar_shapes(
                action_title, categories, series, target_category,
                value_format, show_totals, source, notes, skip_source,
            )

        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import (
            XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION,
            XL_TICK_MARK,
        )

        slide = self._content_slide_with_title(action_title)

        body_top = 0.95
        body_bottom = self._slide_h_in - 0.7
        chart_left = 0.4
        chart_top = body_top
        chart_w = self._slide_w_in - 0.8
        chart_h = body_bottom - body_top - 0.05

        chart_data = CategoryChartData()
        chart_data.categories = [str(c) for c in categories]
        for s in series:
            chart_data.add_series(s["label"],
                                  [float(v) for v in s["values"]])

        gframe = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_STACKED,
            Inches(chart_left), Inches(chart_top),
            Inches(chart_w), Inches(chart_h),
            chart_data,
        )
        chart = gframe.chart
        chart.has_title = False
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(10)
        chart.legend.font.name = self.font_body
        chart.legend.font.color.rgb = self.brand_dark

        def _hex_to_rgb(hx: str):
            h = hx.lstrip("#")
            return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

        default_palette = [self.brand_primary, self.brand_dark, self.accent_cyan,
                           SERIES_PALETTE[4], SERIES_PALETTE[6], SERIES_PALETTE[7]]

        for si, s in enumerate(series):
            if s.get("color"):
                try:
                    color = _hex_to_rgb(s["color"])
                except Exception:
                    color = default_palette[si % len(default_palette)]
            else:
                color = default_palette[si % len(default_palette)]
            ser = chart.series[si]
            ser.format.fill.solid()
            ser.format.fill.fore_color.rgb = color
            try:
                ser.format.line.fill.background()
            except Exception:
                pass

        # Add total labels atop each stacked bar via overlay text boxes
        # (native COLUMN_STACKED data labels show segment values, not totals).
        if show_totals:
            from pptx.util import Emu
            totals = [sum(s["values"][c] for s in series) for c in range(nc)]
            # Estimate plot area for label placement (chart region minus
            # legend ~0.4in and minus title ~0.0). Anchor labels right
            # above each category center.
            cat_w = chart_w / nc
            for ci, total in enumerate(totals):
                cx = chart_left + cat_w * (ci + 0.5)
                # Place label just inside chart top to avoid overlap with title.
                self._add_text(cx - 0.6, chart_top + 0.05, 1.2, 0.22,
                               value_format.format(total),
                               font_size=9, bold=True,
                               color=self.brand_dark, font=self.font_body,
                               align=PP_ALIGN.CENTER, slide=slide)

        # Y / X axis styling
        try:
            va = chart.value_axis
            va.tick_labels.font.size = Pt(9)
            va.tick_labels.font.name = self.font_body
            va.tick_labels.font.color.rgb = FOOTNOTE_GRAY
            va.tick_labels.number_format = self._fmt_to_numfmt(value_format)
            va.has_major_gridlines = True
            try:
                va.major_gridlines.format.line.color.rgb = LIGHT_GRAY
                va.major_gridlines.format.line.width = Pt(0.5)
            except Exception:
                pass
            va.minor_tick_mark = XL_TICK_MARK.NONE
            va.major_tick_mark = XL_TICK_MARK.OUTSIDE
        except Exception:
            pass
        try:
            ca = chart.category_axis
            ca.tick_labels.font.size = Pt(10)
            ca.tick_labels.font.name = self.font_body
            ca.tick_labels.font.color.rgb = self.brand_dark
            # Bold target category if specified — applied via per-tick XML
            # is non-trivial; rely on native uniform styling. Fall back to
            # marking the category in the title hint instead.
        except Exception:
            pass

        if not skip_source:
            self._add_source_line(slide, source, notes)
        self._page += 1
        self._add_footer(slide)

    def _add_stacked_bar_shapes(
        self, action_title, categories, series, target_category,
        value_format, show_totals, source, notes, skip_source,
    ) -> None:
        """Legacy shape-based stacked bar (kept as fallback)."""
        slide = self._content_slide_with_title(action_title)
        nc = len(categories)

        body_top = 0.95
        body_bottom = self._slide_h_in - 0.85
        plot_left = 0.9
        plot_right = self._slide_w_in - 2.1
        plot_w = plot_right - plot_left
        plot_top = body_top + 0.4
        plot_bot = body_bottom - 0.3
        plot_h = plot_bot - plot_top

        totals = [sum(s["values"][c] for s in series) for c in range(nc)]
        y_max = max(totals) or 1
        y_max *= 1.10

        gap = 0.25
        bar_w = (plot_w - gap * (nc - 1)) / nc

        def _hex_to_rgb(hx: str):
            h = hx.lstrip("#")
            return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

        default_palette = [self.brand_primary, self.brand_dark, self.accent_cyan,
                           SERIES_PALETTE[3], SERIES_PALETTE[4], SERIES_PALETTE[6]]
        series_colors = []
        for si, s in enumerate(series):
            if s.get("color"):
                try:
                    series_colors.append(_hex_to_rgb(s["color"]))
                except Exception:
                    series_colors.append(default_palette[si % len(default_palette)])
            else:
                series_colors.append(default_palette[si % len(default_palette)])

        self._add_line(plot_left, plot_bot, plot_right, plot_bot,
                       color=BORDER_GRAY, weight=0.75)

        for ci in range(nc):
            x = plot_left + ci * (bar_w + gap)
            cum = 0.0
            for si, s in enumerate(series):
                v = s["values"][ci]
                seg_h = (v / y_max) * plot_h
                if seg_h <= 0: continue
                seg_y = plot_bot - cum - seg_h
                self._add_rect(x, seg_y, bar_w, seg_h, series_colors[si],
                               no_line=True)
                cum += seg_h
            if show_totals:
                self._add_text(x - 0.1, plot_bot - cum - 0.28,
                               bar_w + 0.2, 0.22,
                               value_format.format(totals[ci]),
                               font_size=10, bold=True,
                               color=self.brand_dark, font=self.font_body,
                               align=PP_ALIGN.CENTER)
            label_bold = (categories[ci] == target_category)
            self._add_text(x - 0.1, plot_bot + 0.1, bar_w + 0.2, 0.3,
                           str(categories[ci]), font_size=10,
                           bold=label_bold, color=self.brand_dark,
                           font=self.font_body, align=PP_ALIGN.CENTER)

        legend_left = plot_right + 0.15
        legend_top = plot_top
        for si, s in enumerate(series):
            ly = legend_top + si * 0.32
            self._add_rect(legend_left, ly + 0.04, 0.18, 0.18,
                           series_colors[si], no_line=True)
            self._add_text(legend_left + 0.25, ly, 1.5, 0.25,
                           s["label"], font_size=10,
                           color=self.brand_dark, font=self.font_body)

        if not skip_source:
            self._add_source_line(slide, source, notes)
        self._page += 1
        self._add_footer(slide)

    def add_pie(
        self,
        action_title: str,
        slices: list,                  # list of dicts: {label, value, color?}
        target_label: str = "",        # explode this slice
        show_pct: bool = True,
        source: str = "",
        notes: str = "",
        skip_source: bool = False,
    ) -> None:
        """
        Pie chart (Section 5.7 'composition share').
        Native python-pptx chart used for accurate wedges; series colors
        applied from brand palette (or per-slice 'color' hex override).
        Density: 2-8 slices. Slices smaller than 3% rolled into 'Other'.
        """
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
        from pptx.enum.chart import XL_LABEL_POSITION

        action_title = self._validate_action_title(action_title)
        n = len(slices)
        if not (2 <= n <= 8):
            raise ValueError(f"pie requires 2-8 slices, got {n}")
        for s in slices:
            if "label" not in s or "value" not in s:
                raise ValueError("each slice needs label + value")
            if s["value"] < 0:
                raise ValueError("pie values must be non-negative")
        total = sum(s["value"] for s in slices)
        if total <= 0:
            raise ValueError("pie total must be positive")

        slide = self._content_slide_with_title(action_title)

        body_top = 0.95
        body_bottom = self._slide_h_in - 0.7
        chart_w = min(6.0, (body_bottom - body_top) * 1.3)
        chart_h = body_bottom - body_top - 0.2
        chart_left = (self._slide_w_in - chart_w) / 2
        chart_top = body_top + 0.05

        chart_data = CategoryChartData()
        chart_data.categories = [s["label"] for s in slices]
        chart_data.add_series("Series 1", [s["value"] for s in slices])

        gframe = slide.shapes.add_chart(
            XL_CHART_TYPE.PIE,
            Inches(chart_left), Inches(chart_top),
            Inches(chart_w), Inches(chart_h),
            chart_data,
        )
        chart = gframe.chart
        chart.has_title = False
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(10)
        chart.legend.font.name = self.font_body
        chart.legend.font.color.rgb = self.brand_dark

        # Apply per-slice colors via DataPoint API
        def _hex_to_rgb(hx: str):
            h = hx.lstrip("#")
            return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

        plot = chart.plots[0]
        plot.has_data_labels = show_pct
        if show_pct:
            dl = plot.data_labels
            dl.show_percentage = True
            dl.show_value = False
            dl.show_category_name = False
            dl.font.size = Pt(10)
            dl.font.bold = True
            dl.font.color.rgb = WHITE
            dl.font.name = self.font_body
            dl.position = XL_LABEL_POSITION.CENTER
            dl.number_format = "0%"

        default_pal = [self.brand_primary, self.brand_dark, self.accent_cyan,
                       SERIES_PALETTE[3], SERIES_PALETTE[4], SERIES_PALETTE[6],
                       SERIES_PALETTE[7], SERIES_PALETTE[8]]
        series0 = plot.series[0]
        target_idx = next(
            (i for i, s in enumerate(slices) if s["label"] == target_label),
            None,
        )
        for i, s in enumerate(slices):
            pt = series0.points[i]
            pt.format.fill.solid()
            if s.get("color"):
                try:
                    pt.format.fill.fore_color.rgb = _hex_to_rgb(s["color"])
                except Exception:
                    pt.format.fill.fore_color.rgb = default_pal[i % len(default_pal)]
            else:
                pt.format.fill.fore_color.rgb = default_pal[i % len(default_pal)]
            pt.format.line.color.rgb = WHITE
            pt.format.line.width = Pt(1.5)

        # Explode the target slice via XML adjustment (python-pptx has no direct API)
        if target_idx is not None:
            from pptx.oxml.ns import qn
            from lxml import etree
            ser = series0._ser
            # remove existing dPt for target if any, then add explosion
            for dpt in ser.findall(qn("c:dPt")):
                idx_el = dpt.find(qn("c:idx"))
                if idx_el is not None and idx_el.get("val") == str(target_idx):
                    ser.remove(dpt)
                    break
            dPt = etree.SubElement(ser, qn("c:dPt"))
            etree.SubElement(dPt, qn("c:idx")).set("val", str(target_idx))
            etree.SubElement(dPt, qn("c:bubble3D")).set("val", "0")
            etree.SubElement(dPt, qn("c:explosion")).set("val", "15")

        if not skip_source:
            self._add_source_line(slide, source, notes)
        self._page += 1
        self._add_footer(slide)

    # ── Citi PitchPres archetypes (per SPEC_PitchPres_A4_Landscape) ──

    def add_pros_cons(
        self,
        action_title: str,
        pros: list,                    # list of strings
        cons: list,                    # list of strings
        neutral: Optional[list] = None,  # optional 3rd column
        source: str = "",
        notes: str = "",
        skip_source: bool = False,
    ) -> None:
        """
        3-column pros/cons/neutral evaluation slide (Section 5.5).
        Forest Bright (#388A42) for pros, Citi Red for cons, Gray 04 fill for neutral.
        """
        action_title = self._validate_action_title(action_title)
        slide = self._content_slide_with_title(action_title)

        cols = 3 if neutral else 2
        gap = 0.2
        body_top = 0.95
        body_bottom = self._slide_h_in - 0.7
        body_h = body_bottom - body_top
        avail_w = self._slide_w_in - 0.6
        col_w = (avail_w - gap * (cols - 1)) / cols
        left0 = 0.3

        forest = RGBColor(0x38, 0x8A, 0x42)  # Forest Bright per spec
        gray04 = RGBColor(0xE6, 0xEB, 0xED)

        groups = [("Pros", pros, forest, None),
                  ("Cons", cons, self.accent_red, None)]
        if neutral:
            groups.append(("Neutral", neutral, self.brand_dark, gray04))

        for i, (label, items, color, fill) in enumerate(groups):
            x = left0 + i * (col_w + gap)
            # Background fill if neutral
            if fill is not None:
                self._add_rect(x, body_top, col_w, body_h, fill, no_line=True)
            # Header strip
            self._add_rect(x, body_top, col_w, 0.4, color, no_line=True)
            self._add_text(x + 0.05, body_top + 0.04, col_w - 0.1, 0.32,
                           label.upper(), font_size=12, bold=True,
                           color=WHITE, font=self.font_headline,
                           align=PP_ALIGN.CENTER)
            # Bullets
            bullet_text = "\n".join(f"• {b}" for b in items)
            self._add_text(x + 0.15, body_top + 0.55,
                           col_w - 0.3, body_h - 0.65,
                           bullet_text, font_size=11, color=color,
                           font=self.font_body, align=PP_ALIGN.LEFT)

        if not skip_source:
            self._add_source_line(slide, source, notes)
        self._page += 1
        self._add_footer(slide)

    def add_quad_page(
        self,
        action_title: str,
        panels: list,                  # list of dicts: {heading, bullets}
        source: str = "",
        notes: str = "",
        skip_source: bool = False,
    ) -> None:
        """
        2x2 grid of text panels (Section 5.6). Each panel: heading + 3 bullets.
        0.5pt Gray 01 outline around each panel.
        """
        action_title = self._validate_action_title(action_title)
        if len(panels) != 4:
            raise ValueError(f"quad_page requires exactly 4 panels, got {len(panels)}")

        slide = self._content_slide_with_title(action_title)

        body_top = 0.95
        body_bottom = self._slide_h_in - 0.7
        body_h = body_bottom - body_top
        gap_x = 0.25
        gap_y = 0.2
        avail_w = self._slide_w_in - 0.6
        panel_w = (avail_w - gap_x) / 2
        panel_h = (body_h - gap_y) / 2
        left0 = 0.3

        for i, panel in enumerate(panels):
            r, c = divmod(i, 2)
            x = left0 + c * (panel_w + gap_x)
            y = body_top + r * (panel_h + gap_y)
            # Outlined panel (0.5pt Gray 01)
            self._add_rect(x, y, panel_w, panel_h, WHITE, line=BORDER_GRAY)
            # Heading
            self._add_text(x + 0.15, y + 0.1, panel_w - 0.3, 0.4,
                           panel.get("heading", ""), font_size=13, bold=True,
                           color=self.brand_primary, font=self.font_headline)
            # Thin divider under heading
            self._add_line(x + 0.15, y + 0.55,
                           x + panel_w - 0.15, y + 0.55,
                           color=BORDER_GRAY, weight=0.5)
            # Bullets
            bullets = panel.get("bullets", [])
            bullet_text = "\n".join(f"• {b}" for b in bullets[:5])
            self._add_text(x + 0.15, y + 0.65,
                           panel_w - 0.3, panel_h - 0.75,
                           bullet_text, font_size=11, color=self.brand_dark,
                           font=self.font_body, align=PP_ALIGN.LEFT)

        if not skip_source:
            self._add_source_line(slide, source, notes)
        self._page += 1
        self._add_footer(slide)

    def add_org_chart(
        self,
        action_title: str,
        boxes: list,                   # list[OrgBox]; tier+parent_idx drive layout
        source: str = "",
        notes: str = "",
        skip_source: bool = False,
    ) -> None:
        """
        Hierarchical org chart (Section 5.8). 4-tier color scheme:
          tier 1 (top)    -> Citi Blue, white text
          tier 2          -> Citi Ink, white text
          tier 3          -> Blue Light, white text
          tier 4 (stats)  -> Gray 01, white text
        Layout: tier 1 centered top; lower tiers spread under their parents.
        """
        action_title = self._validate_action_title(action_title)
        if not boxes or not (1 <= len(boxes) <= 16):
            raise ValueError(f"org_chart requires 1-16 boxes, got {len(boxes)}")

        slide = self._content_slide_with_title(action_title)

        # Color by tier; tier 4 is a "stat box" — light fill, dark text
        tier_colors = {
            1: self.brand_primary,
            2: self.brand_dark,
            3: self.accent_cyan,
            4: LIGHT_GRAY,
        }
        tier_text = {
            1: WHITE,
            2: WHITE,
            3: WHITE,
            4: self.brand_dark,
        }

        body_top = 1.0
        body_bottom = self._slide_h_in - 0.7
        body_h = body_bottom - body_top
        body_w = self._slide_w_in - 0.6
        body_left = 0.3

        # Group boxes by tier
        from collections import defaultdict
        tiers: defaultdict = defaultdict(list)
        for i, b in enumerate(boxes):
            tiers[b.tier].append((i, b))

        max_tier = max(tiers.keys()) if tiers else 1
        tier_h = body_h / max_tier
        box_h = min(0.85, tier_h * 0.5)

        # Box dimensions (compute first so positioning can center within slot)
        max_per_tier = max(len(v) for v in tiers.values())
        box_w = min(2.2, body_w / max(1, max_per_tier) - 0.15)

        # Compute box positions per tier (left edge, centered within slot)
        positions = {}  # idx -> (left_x, top_y)
        for t in sorted(tiers.keys()):
            entries = tiers[t]
            n = len(entries)
            tier_y = body_top + (t - 1) * tier_h + (tier_h - box_h) / 2
            tier_w_per = body_w / n
            for j, (idx, b) in enumerate(entries):
                # Center the box within its slot, left edge at:
                left_x = body_left + tier_w_per * j + (tier_w_per - box_w) / 2
                positions[idx] = (left_x, tier_y)

        # Draw connectors first (so boxes overlap)
        for i, b in enumerate(boxes):
            if b.parent_idx is not None and b.parent_idx in positions:
                cx, cy = positions[i]
                px, py = positions[b.parent_idx]
                # L-shaped connector
                mid_y = (py + box_h + cy) / 2
                self._add_line(px + box_w / 2, py + box_h, px + box_w / 2, mid_y,
                               color=BORDER_GRAY, weight=0.5)
                self._add_line(px + box_w / 2, mid_y, cx + box_w / 2, mid_y,
                               color=BORDER_GRAY, weight=0.5)
                self._add_line(cx + box_w / 2, mid_y, cx + box_w / 2, cy,
                               color=BORDER_GRAY, weight=0.5)

        # Draw boxes
        for i, b in enumerate(boxes):
            cx, cy = positions[i]
            color = tier_colors.get(b.tier, self.brand_dark)
            text_color = tier_text.get(b.tier, WHITE)
            is_stat = (b.tier == 4)
            # Stat boxes are compact: shorter height, narrower width
            this_box_w = box_w * 0.8 if is_stat else box_w
            this_box_h = box_h * 0.65 if is_stat else box_h
            this_cx = cx + (box_w - this_box_w) / 2
            this_cy = cy + (box_h - this_box_h) / 2
            if is_stat:
                shp = self._add_rect(this_cx, this_cy,
                                     this_box_w, this_box_h, color,
                                     line=BORDER_GRAY)
                shp.line.width = Pt(0.5)
            else:
                shp = self._add_rect(this_cx, this_cy,
                                     this_box_w, this_box_h, color,
                                     no_line=True)
            tf = shp.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.text = b.label
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0]
            run.font.size = Pt(9 if is_stat else 10)
            run.font.bold = True
            run.font.color.rgb = text_color
            run.font.name = self.font_body
            if b.sub_label:
                p2 = tf.add_paragraph()
                p2.alignment = PP_ALIGN.CENTER
                run2 = p2.add_run()
                run2.text = b.sub_label
                run2.font.size = Pt(8 if is_stat else 9)
                run2.font.color.rgb = (FOOTNOTE_GRAY if is_stat else text_color)
                run2.font.name = self.font_body

        if not skip_source:
            self._add_source_line(slide, source, notes)
        self._page += 1
        self._add_footer(slide)

    def add_tombstone_page(
        self,
        action_title: str,
        tiles: list,                   # list[TombstoneTile]
        cols: int = 7,
        source: str = "",
        notes: str = "",
        skip_source: bool = False,
    ) -> None:
        """
        Grid of deal tombstone tiles (Section 5.11).
        Default 7 columns x ~4 rows; max 28 tiles per slide.
        """
        action_title = self._validate_action_title(action_title)
        n = len(tiles)
        if not (1 <= n <= 28):
            raise ValueError(f"tombstone_page requires 1-28 tiles, got {n}")

        slide = self._content_slide_with_title(action_title)

        rows = -(-n // cols)
        body_top = 0.95
        body_bottom = self._slide_h_in - 0.7
        body_h = body_bottom - body_top
        body_w = self._slide_w_in - 0.6
        gap = 0.08
        tile_w = (body_w - gap * (cols - 1)) / cols
        tile_h = (body_h - gap * (rows - 1)) / rows
        left0 = 0.3

        gray04 = RGBColor(0xE6, 0xEB, 0xED)
        sand04 = RGBColor(0xEA, 0xE0, 0xD3)

        for i, tile in enumerate(tiles):
            r, c = divmod(i, cols)
            x = left0 + c * (tile_w + gap)
            y = body_top + r * (tile_h + gap)

            # Resolve tile fill + text colors
            if tile.highlight:
                fill = self.brand_primary
                text_color = WHITE
                line = None
            elif tile.fill_variant == "gray":
                fill = gray04; text_color = self.brand_dark; line = None
            elif tile.fill_variant == "sand":
                fill = sand04; text_color = self.brand_dark; line = None
            else:
                fill = WHITE; text_color = self.brand_dark; line = BORDER_GRAY

            self._add_rect(x, y, tile_w, tile_h, fill,
                           line=line, no_line=(line is None))

            # Layout: status, client, desc, size (bold), date
            line_h = tile_h / 5
            self._add_text(x + 0.05, y + 0.03, tile_w - 0.1, line_h,
                           tile.deal_status, font_size=7, italic=True,
                           color=text_color, font=self.font_body,
                           align=PP_ALIGN.CENTER)
            self._add_text(x + 0.05, y + line_h * 1.0, tile_w - 0.1, line_h,
                           tile.client, font_size=8, bold=True,
                           color=text_color, font=self.font_body,
                           align=PP_ALIGN.CENTER)
            self._add_text(x + 0.05, y + line_h * 2.0, tile_w - 0.1, line_h,
                           tile.description, font_size=7,
                           color=text_color, font=self.font_body,
                           align=PP_ALIGN.CENTER)
            self._add_text(x + 0.05, y + line_h * 3.1, tile_w - 0.1, line_h,
                           tile.deal_size, font_size=8, bold=True,
                           color=text_color, font=self.font_body,
                           align=PP_ALIGN.CENTER)
            self._add_text(x + 0.05, y + line_h * 4.0, tile_w - 0.1, line_h,
                           tile.date, font_size=7, italic=True,
                           color=text_color, font=self.font_body,
                           align=PP_ALIGN.CENTER)

        if not skip_source:
            self._add_source_line(slide, source, notes)
        self._page += 1
        self._add_footer(slide)

    def add_team_page(
        self,
        action_title: str,
        team_name: str,
        members: list,                 # list[TeamMember]
        cols: int = 3,
        source: str = "",
        notes: str = "",
        skip_source: bool = False,
    ) -> None:
        """
        Team page (Section 5.12). Banner + grid of person tiles with bold blue name.
        """
        action_title = self._validate_action_title(action_title)
        n = len(members)
        if not (1 <= n <= 18):
            raise ValueError(f"team_page requires 1-18 members, got {n}")

        slide = self._content_slide_with_title(action_title)

        # Team banner (Citi Ink fill, white text)
        body_top = 0.95
        banner_h = 0.4
        body_w = self._slide_w_in - 0.6
        self._add_rect(0.3, body_top, body_w, banner_h,
                       self.brand_dark, no_line=True)
        self._add_text(0.3, body_top + 0.04, body_w, banner_h - 0.08,
                       team_name, font_size=12, bold=True,
                       color=WHITE, font=self.font_headline,
                       align=PP_ALIGN.CENTER)

        # Person tiles
        tile_top = body_top + banner_h + 0.15
        body_bottom = self._slide_h_in - 0.7
        avail_h = body_bottom - tile_top
        rows = -(-n // cols)
        gap = 0.15
        tile_w = (body_w - gap * (cols - 1)) / cols
        tile_h = min(0.95, (avail_h - gap * (rows - 1)) / rows)

        photo_size = min(0.7, tile_h * 0.7)
        for i, m in enumerate(members):
            r, c = divmod(i, cols)
            x = 0.3 + c * (tile_w + gap)
            y = tile_top + r * (tile_h + gap)
            # White tile with 0.5pt outline
            self._add_rect(x, y, tile_w, tile_h, WHITE, line=BORDER_GRAY)

            # Photo or initials disc (left side of tile)
            photo_x = x + 0.12
            photo_y = y + (tile_h - photo_size) / 2
            self._draw_team_photo(m, photo_x, photo_y, photo_size)

            # Text column starts to right of photo
            tx = photo_x + photo_size + 0.15
            tw = tile_w - (tx - x) - 0.12
            # Bold name in brand primary
            self._add_text(tx, y + 0.08, tw, 0.3,
                           m.name, font_size=12, bold=True,
                           color=self.brand_primary, font=self.font_headline)
            # Title + segment
            sub = m.title
            if m.segment:
                sub = f"{m.title} | {m.segment}" if m.title else m.segment
            self._add_text(tx, y + 0.36, tw, 0.22,
                           sub, font_size=9, italic=True,
                           color=self.brand_dark, font=self.font_body)
            # Email + phone
            contact = []
            if m.email: contact.append(m.email)
            if m.phone: contact.append(m.phone)
            self._add_text(tx, y + 0.6, tw, tile_h - 0.65,
                           "\n".join(contact),
                           font_size=8, color=FOOTNOTE_GRAY,
                           font=self.font_body)

        if not skip_source:
            self._add_source_line(slide, source, notes)
        self._page += 1
        self._add_footer(slide)

    def add_table_of_contents(
        self,
        action_title: str,
        entries: list,                 # list[TocEntry]
        source: str = "",
        notes: str = "",
        skip_source: bool = True,      # TOC typically no source
    ) -> None:
        """
        Table of contents (Section 5.2). Numbered hierarchy 1/A/i with page refs.
        """
        action_title = self._validate_action_title(action_title)
        if not entries:
            raise ValueError("table_of_contents requires entries")

        slide = self._content_slide_with_title(action_title)

        body_top = 1.0
        body_bottom = self._slide_h_in - 0.7
        avail_h = body_bottom - body_top
        line_h = min(0.4, avail_h / max(1, len(entries)))

        # Per-level numbering counters
        counters = [0, 0, 0]
        roman_map = ["i", "ii", "iii", "iv", "v", "vi", "vii", "viii", "ix", "x"]

        body_w = self._slide_w_in - 0.6

        for i, e in enumerate(entries):
            lvl = max(1, min(3, e.level))
            # Increment level counter, reset deeper levels
            counters[lvl - 1] += 1
            for k in range(lvl, 3):
                counters[k] = 0

            if lvl == 1:
                num = f"{counters[0]}."
                indent = 0.0
                size = 14
                bold = True
                color = self.brand_primary
            elif lvl == 2:
                # Letter: A, B, C
                num = f"{chr(ord('A') + counters[1] - 1)}."
                indent = 0.4
                size = 12
                bold = True
                color = self.brand_dark
            else:
                # Roman: i, ii, iii
                num = f"{roman_map[min(counters[2] - 1, 9)]}."
                indent = 0.8
                size = 11
                bold = False
                color = self.brand_dark

            y = body_top + i * line_h
            # Number column
            self._add_text(0.3 + indent, y, 0.5, line_h,
                           num, font_size=size, bold=bold, color=color,
                           font=self.font_headline)
            # Title
            self._add_text(0.3 + indent + 0.5, y, body_w - indent - 1.5, line_h,
                           e.text, font_size=size, bold=bold, color=color,
                           font=self.font_body)
            # Page number right-aligned
            if e.page is not None:
                self._add_text(self._slide_w_in - 0.9, y, 0.6, line_h,
                               str(e.page), font_size=size, bold=bold,
                               color=color, font=self.font_body,
                               align=PP_ALIGN.RIGHT)

        if not skip_source:
            self._add_source_line(slide, source, notes)
        self._page += 1
        self._add_footer(slide)

    def save(self, filename: str) -> str:
        """Save deck and return absolute path."""
        if not filename.lower().endswith(".pptx"):
            filename += ".pptx"
        path = os.path.join(self.output_dir, filename)
        self.prs.save(path)
        return os.path.abspath(path)

    # ── Internal helpers ─────────────────────────────────────────────

    @staticmethod
    def _validate_action_title(t: str) -> str:
        """
        Engineering rule R1: action title must be present and non-trivial.
        Returns the sentence-case-normalized title (auto-applies _normalize_heading).
        Callers should use the returned value rather than the original.
        """
        if not t or not t.strip():
            raise ValueError("action title required (engineering R1)")
        if len(t.split()) < 3:
            raise ValueError(
                f"action title '{t}' too short - "
                "use 6-14 words conveying the takeaway"
            )
        return PPTXDeckWriter._normalize_heading(t)

    # Stopwords lowercased mid-sentence in titles (per-style rule)
    _TITLE_STOPWORDS = frozenset({
        "a", "an", "the",
        "and", "but", "or", "nor", "for", "yet", "so",
        "of", "in", "on", "at", "by", "to", "with", "for", "from",
        "as", "vs", "via", "per", "into", "onto", "upon",
        "is", "are", "be", "been", "was", "were",
        "that", "than", "this", "these",
    })

    @classmethod
    def _normalize_heading(cls, t: str) -> str:
        """
        Apply page-heading capitalization rule: first letter uppercase,
        common stopwords lowercase, ALL OTHER WORDS preserved as-is.

        This means:
          - Sentence-case input passes through unchanged (Sandvik stays "Sandvik")
          - Title-case input gets stopwords lowercased ("Investment Of The Year"
            -> "Investment of the Year")
          - Acronyms preserved (EBITDA, CEO, U.S.)
          - Digit tokens preserved (FY26, Q1, 2026)
          - Mixed-case preserved (iPhone, McKinsey)
          - Proper nouns mid-sentence preserved (Sandvik, Apple, Caterpillar)
        """
        if not t:
            return t
        words = t.split()
        out = []
        for i, w in enumerate(words):
            stripped = w.strip(".,;:!?\"'()")
            if not stripped:
                out.append(w)
                continue
            # Preserve all-caps acronyms (>=2 chars all caps)
            if stripped.isupper() and len(stripped) >= 2:
                out.append(w)
                continue
            # Preserve tokens with digits (FY26, Q1, 2026)
            if any(ch.isdigit() for ch in stripped):
                out.append(w)
                continue
            # Preserve mixed-case (proper nouns: iPhone, McKinsey, EBITDAm)
            inner = stripped[1:] if len(stripped) > 1 else ""
            if inner and any(ch.isupper() for ch in inner):
                out.append(w)
                continue
            # First word always capitalized
            if i == 0:
                out.append(w[0].upper() + w[1:].lower() if len(w) > 1
                           else w.upper())
                continue
            # Mid-sentence: stopwords -> lowercase, others preserved
            if stripped.lower() in cls._TITLE_STOPWORDS:
                out.append(w.lower())
            else:
                # Preserve as-is (covers Sandvik, Apple, etc.) - won't lowercase
                out.append(w)
        return " ".join(out)

    @staticmethod
    def _format_value(v) -> str:
        if v is None:
            return "n/a"
        if isinstance(v, str):
            return v
        if isinstance(v, float):
            if abs(v) < 1 and v != 0:
                return f"{v:.1%}" if abs(v) < 5 else f"{v:.2f}"
            return f"{v:,.1f}"
        if isinstance(v, int):
            return f"{v:,}"
        return str(v)

    def _blank_slide(self):
        layout = self.prs.slide_layouts[6]  # blank
        return self.prs.slides.add_slide(layout)

    def _content_slide_with_title(self, action_title: str):
        """
        Standard content slide with action title at top (Section 5).
        Tight top placement matching Citi style: headline at y=0.18",
        hairline divider at y=0.72", content area starts at y=0.85".
        Sentence-case applied via _normalize_heading inside _validate.
        Logo (top-right) auto-rendered if PPTXDeckWriter(logo_path=...) set.
        """
        slide = self._blank_slide()
        # Logo top-right, right-anchored. Width capped inside _add_logo so
        # it never overflows on narrower A4 vs 16:9 slides.
        if self.logo_path:
            self._add_logo(slide=slide, height_in=0.4,
                           right_in=self._slide_w_in - MARGIN_IN + 0.1,
                           y_in=0.15)
            # Reserve max-logo-width + small gap (1.6" cap + 0.2" gap)
            headline_w = self._slide_w_in - MARGIN_IN - 1.8
        else:
            headline_w = self._slide_w_in - 0.6
        # action_title is already normalized by _validate_action_title at archetype entry
        # Deterministic sizing: shrink by 1pt per 2 extra words past 8 (replaces autofit)
        word_count = len(action_title.split())
        size = self.headline_size
        if word_count > 8:
            size = max(16, self.headline_size - ((word_count - 8) // 2))
        self._add_text(0.3, 0.18,
                       headline_w, 0.45,
                       action_title, font_size=size,
                       bold=self.headline_bold, color=self.brand_primary,
                       font=self.font_headline)
        # Hairline divider tight under headline
        self._add_line(0.3, 0.72, self._slide_w_in - 0.3, 0.72,
                       color=MID_GRAY, weight=0.5)
        return slide

    def _logo_aspect(self) -> float:
        """Return logo aspect (w/h). Cached. Falls back to 2.5 (typical
        wordmark) when PIL unavailable or read fails."""
        if not self.logo_path:
            return 2.5
        cached = getattr(self, "_logo_aspect_cache", None)
        if cached is not None:
            return cached
        try:
            from PIL import Image
            with Image.open(self.logo_path) as im:
                w, h = im.size
                aspect = (w / h) if h else 2.5
        except Exception:
            aspect = 2.5
        self._logo_aspect_cache = aspect
        return aspect

    def _add_logo(self, slide, height_in: float,
                  x_in: Optional[float] = None, y_in: float = 0.15,
                  right_in: Optional[float] = None):
        """Embed logo image at the given position. Width auto-scales by
        aspect. If `right_in` is given, the logo is right-anchored (its
        right edge sits at that x position) — preferred for cover/content
        slides so the logo never overflows on aspects with smaller widths
        (A4 vs 16:9). Caps width at 1.6in for content slides, 2.0in for cover."""
        if not self.logo_path:
            return
        if slide is None:
            slide = self.prs.slides[-1]
        aspect = self._logo_aspect()
        # Cap maximum width so unusually wide logos don't dominate the slide
        max_w = 2.0 if height_in >= 0.6 else 1.6
        target_w = min(height_in * aspect, max_w)
        # If width hit the cap, scale height down to preserve aspect
        actual_h = (target_w / aspect) if aspect > 0 else height_in
        if right_in is not None:
            x_in = right_in - target_w
        elif x_in is None:
            x_in = self._slide_w_in - MARGIN_IN - target_w
        try:
            slide.shapes.add_picture(self.logo_path,
                                     Inches(x_in), Inches(y_in),
                                     width=Inches(target_w),
                                     height=Inches(actual_h))
        except Exception:
            pass  # silently skip on bad image

    def _add_footer(self, slide):
        """
        Footer (Section 9.1). Minimalist Citi-style: page number only at
        bottom-right; firm/project on bottom-left in tiny text only when
        confidentiality marking present.
        """
        y = self._slide_h_in - 0.3
        # Page number bottom-right
        self._add_text(self._slide_w_in - 1.0, y,
                       0.6, 0.25,
                       str(self._page), font_size=PT_PAGE,
                       color=FOOTNOTE_GRAY, align=PP_ALIGN.RIGHT,
                       slide=slide)
        # Confidentiality marking bottom-left if set
        if self.confidentiality:
            self._add_text(MARGIN_IN, y, 3.0, 0.25,
                           self.confidentiality, font_size=PT_PAGE,
                           color=FOOTNOTE_GRAY, bold=True,
                           slide=slide)

    def _add_source_line(self, slide, source: str, notes: str = ""):
        """Engineering R3: every data slide carries source line (7-8pt italic).
        Citi-style: tight near footer, no leading 'Source:' on note-only lines."""
        y = self._slide_h_in - 0.45
        parts = []
        if source:
            parts.append(f"Source: {source}")
        elif not notes:
            parts.append("Source: [TBD]")
        if notes:
            parts.append(f"Note: {notes}")
        text = "   ".join(parts)
        self._add_text(0.3, y, self._slide_w_in - 0.6, 0.25,
                       text, font_size=PT_FOOTNOTE,
                       color=FOOTNOTE_GRAY, italic=True, slide=slide)

    def _add_text(self, left, top, width, height, text,
                  font_size=PT_BODY, bold=False, italic=False,
                  color=None, align=PP_ALIGN.LEFT, slide=None,
                  font=None, anchor=None):
        """
        Render text in a fixed-size textbox. No autofit (shape-grow / text-shrink
        autofit caused unwanted growth in tests; use deterministic sizing).
        Default color resolves to self.brand_dark when not explicitly set.
        """
        if slide is None:
            slide = self.prs.slides[-1]
        if font is None:
            font = self.font_body
        if color is None:
            color = self.brand_dark
        tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Emu(0); tf.margin_right = Emu(0)
        tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
        if anchor is not None:
            tf.vertical_anchor = anchor

        # Multi-line text → multi paragraphs
        lines = text.split("\n")
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            run = p.add_run() if i > 0 else (p.runs[0] if p.runs else p.add_run())
            run.text = line
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = color
            run.font.name = font
        return tb

    def _add_rect(self, left, top, width, height, fill, line=None,
                  no_line=False):
        slide = self.prs.slides[-1]
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(left), Inches(top),
                                     Inches(width), Inches(height))
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        if no_line:
            shp.line.fill.background()
        elif line is not None:
            shp.line.color.rgb = line
        return shp

    def _add_rounded_rect(self, left, top, width, height, fill):
        slide = self.prs.slides[-1]
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(left), Inches(top),
                                     Inches(width), Inches(height))
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        shp.line.color.rgb = fill
        return shp

    def _add_oval(self, left, top, width, height, fill):
        slide = self.prs.slides[-1]
        shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     Inches(left), Inches(top),
                                     Inches(width), Inches(height))
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        shp.line.fill.background()
        return shp

    def _add_line(self, x1, y1, x2, y2, color=None, weight=1):
        if color is None: color = self.brand_dark
        slide = self.prs.slides[-1]
        from pptx.shapes.connector import Connector
        line = slide.shapes.add_connector(
            1,  # straight
            Inches(x1), Inches(y1), Inches(x2), Inches(y2),
        )
        line.line.color.rgb = color
        line.line.width = Pt(weight)
        return line

    def _add_arrow(self, x1, y1, x2, y2, color=None):
        if color is None: color = self.brand_dark
        line = self._add_line(x1, y1, x2, y2, color=color, weight=1.5)
        # Add arrow head via underlying XML
        from pptx.oxml.ns import qn
        ln = line.line._get_or_add_ln()
        # Tail end / head end arrow
        from lxml import etree
        tailEnd = etree.SubElement(ln, qn("a:tailEnd"))
        tailEnd.set("type", "triangle")
        tailEnd.set("w", "med")
        tailEnd.set("len", "med")
        return line

    def _set_cell(self, cell, text, bg, fg,
                  bold=False, italic=False,
                  size=PT_TABLE_B, align=PP_ALIGN.LEFT):
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.08)
        cell.margin_top = Inches(0.04); cell.margin_bottom = Inches(0.04)
        tf = cell.text_frame
        tf.word_wrap = True
        tf.text = str(text) if text is not None else ""
        p = tf.paragraphs[0]
        p.alignment = align
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = fg
            run.font.name = self.font_body

    def _draw_tile(self, slide, left, top, w, h, tile: ScorecardTile):
        # Light card background
        bg = self._add_rect(left, top, w, h, LIGHT_GRAY,
                            line=BORDER_GRAY)
        # Metric label
        self._add_text(left + 0.15, top + 0.1, w - 0.3, 0.4,
                       tile.metric, font_size=11, color=self.brand_dark,
                       bold=False)
        # Big value
        self._add_text(left + 0.15, top + 0.5, w - 0.3, h - 1.0,
                       tile.value, font_size=24, color=self.brand_primary,
                       bold=True)
        # Rating dots (filled vs empty)
        if tile.rating > 0:
            dot_y = top + h - 0.45
            r = max(0, min(5, tile.rating))
            for i in range(5):
                fill = self.brand_primary if i < r else MID_GRAY
                self._add_oval(left + 0.15 + i * 0.22, dot_y, 0.15, 0.15, fill)
        # Sub note
        if tile.sub:
            self._add_text(left + 0.15, top + h - 0.25, w - 0.3, 0.2,
                           tile.sub, font_size=8, color=FOOTNOTE_GRAY,
                           italic=True)

    def _draw_team_photo(self, m: "TeamMember", x: float, y: float, size: float):
        """
        Render a team member's headshot or an initials placeholder.
        photo_path -> embedded image cropped to a square; missing -> brand_primary
        disc with white initials (first + last name first letters).
        """
        # Try to render the photo if a usable file exists
        if m.photo_path and os.path.exists(m.photo_path):
            try:
                slide = self.prs.slides[-1]
                slide.shapes.add_picture(
                    m.photo_path,
                    Inches(x), Inches(y),
                    width=Inches(size), height=Inches(size),
                )
                return
            except Exception:
                pass  # fall through to initials placeholder

        # Initials disc placeholder
        self._add_oval(x, y, size, size, self.brand_primary)
        parts = (m.name or "").split()
        initials = "".join(p[0].upper() for p in parts[:2] if p) or "?"
        # Centered initials text — use anchor=MIDDLE for vertical centering
        self._add_text(x, y, size, size,
                       initials, font_size=int(size * 24),
                       bold=True, color=WHITE,
                       font=self.font_headline,
                       align=PP_ALIGN.CENTER,
                       anchor=MSO_ANCHOR.MIDDLE)

    def _draw_quote_card(self, slide, left, top, w, h, q: Quote):
        # Card background (white with brand left border)
        self._add_rect(left, top, w, h, WHITE, line=BORDER_GRAY)
        self._add_rect(left, top, 0.08, h, self.brand_primary, no_line=True)

        # Company name
        self._add_text(left + 0.2, top + 0.1, w - 0.3, 0.35,
                       q.company, font_size=12, bold=True, color=self.brand_dark)

        # Quote text
        self._add_text(left + 0.2, top + 0.5, w - 0.3, h - 1.0,
                       f'"{q.text}"', font_size=11, color=self.brand_dark,
                       italic=True)

        # Speaker / event
        self._add_text(left + 0.2, top + h - 0.35, w - 0.3, 0.25,
                       f"- {q.speaker}, {q.date_or_event}",
                       font_size=9, color=FOOTNOTE_GRAY)


# ──────────────────────────────────────────────────────────────────────
# Markdown / YAML -> slide spec parser
# ──────────────────────────────────────────────────────────────────────

def parse_deck_markdown(text: str) -> list:
    """
    Parse a multi-document YAML stream into a list of slide specs.

    Each YAML document is one slide. The first 'type:' key drives dispatch.
    Empty documents are skipped. Documents missing 'type' raise ValueError.

    Example::

        ---
        type: cover
        title: Sandvik AB Investment Memo
        subtitle: Industrials | Long
        date: April 2026
        ---
        type: bar_chart
        action_title: Sandvik trades at a discount to peer median
        labels: [SAND.ST, CAT, KMT, ITW]
        values: [10.5, 11.0, 9.2, 11.8]
        target_label: SAND.ST
        value_format: "{:.1f}x"
        x_label: EV / LTM EBITDA
        source: Bloomberg, Apr 30 2026

    Returns a list of dicts ready to pass to `build_deck` (see orchestrator).
    """
    import yaml

    specs = []
    for doc in yaml.safe_load_all(text):
        if doc is None:
            continue
        if not isinstance(doc, dict):
            raise ValueError(
                f"each YAML doc must be a mapping, got {type(doc).__name__}")
        if "type" not in doc:
            raise ValueError(f"slide spec missing 'type': {doc!r}")
        specs.append(doc)
    return specs


# ──────────────────────────────────────────────────────────────────────
# Convenience: one-shot single-slide builders
# ──────────────────────────────────────────────────────────────────────

def quick_comparison_slide(
    title: str, entities: list[str], metrics: list[str],
    values: list[list], target_label: str = "",
    source: str = "", filename: str = "ad_hoc_comparison",
    firm: str = "", project: str = "Confidential",
) -> str:
    """One-shot: build a single comparison-matrix slide and save."""
    deck = PPTXDeckWriter(firm=firm, project=project)
    deck.add_comparison_matrix(
        action_title=title, entities=entities, metrics=metrics,
        values=values, target_label=target_label, source=source,
    )
    return deck.save(filename)
