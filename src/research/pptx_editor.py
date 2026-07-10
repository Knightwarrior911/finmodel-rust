"""
In-place PowerPoint edit helpers.

Companion to pptx_inspector.py (read) and pptx_writer.py (rebuild).
Implements the inspect -> plan -> edit -> render -> verify loop described in
docs/SPEC_powerpoint_editing.md.

Edit categories covered here:
  A. Text content swaps          replace_text_in_slide / replace_text_in_deck
  B. Image swaps                 replace_picture
  D. Layout/shape iteration      iter_named_shapes / find_shape_by_id
  E. Slide structure (OOXML)     duplicate_slide / delete_slide / reorder_slides

Slide-management helpers are OOXML-aware: they manipulate
[Content_Types].xml, ppt/presentation.xml, ppt/_rels/presentation.xml.rels,
and the slide parts directly via zipfile + lxml so the resulting package has
no orphan parts and no duplicate-name warnings on save.

Charts are out-of-scope for this module. For chart data edits use
pptx_writer.py chart APIs in a single open/edit/save pass; never round-trip
chart-bearing files through python-pptx after.
"""

from __future__ import annotations

import functools
import inspect as _inspect
import io
import json
import shutil
import threading
import zipfile
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Callable, Iterator, Optional

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.util import Emu, Inches, Pt


# ── XML namespaces ───────────────────────────────────────────────────────────

NS = {
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "ct": "http://schemas.openxmlformats.org/package/2006/content-types",
    "pr": "http://schemas.openxmlformats.org/package/2006/relationships",
}

THEME_SLOTS = (
    "dk1", "lt1", "dk2", "lt2",
    "accent1", "accent2", "accent3", "accent4", "accent5", "accent6",
    "hlink", "folHlink",
)


# ─────────────────────────────────────────────────────────────────────────────
# Edit log (Phase 5 — cross-call/replay memory)
#
# Every top-level edit appends a JSONL line to <deck>.edit_log.jsonl next to
# the deck. Nested calls (a macro that internally invokes primitives) only
# log the macro, not the primitives, via a thread-local depth counter.
# ─────────────────────────────────────────────────────────────────────────────

_log_state = threading.local()


def _log_path_for(deck_path: str | Path) -> Path:
    p = Path(str(deck_path))
    return p.with_suffix(p.suffix + ".edit_log.jsonl")


def _log_edit(deck_path: str | Path, op: str, params: dict) -> None:
    """Append an edit entry to the deck's edit log. Never raises."""
    try:
        log = _log_path_for(deck_path)
        log.parent.mkdir(parents=True, exist_ok=True)
        cleaned = {}
        for k, v in params.items():
            if k.startswith("_"):
                continue
            if isinstance(v, Path):
                cleaned[k] = str(v)
            else:
                cleaned[k] = v
        entry = {
            "ts": datetime.now(timezone.utc).isoformat(),
            "op": op,
            "params": cleaned,
        }
        with open(log, "a", encoding="utf-8") as f:
            f.write(json.dumps(entry, default=str) + "\n")
    except Exception:
        pass  # logging never blocks the actual edit


def _safe_params(fn: Callable, args: tuple, kwargs: dict) -> dict:
    """Bind args+kwargs to fn signature; return as plain dict."""
    try:
        sig = _inspect.signature(fn)
        bound = sig.bind_partial(*args, **kwargs)
        bound.apply_defaults()
        params = dict(bound.arguments)
        params.pop("deck_path", None)
        return params
    except Exception:
        return {**kwargs}


def _logged(op_name: str):
    """
    Decorator: log the operation if at top-level call depth.
    Nested calls (macro -> primitive) do not log; only the outermost wrapper
    runs the log write. Errors in logging never propagate.
    """
    def deco(fn):
        @functools.wraps(fn)
        def wrapper(*args, **kwargs):
            depth = getattr(_log_state, "depth", 0)
            _log_state.depth = depth + 1
            try:
                result = fn(*args, **kwargs)
            finally:
                _log_state.depth = depth
            if depth == 0:
                try:
                    deck_path = args[0] if args else kwargs.get("deck_path")
                    target = kwargs.get("output_path") or deck_path
                    params = _safe_params(fn, args, kwargs)
                    _log_edit(target, op_name, params)
                except Exception:
                    pass
            return result
        return wrapper
    return deco


def get_edit_history(
    deck_path: str | Path,
    *,
    last_n: int = 20,
    op_filter: Optional[list[str]] = None,
) -> list[dict]:
    """
    Read the last_n entries from the deck's edit log.
    Optional op_filter restricts to specific operation names.
    """
    log = _log_path_for(deck_path)
    if not log.exists():
        return []
    lines = log.read_text(encoding="utf-8").splitlines()
    out: list[dict] = []
    for line in lines:
        try:
            entry = json.loads(line)
        except json.JSONDecodeError:
            continue
        if op_filter is not None and entry.get("op") not in op_filter:
            continue
        out.append(entry)
    return out[-last_n:]


def clear_edit_history(deck_path: str | Path) -> None:
    """Delete the edit log for a deck."""
    log = _log_path_for(deck_path)
    if log.exists():
        log.unlink()

PRES_XML = "ppt/presentation.xml"
PRES_RELS = "ppt/_rels/presentation.xml.rels"
CONTENT_TYPES = "[Content_Types].xml"
SLIDE_REL_TYPE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide"
)
SLIDE_CONTENT_TYPE = (
    "application/vnd.openxmlformats-officedocument.presentationml.slide+xml"
)


# ─────────────────────────────────────────────────────────────────────────────
# Text edits
# ─────────────────────────────────────────────────────────────────────────────

def replace_text_in_slide(slide, old: str, new: str) -> int:
    """
    Replace text on a single slide while preserving run-level formatting
    (font, size, color, bold, italic). Returns the number of substitutions.
    """
    count = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if old in run.text:
                    run.text = run.text.replace(old, new)
                    count += 1
    return count


@_logged("replace_text_in_deck")
def replace_text_in_deck(
    deck_path: str | Path,
    replacements: dict[str, str],
    output_path: Optional[str | Path] = None,
) -> str:
    """
    Apply a {old: new} substitution map across every slide. Format-preserving.
    Returns the saved path (overwrites in place when output_path is None).
    """
    deck_path = str(deck_path)
    output_path = str(output_path or deck_path)
    prs = Presentation(deck_path)
    for slide in prs.slides:
        for old, new in replacements.items():
            replace_text_in_slide(slide, old, new)
    prs.save(output_path)
    return output_path


def set_placeholder_text(slide, idx: int, text: str) -> None:
    """
    Set a placeholder's text by its idx, preserving the first run's formatting.
    Raises KeyError if no placeholder with that idx exists on the slide.
    """
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            tf = ph.text_frame
            if tf.paragraphs and tf.paragraphs[0].runs:
                tf.paragraphs[0].runs[0].text = text
                for run in tf.paragraphs[0].runs[1:]:
                    run.text = ""
                for extra in tf.paragraphs[1:]:
                    extra._pPr = None
            else:
                tf.text = text
            return
    raise KeyError(f"No placeholder with idx={idx} on this slide")


# ─────────────────────────────────────────────────────────────────────────────
# Image edits
# ─────────────────────────────────────────────────────────────────────────────

@_logged("replace_picture")
def replace_picture(
    deck_path: str | Path,
    slide_index: int,
    new_image_path: str | Path,
    *,
    shape_name: Optional[str] = None,
    shape_id: Optional[int] = None,
    output_path: Optional[str | Path] = None,
) -> str:
    """
    Replace an existing picture shape with a new image at the same box.

    Identify the target by either shape_name or shape_id (one is required).
    Position and size of the original shape are preserved exactly.
    """
    if shape_name is None and shape_id is None:
        raise ValueError("Pass shape_name or shape_id to identify the target")

    deck_path = str(deck_path)
    output_path = str(output_path or deck_path)
    new_image_path = str(new_image_path)

    prs = Presentation(deck_path)
    slide = prs.slides[slide_index]

    target = None
    for shape in slide.shapes:
        if shape_id is not None and shape.shape_id == shape_id:
            target = shape
            break
        if shape_name is not None and shape.name == shape_name:
            target = shape
            break
    if target is None:
        raise LookupError(
            f"No shape with name={shape_name!r} or id={shape_id!r} "
            f"on slide {slide_index}"
        )

    left, top, width, height = target.left, target.top, target.width, target.height
    target._element.getparent().remove(target._element)
    slide.shapes.add_picture(new_image_path, left, top, width=width, height=height)

    prs.save(output_path)
    return output_path


# ─────────────────────────────────────────────────────────────────────────────
# Shape iteration / lookup
# ─────────────────────────────────────────────────────────────────────────────

def iter_named_shapes(
    slide,
    *,
    prefix: Optional[str] = None,
    suffix: Optional[str] = None,
    contains: Optional[str] = None,
) -> Iterator:
    """Yield shapes whose .name matches the given prefix/suffix/contains filter."""
    for shape in slide.shapes:
        name = shape.name or ""
        if prefix is not None and not name.startswith(prefix):
            continue
        if suffix is not None and not name.endswith(suffix):
            continue
        if contains is not None and contains not in name:
            continue
        yield shape


def find_shape_by_id(slide, shape_id: int):
    """Return the shape on slide with the given shape_id, or None."""
    for shape in slide.shapes:
        if shape.shape_id == shape_id:
            return shape
    return None


# ─────────────────────────────────────────────────────────────────────────────
# Theme recolor (Category C — palette rebrand)
#
# Modifies ppt/theme/themeN.xml accent slots so every shape that references a
# theme color updates automatically. For shapes that hard-code RGB, also pass
# `also_replace_hardcoded={old_hex: new_hex}` to swap srgbClr values across
# all slide XML.
# ─────────────────────────────────────────────────────────────────────────────

def _normalise_hex(value: str) -> str:
    """Accept '#RRGGBB' or 'RRGGBB' (case-insensitive), return uppercase 'RRGGBB'."""
    v = value.lstrip("#").strip()
    if len(v) != 6 or any(c not in "0123456789abcdefABCDEF" for c in v):
        raise ValueError(f"Invalid hex color: {value!r}")
    return v.upper()


def _recolor_clr_scheme(theme_xml: bytes, palette: dict[str, str]) -> bytes:
    """Replace accent/dk/lt/hlink slots in <a:clrScheme>. Returns updated bytes."""
    root = _parse(theme_xml)
    a = NS["a"]
    scheme = root.find(f".//{{{a}}}themeElements/{{{a}}}clrScheme")
    if scheme is None:
        return theme_xml

    for slot, new_hex in palette.items():
        if slot not in THEME_SLOTS:
            raise KeyError(
                f"Unknown theme slot {slot!r}. Valid slots: {THEME_SLOTS}"
            )
        slot_el = scheme.find(f"{{{a}}}{slot}")
        if slot_el is None:
            continue
        for child in list(slot_el):
            slot_el.remove(child)
        srgb = etree.SubElement(slot_el, f"{{{a}}}srgbClr")
        srgb.set("val", _normalise_hex(new_hex))

    return _serialise(root)


def _replace_srgb_in_xml(xml_bytes: bytes, swaps: dict[str, str]) -> bytes:
    """Replace <a:srgbClr val="OLD"/> with NEW everywhere, case-insensitive."""
    root = _parse(xml_bytes)
    norm = {_normalise_hex(o): _normalise_hex(n) for o, n in swaps.items()}
    a = NS["a"]
    for el in root.iter(f"{{{a}}}srgbClr"):
        val = el.get("val", "").upper()
        if val in norm:
            el.set("val", norm[val])
    return _serialise(root)


@_logged("recolor_theme")
def recolor_theme(
    deck_path: str | Path,
    palette: dict[str, str],
    *,
    also_replace_hardcoded: Optional[dict[str, str]] = None,
    output_path: Optional[str | Path] = None,
) -> str:
    """
    Recolor a deck's theme accent slots in place. Shapes that reference theme
    colors update automatically; hard-coded RGBs need `also_replace_hardcoded`.

    palette: {"accent1": "#255BE3", "accent2": "#0F1632", "dk1": "000000", ...}
    also_replace_hardcoded: {"4472C4": "255BE3", ...} optional global RGB swap.

    Returns the saved path.
    """
    deck_path = str(deck_path)
    out_path = _resolve_output(deck_path, output_path)
    members = _read_zip(deck_path)

    theme_parts = [n for n in members if n.startswith("ppt/theme/") and n.endswith(".xml")]
    if not theme_parts:
        raise RuntimeError("No theme XML found in deck")

    for name in theme_parts:
        members[name] = _recolor_clr_scheme(members[name], palette)

    if also_replace_hardcoded:
        for name in list(members):
            if name.endswith(".xml") and (
                name.startswith("ppt/slides/")
                or name.startswith("ppt/slideLayouts/")
                or name.startswith("ppt/slideMasters/")
                or name.startswith("ppt/theme/")
            ):
                members[name] = _replace_srgb_in_xml(
                    members[name], also_replace_hardcoded
                )

    _write_zip(out_path, members)
    return out_path


# ─────────────────────────────────────────────────────────────────────────────
# Slide structure (OOXML-aware)
#
# python-pptx's sldIdLst.remove() leaves orphan parts inside the OOXML
# package, producing duplicate-name warnings on save. The helpers below
# operate on the .pptx zip directly so the resulting package is clean.
# ─────────────────────────────────────────────────────────────────────────────

def _read_zip(path: str) -> dict[str, bytes]:
    with zipfile.ZipFile(path, "r") as zf:
        return {name: zf.read(name) for name in zf.namelist()}


def _write_zip(path: str, members: dict[str, bytes]) -> None:
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as zf:
        for name, data in members.items():
            zf.writestr(name, data)


def _parse(xml_bytes: bytes):
    return etree.fromstring(xml_bytes)


def _serialise(elem) -> bytes:
    return etree.tostring(
        elem, xml_declaration=True, encoding="UTF-8", standalone=True
    )


def _slide_part_names(members: dict[str, bytes]) -> list[str]:
    """All ppt/slides/slideN.xml entries, sorted by N."""
    return sorted(
        (n for n in members if n.startswith("ppt/slides/slide") and n.endswith(".xml")),
        key=lambda n: int(n.rsplit("slide", 1)[1].rsplit(".", 1)[0]),
    )


def _next_slide_index(members: dict[str, bytes]) -> int:
    used = set()
    for n in _slide_part_names(members):
        used.add(int(n.rsplit("slide", 1)[1].rsplit(".", 1)[0]))
    i = 1
    while i in used:
        i += 1
    return i


def _rewrite_presentation_rels(
    rels_xml: bytes, ordered_slide_targets: list[str]
) -> bytes:
    """
    Rewrite ppt/_rels/presentation.xml.rels so slide rIds match the new order.
    ordered_slide_targets is a list of relative targets like "slides/slide3.xml".
    Non-slide relationships are preserved untouched.
    """
    root = _parse(rels_xml)
    ns_pr = NS["pr"]
    nsmap = {None: ns_pr}

    keep = []
    for rel in root.findall(f"{{{ns_pr}}}Relationship"):
        if rel.get("Type") != SLIDE_REL_TYPE:
            keep.append(rel)

    used_ids = {r.get("Id") for r in keep}

    def next_id() -> str:
        n = 1
        while f"rId{n}" in used_ids:
            n += 1
        used_ids.add(f"rId{n}")
        return f"rId{n}"

    new_root = etree.Element(f"{{{ns_pr}}}Relationships", nsmap=nsmap)
    for rel in keep:
        new_root.append(rel)

    slide_id_to_rid: dict[str, str] = {}
    for target in ordered_slide_targets:
        rid = next_id()
        slide_id_to_rid[target] = rid
        rel = etree.SubElement(new_root, f"{{{ns_pr}}}Relationship")
        rel.set("Id", rid)
        rel.set("Type", SLIDE_REL_TYPE)
        rel.set("Target", target)

    return _serialise(new_root), slide_id_to_rid


def _rewrite_presentation_xml(
    pres_xml: bytes,
    ordered_slide_rids: list[str],
) -> bytes:
    """Replace <p:sldIdLst> with new <p:sldId r:id="rIdN"/> entries."""
    root = _parse(pres_xml)
    sld_id_lst = root.find(f"{{{NS['p']}}}sldIdLst")
    if sld_id_lst is None:
        raise ValueError("presentation.xml has no <p:sldIdLst>")

    # Preserve max existing id for stable ordering, allocate new sequential
    base = 256
    for child in list(sld_id_lst):
        sld_id_lst.remove(child)
    for i, rid in enumerate(ordered_slide_rids):
        sld_id = etree.SubElement(sld_id_lst, f"{{{NS['p']}}}sldId")
        sld_id.set("id", str(base + i))
        sld_id.set(f"{{{NS['r']}}}id", rid)
    return _serialise(root)


def _rewrite_content_types(
    ct_xml: bytes, slide_part_names: list[str]
) -> bytes:
    """Ensure [Content_Types].xml lists every current slide as Override."""
    root = _parse(ct_xml)
    ns_ct = NS["ct"]
    # Drop existing slide overrides; re-add for kept slides.
    for child in list(root):
        if child.tag == f"{{{ns_ct}}}Override":
            ct = child.get("ContentType", "")
            part = child.get("PartName", "")
            if ct == SLIDE_CONTENT_TYPE or part.startswith("/ppt/slides/slide"):
                root.remove(child)
    for name in slide_part_names:
        ov = etree.SubElement(root, f"{{{ns_ct}}}Override")
        ov.set("PartName", "/" + name)
        ov.set("ContentType", SLIDE_CONTENT_TYPE)
    return _serialise(root)


def _renumber_and_reorder(
    members: dict[str, bytes], desired_order: list[int]
) -> dict[str, bytes]:
    """
    Given the original deck member dict and a desired_order of original
    1-based slide indices, produce a new member dict with slides renamed
    contiguously (slide1.xml, slide2.xml, ...) and presentation/rels/CT
    rewritten to match.
    """
    out = dict(members)
    original_parts = _slide_part_names(members)

    # Map original index -> original part name, rels name
    by_idx = {
        int(n.rsplit("slide", 1)[1].rsplit(".", 1)[0]): n for n in original_parts
    }

    # Drop all old slide xml + rels from out
    for old_name in original_parts:
        out.pop(old_name, None)
        rels_name = old_name.replace("ppt/slides/", "ppt/slides/_rels/") + ".rels"
        if rels_name in members:
            out.pop(rels_name, None)

    # Re-add in new order with new contiguous names
    new_part_names = []
    new_rel_targets = []
    for new_idx, original_idx in enumerate(desired_order, start=1):
        old_part = by_idx[original_idx]
        new_part = f"ppt/slides/slide{new_idx}.xml"
        out[new_part] = members[old_part]
        new_part_names.append(new_part)
        new_rel_targets.append(f"slides/slide{new_idx}.xml")

        old_rels = old_part.replace("ppt/slides/", "ppt/slides/_rels/") + ".rels"
        new_rels = f"ppt/slides/_rels/slide{new_idx}.xml.rels"
        if old_rels in members:
            out[new_rels] = members[old_rels]

    # Rewrite presentation.xml.rels
    rels_xml, target_to_rid = _rewrite_presentation_rels(
        members[PRES_RELS], new_rel_targets
    )
    out[PRES_RELS] = rels_xml

    ordered_rids = [target_to_rid[t] for t in new_rel_targets]
    out[PRES_XML] = _rewrite_presentation_xml(members[PRES_XML], ordered_rids)
    out[CONTENT_TYPES] = _rewrite_content_types(members[CONTENT_TYPES], new_part_names)

    return out


def _resolve_output(deck_path: str, output_path: Optional[str | Path]) -> str:
    if output_path is None:
        return deck_path
    out = str(output_path)
    if out != deck_path:
        Path(out).parent.mkdir(parents=True, exist_ok=True)
    return out


@_logged("duplicate_slide")
def duplicate_slide(
    deck_path: str | Path,
    slide_index: int,
    *,
    position: Optional[int] = None,
    output_path: Optional[str | Path] = None,
) -> str:
    """
    Duplicate the slide at slide_index (0-based). New slide is inserted at
    `position` (0-based) or appended to the end when position is None.
    """
    deck_path = str(deck_path)
    out_path = _resolve_output(deck_path, output_path)
    members = _read_zip(deck_path)
    parts = _slide_part_names(members)
    n = len(parts)
    if not 0 <= slide_index < n:
        raise IndexError(f"slide_index {slide_index} out of range (0..{n-1})")

    original_order = list(range(1, n + 1))  # 1-based original indices
    src_original_idx = original_order[slide_index]

    # We need to copy the underlying part. Re-read the original zip to
    # synthesize a "new" original index for the duplicate that points at the
    # same bytes. Simulate by inserting src_original_idx again into the order
    # — we then remap below to unique part names.
    insert_at = n if position is None else max(0, min(position, n))
    new_order = list(original_order)
    new_order.insert(insert_at, src_original_idx)

    # Materialise the duplicate: assign it a fresh original index so the
    # renumber step can copy bytes. Allocate next unused index.
    free_idx = _next_slide_index(members)
    src_part = f"ppt/slides/slide{src_original_idx}.xml"
    src_rels = f"ppt/slides/_rels/slide{src_original_idx}.xml.rels"
    dup_part = f"ppt/slides/slide{free_idx}.xml"
    dup_rels = f"ppt/slides/_rels/slide{free_idx}.xml.rels"
    members[dup_part] = members[src_part]
    if src_rels in members:
        members[dup_rels] = members[src_rels]

    # Rebuild order using the new free_idx in place of the second occurrence
    # of src_original_idx (the inserted one).
    order: list[int] = []
    seen = False
    for orig in new_order:
        if orig == src_original_idx and not seen:
            order.append(orig)
            seen = True
        elif orig == src_original_idx and seen:
            order.append(free_idx)
        else:
            order.append(orig)

    out = _renumber_and_reorder(members, order)
    _write_zip(out_path, out)
    return out_path


@_logged("delete_slide")
def delete_slide(
    deck_path: str | Path,
    slide_index: int,
    output_path: Optional[str | Path] = None,
) -> str:
    """Delete the slide at slide_index (0-based), cleaning up orphan parts."""
    deck_path = str(deck_path)
    out_path = _resolve_output(deck_path, output_path)
    members = _read_zip(deck_path)
    parts = _slide_part_names(members)
    n = len(parts)
    if not 0 <= slide_index < n:
        raise IndexError(f"slide_index {slide_index} out of range (0..{n-1})")

    keep_original = [
        int(p.rsplit("slide", 1)[1].rsplit(".", 1)[0])
        for i, p in enumerate(parts)
        if i != slide_index
    ]
    out = _renumber_and_reorder(members, keep_original)
    _write_zip(out_path, out)
    return out_path


@_logged("reorder_slides")
def reorder_slides(
    deck_path: str | Path,
    new_order: list[int],
    output_path: Optional[str | Path] = None,
) -> str:
    """
    Reorder slides. new_order is a list of 0-based indices that must be a
    permutation of range(slide_count).
    """
    deck_path = str(deck_path)
    out_path = _resolve_output(deck_path, output_path)
    members = _read_zip(deck_path)
    parts = _slide_part_names(members)
    n = len(parts)
    if sorted(new_order) != list(range(n)):
        raise ValueError(
            f"new_order must be a permutation of 0..{n-1}; got {new_order}"
        )
    by_idx = [int(p.rsplit("slide", 1)[1].rsplit(".", 1)[0]) for p in parts]
    desired = [by_idx[i] for i in new_order]
    out = _renumber_and_reorder(members, desired)
    _write_zip(out_path, out)
    return out_path


# ─────────────────────────────────────────────────────────────────────────────
# Granular shape primitives (Phase 1 — MD-style edit support)
#
# Each primitive opens the deck, edits one shape (or adds one), saves. All
# coordinate inputs are in inches (floats). Colors are hex strings ("#RRGGBB"
# or "RRGGBB"). Each call returns the saved path.
#
# Pair these with inspect_pptx_with_preview() so the LLM can map fuzzy MD
# language ("the chart on the right", "the Falcon column") to specific
# shape IDs before calling these.
# ─────────────────────────────────────────────────────────────────────────────

DASH_STYLES = {
    "solid": MSO_LINE_DASH_STYLE.SOLID,
    "dash": MSO_LINE_DASH_STYLE.DASH,
    "dot": MSO_LINE_DASH_STYLE.ROUND_DOT,
    "dashdot": MSO_LINE_DASH_STYLE.DASH_DOT,
    "longdash": MSO_LINE_DASH_STYLE.LONG_DASH,
}

SHAPE_KINDS = {
    "rect": MSO_SHAPE.RECTANGLE,
    "rectangle": MSO_SHAPE.RECTANGLE,
    "rrect": MSO_SHAPE.ROUNDED_RECTANGLE,
    "rounded": MSO_SHAPE.ROUNDED_RECTANGLE,
    "oval": MSO_SHAPE.OVAL,
    "circle": MSO_SHAPE.OVAL,
    "capsule": MSO_SHAPE.ROUNDED_RECTANGLE,
    "line": MSO_SHAPE.LINE_INVERSE,
    "arrow": MSO_SHAPE.RIGHT_ARROW,
}


def _hex_to_rgb(color: Optional[str]) -> Optional[RGBColor]:
    if color is None:
        return None
    h = _normalise_hex(color)
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _find_shape(slide, *, shape_id: Optional[int] = None,
                shape_name: Optional[str] = None):
    if shape_id is None and shape_name is None:
        raise ValueError("Pass shape_id or shape_name to identify the shape")
    for sh in slide.shapes:
        if shape_id is not None and sh.shape_id == shape_id:
            return sh
        if shape_name is not None and sh.name == shape_name:
            return sh
    raise LookupError(
        f"No shape with id={shape_id!r} name={shape_name!r} on slide"
    )


def _open_for_edit(deck_path: str | Path, output_path: Optional[str | Path]):
    deck_path = str(deck_path)
    out = str(output_path or deck_path)
    prs = Presentation(deck_path)
    return prs, out


# ── Move / resize ────────────────────────────────────────────────────────────

@_logged("move_shape")
def move_shape(
    deck_path: str | Path,
    slide_index: int,
    *,
    shape_id: Optional[int] = None,
    shape_name: Optional[str] = None,
    left: Optional[float] = None,
    top: Optional[float] = None,
    dx: Optional[float] = None,
    dy: Optional[float] = None,
    output_path: Optional[str | Path] = None,
) -> str:
    """
    Move a shape. Pass `left`/`top` for absolute (inches) or `dx`/`dy` for
    relative deltas (inches). At least one of left/top/dx/dy is required.
    """
    if all(v is None for v in (left, top, dx, dy)):
        raise ValueError("Pass at least one of left/top/dx/dy")
    prs, out = _open_for_edit(deck_path, output_path)
    shape = _find_shape(prs.slides[slide_index],
                        shape_id=shape_id, shape_name=shape_name)
    if left is not None:
        shape.left = Inches(left)
    elif dx is not None:
        shape.left = shape.left + Inches(dx)
    if top is not None:
        shape.top = Inches(top)
    elif dy is not None:
        shape.top = shape.top + Inches(dy)
    prs.save(out)
    return out


@_logged("resize_shape")
def resize_shape(
    deck_path: str | Path,
    slide_index: int,
    *,
    shape_id: Optional[int] = None,
    shape_name: Optional[str] = None,
    width: Optional[float] = None,
    height: Optional[float] = None,
    output_path: Optional[str | Path] = None,
) -> str:
    """Resize a shape to width/height (inches). Either or both."""
    if width is None and height is None:
        raise ValueError("Pass at least one of width/height")
    prs, out = _open_for_edit(deck_path, output_path)
    shape = _find_shape(prs.slides[slide_index],
                        shape_id=shape_id, shape_name=shape_name)
    if width is not None:
        shape.width = Inches(width)
    if height is not None:
        shape.height = Inches(height)
    prs.save(out)
    return out


# ── Fill / line ──────────────────────────────────────────────────────────────

@_logged("set_shape_fill")
def set_shape_fill(
    deck_path: str | Path,
    slide_index: int,
    *,
    shape_id: Optional[int] = None,
    shape_name: Optional[str] = None,
    color: Optional[str] = None,
    no_fill: bool = False,
    output_path: Optional[str | Path] = None,
) -> str:
    """Set solid fill color (hex). Pass no_fill=True to clear fill."""
    prs, out = _open_for_edit(deck_path, output_path)
    shape = _find_shape(prs.slides[slide_index],
                        shape_id=shape_id, shape_name=shape_name)
    if no_fill:
        shape.fill.background()
    elif color is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb(color)
    else:
        raise ValueError("Pass color or no_fill=True")
    prs.save(out)
    return out


@_logged("set_shape_line")
def set_shape_line(
    deck_path: str | Path,
    slide_index: int,
    *,
    shape_id: Optional[int] = None,
    shape_name: Optional[str] = None,
    color: Optional[str] = None,
    width: Optional[float] = None,
    dash: Optional[str] = None,
    no_line: bool = False,
    output_path: Optional[str | Path] = None,
) -> str:
    """Set line color (hex), width (pt), and dash ('solid'|'dash'|'dot'|...)."""
    prs, out = _open_for_edit(deck_path, output_path)
    shape = _find_shape(prs.slides[slide_index],
                        shape_id=shape_id, shape_name=shape_name)
    if no_line:
        shape.line.fill.background()
    else:
        if color is not None:
            shape.line.color.rgb = _hex_to_rgb(color)
        if width is not None:
            shape.line.width = Pt(width)
        if dash is not None:
            if dash not in DASH_STYLES:
                raise ValueError(
                    f"Unknown dash {dash!r}. Valid: {list(DASH_STYLES)}"
                )
            shape.line.dash_style = DASH_STYLES[dash]
    prs.save(out)
    return out


# ── Text styling ─────────────────────────────────────────────────────────────

@_logged("set_text_style")
def set_text_style(
    deck_path: str | Path,
    slide_index: int,
    *,
    shape_id: Optional[int] = None,
    shape_name: Optional[str] = None,
    paragraph_index: Optional[int] = None,
    run_index: Optional[int] = None,
    bold: Optional[bool] = None,
    italic: Optional[bool] = None,
    underline: Optional[bool] = None,
    color: Optional[str] = None,
    size: Optional[float] = None,
    font_name: Optional[str] = None,
    text: Optional[str] = None,
    output_path: Optional[str | Path] = None,
) -> str:
    """
    Restyle text in a shape. Targets paragraph_index + run_index (both 0-based)
    when given, else applies to all runs in the shape. Any of
    bold/italic/underline/color/size/font_name/text may be set.
    """
    prs, out = _open_for_edit(deck_path, output_path)
    shape = _find_shape(prs.slides[slide_index],
                        shape_id=shape_id, shape_name=shape_name)
    if not shape.has_text_frame:
        raise ValueError("Shape has no text frame")

    paras = shape.text_frame.paragraphs
    target_paras = (
        [paras[paragraph_index]] if paragraph_index is not None else paras
    )

    for para in target_paras:
        runs = para.runs
        target_runs = (
            [runs[run_index]] if (run_index is not None and runs) else runs
        )
        for run in target_runs:
            if text is not None:
                run.text = text
            if bold is not None:
                run.font.bold = bold
            if italic is not None:
                run.font.italic = italic
            if underline is not None:
                run.font.underline = underline
            if color is not None:
                run.font.color.rgb = _hex_to_rgb(color)
            if size is not None:
                run.font.size = Pt(size)
            if font_name is not None:
                run.font.name = font_name

    prs.save(out)
    return out


# ── Add new shapes ───────────────────────────────────────────────────────────

@_logged("add_textbox")
def add_textbox(
    deck_path: str | Path,
    slide_index: int,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str = "",
    bold: bool = False,
    italic: bool = False,
    color: Optional[str] = None,
    size: Optional[float] = None,
    font_name: Optional[str] = None,
    name: Optional[str] = None,
    output_path: Optional[str | Path] = None,
) -> str:
    """Add a textbox at (left, top) sized (width, height) inches."""
    prs, out = _open_for_edit(deck_path, output_path)
    slide = prs.slides[slide_index]
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    if name:
        tb.name = name
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.bold = bold
    run.font.italic = italic
    if color is not None:
        run.font.color.rgb = _hex_to_rgb(color)
    if size is not None:
        run.font.size = Pt(size)
    if font_name:
        run.font.name = font_name
    prs.save(out)
    return out


@_logged("add_line")
def add_line(
    deck_path: str | Path,
    slide_index: int,
    *,
    x1: float, y1: float, x2: float, y2: float,
    color: str = "#000000",
    width: float = 1.0,
    dash: str = "solid",
    name: Optional[str] = None,
    output_path: Optional[str | Path] = None,
) -> str:
    """Add a straight line connector from (x1,y1) to (x2,y2) inches."""
    prs, out = _open_for_edit(deck_path, output_path)
    slide = prs.slides[slide_index]
    from pptx.enum.shapes import MSO_CONNECTOR
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2),
    )
    if name:
        line.name = name
    line.line.color.rgb = _hex_to_rgb(color)
    line.line.width = Pt(width)
    if dash in DASH_STYLES:
        line.line.dash_style = DASH_STYLES[dash]
    prs.save(out)
    return out


@_logged("add_shape_box")
def add_shape_box(
    deck_path: str | Path,
    slide_index: int,
    *,
    kind: str,
    left: float, top: float, width: float, height: float,
    fill: Optional[str] = None,
    line: Optional[str] = None,
    line_width: Optional[float] = None,
    corner: Optional[float] = None,
    text: str = "",
    text_color: Optional[str] = None,
    text_size: Optional[float] = None,
    bold: bool = False,
    name: Optional[str] = None,
    output_path: Optional[str | Path] = None,
) -> str:
    """
    Add an autoshape (rect / rrect / oval / circle / capsule / arrow).
    Optional fill / line color (hex), optional inline text.
    """
    if kind not in SHAPE_KINDS:
        raise ValueError(f"Unknown kind {kind!r}. Valid: {list(SHAPE_KINDS)}")
    prs, out = _open_for_edit(deck_path, output_path)
    slide = prs.slides[slide_index]
    sh = slide.shapes.add_shape(
        SHAPE_KINDS[kind],
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    if name:
        sh.name = name
    if kind == "capsule":
        sh.adjustments[0] = 0.5
    elif corner is not None and SHAPE_KINDS[kind] == MSO_SHAPE.ROUNDED_RECTANGLE:
        sh.adjustments[0] = corner
    if fill is not None:
        sh.fill.solid()
        sh.fill.fore_color.rgb = _hex_to_rgb(fill)
    else:
        sh.fill.background()
    if line is not None:
        sh.line.color.rgb = _hex_to_rgb(line)
        if line_width is not None:
            sh.line.width = Pt(line_width)
    else:
        sh.line.fill.background()
    if text:
        tf = sh.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.bold = bold
        if text_color is not None:
            run.font.color.rgb = _hex_to_rgb(text_color)
        if text_size is not None:
            run.font.size = Pt(text_size)
    prs.save(out)
    return out


# ── Delete / align / distribute / copy-style ─────────────────────────────────

@_logged("delete_shape")
def delete_shape(
    deck_path: str | Path,
    slide_index: int,
    *,
    shape_id: Optional[int] = None,
    shape_name: Optional[str] = None,
    output_path: Optional[str | Path] = None,
) -> str:
    """Remove a shape from a slide."""
    prs, out = _open_for_edit(deck_path, output_path)
    shape = _find_shape(prs.slides[slide_index],
                        shape_id=shape_id, shape_name=shape_name)
    shape._element.getparent().remove(shape._element)
    prs.save(out)
    return out


@_logged("align_shapes")
def align_shapes(
    deck_path: str | Path,
    slide_index: int,
    shape_ids: list[int],
    how: str,
    output_path: Optional[str | Path] = None,
) -> str:
    """
    Align a group of shapes. how: 'left', 'right', 'center' (horizontal),
    'top', 'bottom', 'middle' (vertical).
    """
    valid = {"left", "right", "center", "top", "bottom", "middle"}
    if how not in valid:
        raise ValueError(f"how must be one of {valid}")
    prs, out = _open_for_edit(deck_path, output_path)
    slide = prs.slides[slide_index]
    shapes = [_find_shape(slide, shape_id=sid) for sid in shape_ids]
    if not shapes:
        return out

    if how == "left":
        target = min(s.left for s in shapes)
        for s in shapes:
            s.left = target
    elif how == "right":
        target = max(s.left + s.width for s in shapes)
        for s in shapes:
            s.left = target - s.width
    elif how == "center":
        target = sum(s.left + s.width // 2 for s in shapes) // len(shapes)
        for s in shapes:
            s.left = target - s.width // 2
    elif how == "top":
        target = min(s.top for s in shapes)
        for s in shapes:
            s.top = target
    elif how == "bottom":
        target = max(s.top + s.height for s in shapes)
        for s in shapes:
            s.top = target - s.height
    elif how == "middle":
        target = sum(s.top + s.height // 2 for s in shapes) // len(shapes)
        for s in shapes:
            s.top = target - s.height // 2

    prs.save(out)
    return out


@_logged("distribute_shapes")
def distribute_shapes(
    deck_path: str | Path,
    slide_index: int,
    shape_ids: list[int],
    axis: str,
    output_path: Optional[str | Path] = None,
) -> str:
    """Even-distribute shapes along x or y axis."""
    if axis not in {"x", "y"}:
        raise ValueError("axis must be 'x' or 'y'")
    if len(shape_ids) < 3:
        raise ValueError("Need at least 3 shapes to distribute")

    prs, out = _open_for_edit(deck_path, output_path)
    slide = prs.slides[slide_index]
    shapes = [_find_shape(slide, shape_id=sid) for sid in shape_ids]

    if axis == "x":
        shapes.sort(key=lambda s: s.left)
        first, last = shapes[0], shapes[-1]
        span = (last.left + last.width // 2) - (first.left + first.width // 2)
        step = span // (len(shapes) - 1)
        center0 = first.left + first.width // 2
        for i, s in enumerate(shapes[1:-1], start=1):
            s.left = (center0 + step * i) - s.width // 2
    else:
        shapes.sort(key=lambda s: s.top)
        first, last = shapes[0], shapes[-1]
        span = (last.top + last.height // 2) - (first.top + first.height // 2)
        step = span // (len(shapes) - 1)
        center0 = first.top + first.height // 2
        for i, s in enumerate(shapes[1:-1], start=1):
            s.top = (center0 + step * i) - s.height // 2

    prs.save(out)
    return out


@_logged("copy_style")
def copy_style(
    deck_path: str | Path,
    slide_index: int,
    *,
    source_shape_id: int,
    target_shape_id: int,
    output_path: Optional[str | Path] = None,
) -> str:
    """Copy fill, line, and first-run font properties from source to target."""
    prs, out = _open_for_edit(deck_path, output_path)
    slide = prs.slides[slide_index]
    src = _find_shape(slide, shape_id=source_shape_id)
    tgt = _find_shape(slide, shape_id=target_shape_id)

    try:
        if src.fill.type is not None:
            try:
                rgb = src.fill.fore_color.rgb
                tgt.fill.solid()
                tgt.fill.fore_color.rgb = rgb
            except Exception:
                pass
    except Exception:
        pass

    try:
        if src.line.color.rgb is not None:
            tgt.line.color.rgb = src.line.color.rgb
        if src.line.width is not None:
            tgt.line.width = src.line.width
    except Exception:
        pass

    if src.has_text_frame and tgt.has_text_frame:
        if src.text_frame.paragraphs and src.text_frame.paragraphs[0].runs:
            src_run = src.text_frame.paragraphs[0].runs[0]
            for tp in tgt.text_frame.paragraphs:
                for tr in tp.runs:
                    if src_run.font.bold is not None:
                        tr.font.bold = src_run.font.bold
                    if src_run.font.italic is not None:
                        tr.font.italic = src_run.font.italic
                    if src_run.font.size is not None:
                        tr.font.size = src_run.font.size
                    if src_run.font.name is not None:
                        tr.font.name = src_run.font.name
                    try:
                        if src_run.font.color.rgb is not None:
                            tr.font.color.rgb = src_run.font.color.rgb
                    except Exception:
                        pass

    prs.save(out)
    return out


# ─────────────────────────────────────────────────────────────────────────────
# Vision-augmented inspect (Phase 2)
#
# Returns the JSON descriptor + paths to per-slide PNG previews. The
# orchestrator uses these to feed images to the LLM so it can map fuzzy MD
# language ("the chart on the right") to concrete shape IDs.
# ─────────────────────────────────────────────────────────────────────────────

def inspect_with_preview(
    deck_path: str | Path,
    *,
    slide_indices: Optional[list[int]] = None,
    out_dir: Optional[str | Path] = None,
    dpi: int = 120,
) -> dict:
    """
    Returns:
      {
        "json": {<inspect_pptx output>},
        "previews": [{"slide_index": 0, "png": "/abs/path.png"}, ...]
      }
    Renders the deck once via pptx_render.render_deck and pairs PNGs to slides.
    """
    from src.research.pptx_inspector import inspect_pptx
    from src.research.pptx_render import render_deck

    deck_path = Path(deck_path).resolve()
    desc = inspect_pptx(str(deck_path), include_raw_xml=False)

    rendered = render_deck(deck_path, out_dir=out_dir, dpi=dpi)
    pngs = [p for p in rendered if p.suffix.lower() == ".png"]
    pngs.sort(key=lambda p: p.name.lower())

    n_slides = desc["slideCount"]
    if len(pngs) < n_slides:
        # Render returned fewer PNGs than slides — pair what we have
        n_slides = len(pngs)

    indices = (
        list(range(n_slides)) if slide_indices is None
        else [i for i in slide_indices if 0 <= i < n_slides]
    )
    previews = [{"slide_index": i, "png": str(pngs[i])} for i in indices]

    return {"json": desc, "previews": previews}


# ─────────────────────────────────────────────────────────────────────────────
# Native-table operations
#
# python-pptx tables expose cells but not easy column/row reorder. These
# helpers manipulate the underlying <a:tbl> XML directly: <a:tblGrid> for
# column widths and <a:tr>/<a:tc> for cells.
# ─────────────────────────────────────────────────────────────────────────────

def _get_table_element(slide, *, shape_id=None, shape_name=None):
    target = _find_shape(slide, shape_id=shape_id, shape_name=shape_name)
    if not getattr(target, "has_table", False):
        raise ValueError(
            f"Shape {target.name!r} (id={target.shape_id}) is not a table"
        )
    return target, target.table._tbl


def _reorder_xml_children(parent, tag: str, new_order: list[int]) -> None:
    """Reorder `parent`'s children matching `tag` according to new_order indices."""
    children = list(parent.findall(tag))
    if sorted(new_order) != list(range(len(children))):
        raise ValueError(
            f"new_order must be a permutation of 0..{len(children)-1}; "
            f"got {new_order}"
        )
    reordered = [children[i] for i in new_order]
    for c in children:
        parent.remove(c)
    for c in reordered:
        parent.append(c)


@_logged("swap_table_columns")
def swap_table_columns(
    deck_path: str | Path,
    slide_index: int,
    *,
    shape_id: Optional[int] = None,
    shape_name: Optional[str] = None,
    col_a: int,
    col_b: int,
    output_path: Optional[str | Path] = None,
) -> str:
    """
    Swap two columns of a native PPTX table by 0-based index.
    Updates <a:tblGrid> column widths and every <a:tr>'s <a:tc> cells.
    """
    out = str(output_path or deck_path)
    prs = Presentation(str(deck_path))
    target, tbl = _get_table_element(
        prs.slides[slide_index], shape_id=shape_id, shape_name=shape_name,
    )
    n_cols = len(target.table.columns)
    if not (0 <= col_a < n_cols and 0 <= col_b < n_cols):
        raise IndexError(f"Column indices out of range (0..{n_cols-1})")
    if col_a == col_b:
        prs.save(out)
        return out

    a = NS["a"]
    new_order = list(range(n_cols))
    new_order[col_a], new_order[col_b] = new_order[col_b], new_order[col_a]

    grid = tbl.find(f"{{{a}}}tblGrid")
    if grid is not None:
        _reorder_xml_children(grid, f"{{{a}}}gridCol", new_order)

    for tr in tbl.findall(f"{{{a}}}tr"):
        _reorder_xml_children(tr, f"{{{a}}}tc", new_order)

    prs.save(out)
    return out


@_logged("move_table_column")
def move_table_column(
    deck_path: str | Path,
    slide_index: int,
    *,
    shape_id: Optional[int] = None,
    shape_name: Optional[str] = None,
    col_from: int,
    col_to: int,
    output_path: Optional[str | Path] = None,
) -> str:
    """
    Move a table column from `col_from` to `col_to` (0-based).
    Other columns shift to fill the gap. col_to is the destination index
    AFTER removal, so move(0, 2) moves first column past the next two.
    """
    out = str(output_path or deck_path)
    prs = Presentation(str(deck_path))
    target, tbl = _get_table_element(
        prs.slides[slide_index], shape_id=shape_id, shape_name=shape_name,
    )
    n_cols = len(target.table.columns)
    if not (0 <= col_from < n_cols and 0 <= col_to < n_cols):
        raise IndexError(f"Column indices out of range (0..{n_cols-1})")
    if col_from == col_to:
        prs.save(out)
        return out

    a = NS["a"]
    new_order = list(range(n_cols))
    new_order.pop(col_from)
    new_order.insert(col_to, col_from)

    grid = tbl.find(f"{{{a}}}tblGrid")
    if grid is not None:
        _reorder_xml_children(grid, f"{{{a}}}gridCol", new_order)
    for tr in tbl.findall(f"{{{a}}}tr"):
        _reorder_xml_children(tr, f"{{{a}}}tc", new_order)

    prs.save(out)
    return out


@_logged("swap_table_rows")
def swap_table_rows(
    deck_path: str | Path,
    slide_index: int,
    *,
    shape_id: Optional[int] = None,
    shape_name: Optional[str] = None,
    row_a: int,
    row_b: int,
    output_path: Optional[str | Path] = None,
) -> str:
    """Swap two rows of a native PPTX table by 0-based index."""
    out = str(output_path or deck_path)
    prs = Presentation(str(deck_path))
    target, tbl = _get_table_element(
        prs.slides[slide_index], shape_id=shape_id, shape_name=shape_name,
    )
    n_rows = len(target.table.rows)
    if not (0 <= row_a < n_rows and 0 <= row_b < n_rows):
        raise IndexError(f"Row indices out of range (0..{n_rows-1})")
    if row_a == row_b:
        prs.save(out)
        return out

    a = NS["a"]
    new_order = list(range(n_rows))
    new_order[row_a], new_order[row_b] = new_order[row_b], new_order[row_a]
    _reorder_xml_children(tbl, f"{{{a}}}tr", new_order)
    prs.save(out)
    return out


# ─────────────────────────────────────────────────────────────────────────────
# Phase 4 — Semantic macros (named recipes built on Phase 1 primitives)
#
# Each macro is decorated with @_logged so the edit log records the macro
# name (e.g., "emphasize"), not the inner primitive calls. This keeps the
# replay history at the right abstraction level for cross-slide commands
# like "do the same on slide 4".
# ─────────────────────────────────────────────────────────────────────────────

def _slide_dims_inches(deck_path: str | Path) -> tuple[float, float]:
    prs = Presentation(str(deck_path))
    return prs.slide_width / 914400, prs.slide_height / 914400


def _shape_geometry_inches(deck_path: str | Path, slide_index: int,
                           *, shape_id=None, shape_name=None
                           ) -> tuple[float, float, float, float]:
    prs = Presentation(str(deck_path))
    sh = _find_shape(prs.slides[slide_index],
                     shape_id=shape_id, shape_name=shape_name)
    return (sh.left / 914400, sh.top / 914400,
            sh.width / 914400, sh.height / 914400)


def _current_run_size_pt(deck_path: str | Path, slide_index: int,
                         *, shape_id=None, shape_name=None
                         ) -> Optional[float]:
    prs = Presentation(str(deck_path))
    sh = _find_shape(prs.slides[slide_index],
                     shape_id=shape_id, shape_name=shape_name)
    if not sh.has_text_frame:
        return None
    for para in sh.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size is not None:
                return run.font.size.pt
    return None


@_logged("emphasize")
def emphasize(
    deck_path: str | Path,
    slide_index: int,
    *,
    shape_id: Optional[int] = None,
    shape_name: Optional[str] = None,
    brand_color: str = "#255BE3",
    scale: float = 1.25,
    bold: bool = True,
    output_path: Optional[str | Path] = None,
) -> str:
    """Make a shape stand out: bold + scale font + brand color."""
    out = str(output_path or deck_path)
    cur_size = _current_run_size_pt(
        deck_path, slide_index,
        shape_id=shape_id, shape_name=shape_name,
    )
    new_size = (cur_size or 14) * scale
    set_text_style(
        deck_path, slide_index,
        shape_id=shape_id, shape_name=shape_name,
        bold=bold, color=brand_color, size=new_size,
        output_path=out,
    )
    return out


@_logged("de_emphasize")
def de_emphasize(
    deck_path: str | Path,
    slide_index: int,
    *,
    shape_id: Optional[int] = None,
    shape_name: Optional[str] = None,
    mute_color: str = "#999999",
    scale: float = 0.85,
    output_path: Optional[str | Path] = None,
) -> str:
    """Tone down a shape: gray + smaller."""
    out = str(output_path or deck_path)
    cur_size = _current_run_size_pt(
        deck_path, slide_index,
        shape_id=shape_id, shape_name=shape_name,
    )
    new_size = (cur_size or 14) * scale
    set_text_style(
        deck_path, slide_index,
        shape_id=shape_id, shape_name=shape_name,
        bold=False, color=mute_color, size=new_size,
        output_path=out,
    )
    return out


@_logged("highlight_row")
def highlight_row(
    deck_path: str | Path,
    slide_index: int,
    shape_ids: list[int],
    *,
    fill_color: str = "#255BE3",
    text_color: str = "#FFFFFF",
    bold: bool = True,
    output_path: Optional[str | Path] = None,
) -> str:
    """
    Fill a list of shapes with brand color and set their text white+bold.
    Useful for highlighting a target column/row in a comparison matrix.
    """
    out = str(output_path or deck_path)
    src = str(deck_path)
    for sid in shape_ids:
        try:
            set_shape_fill(src, slide_index, shape_id=sid,
                           color=fill_color, output_path=out)
        except (AttributeError, ValueError):
            continue  # shape doesn't support fill (e.g. connector)
        try:
            set_text_style(out, slide_index, shape_id=sid,
                           color=text_color, bold=bold, output_path=out)
        except (AttributeError, ValueError):
            pass  # shape has no text frame
        src = out  # subsequent ops read latest state
    return out


@_logged("add_footnote")
def add_footnote(
    deck_path: str | Path,
    slide_index: int,
    text: str,
    *,
    color: str = "#666666",
    size: float = 9.0,
    rule_color: str = "#CCCCCC",
    output_path: Optional[str | Path] = None,
) -> str:
    """
    Add a footnote at bottom-left with a thin horizontal rule above it.
    Position auto-computed from slide dimensions.
    """
    out = str(output_path or deck_path)
    slide_w, slide_h = _slide_dims_inches(deck_path)
    margin = 0.5
    fn_top = slide_h - 0.55
    rule_top = fn_top - 0.06

    add_line(
        deck_path, slide_index,
        x1=margin, y1=rule_top,
        x2=slide_w - margin, y2=rule_top,
        color=rule_color, width=0.75, dash="solid",
        name="FootnoteRule", output_path=out,
    )
    add_textbox(
        out, slide_index,
        left=margin, top=fn_top,
        width=slide_w - 2 * margin, height=0.4,
        text=text, italic=True, color=color, size=size,
        name="Footnote", output_path=out,
    )
    return out


@_logged("add_section_label")
def add_section_label(
    deck_path: str | Path,
    slide_index: int,
    text: str,
    *,
    position: str = "top-left",
    fill: str = "#255BE3",
    text_color: str = "#FFFFFF",
    size: float = 10.0,
    output_path: Optional[str | Path] = None,
) -> str:
    """Add a small badge label. position: top-left / top-right / bottom-left / bottom-right."""
    out = str(output_path or deck_path)
    slide_w, slide_h = _slide_dims_inches(deck_path)
    badge_w, badge_h = 1.8, 0.35
    margin = 0.3
    if position == "top-left":
        left, top = margin, margin
    elif position == "top-right":
        left, top = slide_w - badge_w - margin, margin
    elif position == "bottom-left":
        left, top = margin, slide_h - badge_h - margin
    elif position == "bottom-right":
        left, top = slide_w - badge_w - margin, slide_h - badge_h - margin
    else:
        raise ValueError(
            f"Unknown position {position!r}. "
            "Use top-left / top-right / bottom-left / bottom-right"
        )
    add_shape_box(
        deck_path, slide_index,
        kind="capsule",
        left=left, top=top, width=badge_w, height=badge_h,
        fill=fill, text=text, text_color=text_color,
        text_size=size, bold=True, name="SectionLabel",
        output_path=out,
    )
    return out


@_logged("make_callout")
def make_callout(
    deck_path: str | Path,
    slide_index: int,
    target_shape_id: int,
    text: str,
    *,
    side: str = "right",
    brand_color: str = "#255BE3",
    width: float = 2.5,
    height: float = 0.5,
    output_path: Optional[str | Path] = None,
) -> str:
    """
    Add a capsule-shaped callout near a target shape with an arrow pointing
    at it. side: 'right' / 'left' / 'top' / 'bottom'.
    """
    out = str(output_path or deck_path)
    slide_w, slide_h = _slide_dims_inches(deck_path)
    tx, ty, tw, th = _shape_geometry_inches(
        deck_path, slide_index, shape_id=target_shape_id,
    )

    target_cx = tx + tw / 2
    target_cy = ty + th / 2

    if side == "right":
        callout_left = min(tx + tw + 0.5, slide_w - width - 0.2)
        callout_top = max(ty + th / 2 - height / 2, 0.2)
        line_x1, line_y1 = callout_left, callout_top + height / 2
        line_x2, line_y2 = tx + tw + 0.05, target_cy
    elif side == "left":
        callout_left = max(tx - 0.5 - width, 0.2)
        callout_top = max(ty + th / 2 - height / 2, 0.2)
        line_x1, line_y1 = callout_left + width, callout_top + height / 2
        line_x2, line_y2 = tx - 0.05, target_cy
    elif side == "top":
        callout_left = max(tx + tw / 2 - width / 2, 0.2)
        callout_top = max(ty - 0.6 - height, 0.2)
        line_x1, line_y1 = callout_left + width / 2, callout_top + height
        line_x2, line_y2 = target_cx, ty - 0.05
    elif side == "bottom":
        callout_left = max(tx + tw / 2 - width / 2, 0.2)
        callout_top = min(ty + th + 0.5, slide_h - height - 0.2)
        line_x1, line_y1 = callout_left + width / 2, callout_top
        line_x2, line_y2 = target_cx, ty + th + 0.05
    else:
        raise ValueError(
            f"Unknown side {side!r}. Use right / left / top / bottom"
        )

    add_shape_box(
        deck_path, slide_index,
        kind="capsule",
        left=callout_left, top=callout_top,
        width=width, height=height,
        fill=brand_color, text=text,
        text_color="#FFFFFF", text_size=11, bold=True,
        name="Callout", output_path=out,
    )
    add_line(
        out, slide_index,
        x1=line_x1, y1=line_y1, x2=line_x2, y2=line_y2,
        color=brand_color, width=1.5, dash="solid",
        name="CalloutArrow", output_path=out,
    )
    return out


@_logged("match_brand_style")
def match_brand_style(
    deck_path: str | Path,
    ref_deck_path: str | Path,
    *,
    output_path: Optional[str | Path] = None,
) -> str:
    """
    Apply theme palette from `ref_deck_path` to `deck_path`. Reads the
    reference deck's clrScheme and recolors the target's accent slots.
    Fonts are not migrated (they live in the master template).
    """
    out = str(output_path or deck_path)
    ref_members = _read_zip(str(ref_deck_path))
    theme_parts = [n for n in ref_members
                   if n.startswith("ppt/theme/") and n.endswith(".xml")]
    if not theme_parts:
        raise RuntimeError(
            f"Reference deck has no theme XML: {ref_deck_path}"
        )

    ref_root = _parse(ref_members[theme_parts[0]])
    a = NS["a"]
    scheme = ref_root.find(f".//{{{a}}}themeElements/{{{a}}}clrScheme")
    if scheme is None:
        raise RuntimeError(
            f"Reference deck theme has no clrScheme: {ref_deck_path}"
        )

    palette: dict[str, str] = {}
    for slot in THEME_SLOTS:
        slot_el = scheme.find(f"{{{a}}}{slot}")
        if slot_el is None:
            continue
        srgb = slot_el.find(f"{{{a}}}srgbClr")
        if srgb is not None and srgb.get("val"):
            palette[slot] = srgb.get("val")
            continue
        sys_el = slot_el.find(f"{{{a}}}sysClr")
        if sys_el is not None:
            last = sys_el.get("lastClr")
            if last:
                palette[slot] = last

    if not palette:
        raise RuntimeError(
            f"Could not extract any theme colors from {ref_deck_path}"
        )

    return recolor_theme(deck_path, palette, output_path=out)


# ─────────────────────────────────────────────────────────────────────────────
# CLI
# ─────────────────────────────────────────────────────────────────────────────

def _main() -> None:
    import argparse
    import json
    import sys

    parser = argparse.ArgumentParser(description="In-place .pptx editor")
    sub = parser.add_subparsers(dest="cmd", required=True)

    p_text = sub.add_parser("replace-text", help="Bulk text replace across deck")
    p_text.add_argument("path")
    p_text.add_argument("--map", required=True,
                        help='JSON of {"old":"new",...}')
    p_text.add_argument("--out", default=None)

    p_img = sub.add_parser("replace-image", help="Swap a picture shape")
    p_img.add_argument("path")
    p_img.add_argument("--slide", type=int, required=True)
    p_img.add_argument("--shape-name", default=None)
    p_img.add_argument("--shape-id", type=int, default=None)
    p_img.add_argument("--image", required=True)
    p_img.add_argument("--out", default=None)

    p_dup = sub.add_parser("duplicate-slide")
    p_dup.add_argument("path")
    p_dup.add_argument("--slide", type=int, required=True)
    p_dup.add_argument("--at", type=int, default=None)
    p_dup.add_argument("--out", default=None)

    p_del = sub.add_parser("delete-slide")
    p_del.add_argument("path")
    p_del.add_argument("--slide", type=int, required=True)
    p_del.add_argument("--out", default=None)

    p_re = sub.add_parser("reorder-slides")
    p_re.add_argument("path")
    p_re.add_argument("--order", required=True,
                      help="Comma-separated new order (0-based)")
    p_re.add_argument("--out", default=None)

    p_rc = sub.add_parser("recolor-theme")
    p_rc.add_argument("path")
    p_rc.add_argument("--palette", required=True,
                      help='JSON of {"accent1":"#RRGGBB",...}')
    p_rc.add_argument("--swap-hardcoded", default=None,
                      help='Optional JSON of {"OLDHEX":"NEWHEX",...}')
    p_rc.add_argument("--out", default=None)

    args = parser.parse_args()

    if args.cmd == "replace-text":
        out = replace_text_in_deck(args.path, json.loads(args.map), args.out)
    elif args.cmd == "replace-image":
        out = replace_picture(
            args.path, args.slide, args.image,
            shape_name=args.shape_name, shape_id=args.shape_id,
            output_path=args.out,
        )
    elif args.cmd == "duplicate-slide":
        out = duplicate_slide(args.path, args.slide,
                              position=args.at, output_path=args.out)
    elif args.cmd == "delete-slide":
        out = delete_slide(args.path, args.slide, args.out)
    elif args.cmd == "reorder-slides":
        order = [int(x) for x in args.order.split(",")]
        out = reorder_slides(args.path, order, args.out)
    elif args.cmd == "recolor-theme":
        swap = json.loads(args.swap_hardcoded) if args.swap_hardcoded else None
        out = recolor_theme(
            args.path, json.loads(args.palette),
            also_replace_hardcoded=swap, output_path=args.out,
        )
    else:
        parser.error("unknown command")
        return
    print(out)


if __name__ == "__main__":
    _main()
