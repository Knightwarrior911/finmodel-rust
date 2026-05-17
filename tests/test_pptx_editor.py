"""
Tests for src/research/pptx_editor.py.

Builds tiny synthetic decks via python-pptx, runs editor ops, then verifies
results structurally via pptx_inspector.diff_decks() and direct shape checks.
"""

from __future__ import annotations

import json
import shutil
import zipfile
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from src.research.pptx_editor import (
    THEME_SLOTS,
    add_footnote,
    add_line,
    add_section_label,
    add_shape_box,
    add_textbox,
    align_shapes,
    clear_edit_history,
    copy_style,
    de_emphasize,
    delete_shape,
    delete_slide,
    distribute_shapes,
    duplicate_slide,
    emphasize,
    find_shape_by_id,
    get_edit_history,
    highlight_row,
    iter_named_shapes,
    make_callout,
    match_brand_style,
    move_shape,
    move_table_column,
    recolor_theme,
    reorder_slides,
    replace_picture,
    replace_text_in_deck,
    replace_text_in_slide,
    resize_shape,
    set_shape_fill,
    set_shape_line,
    set_text_style,
    swap_table_columns,
    swap_table_rows,
)
from src.research.pptx_inspector import diff_decks, inspect_pptx


# ─────────────────────────────────────────────────────────────────────────────
# Fixtures
# ─────────────────────────────────────────────────────────────────────────────

def _make_three_slide_deck(path: Path) -> Path:
    """Slide 0: title. Slide 1: text+image. Slide 2: text only."""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]

    s0 = prs.slides.add_slide(blank)
    tb0 = s0.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    tb0.name = "Title"
    tb0.text_frame.text = "Cover Title"

    s1 = prs.slides.add_slide(blank)
    tb1 = s1.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    tb1.name = "Body"
    tb1.text_frame.text = "Q1 revenue grew strongly"

    img_path = _make_test_png(path.parent / "_logo.png")
    pic = s1.shapes.add_picture(
        str(img_path), Inches(8), Inches(0.5), Inches(1.5), Inches(1.5)
    )
    pic.name = "Logo"

    s2 = prs.slides.add_slide(blank)
    tb2 = s2.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    tb2.name = "Outro"
    tb2.text_frame.text = "Thank you — Q1 wrap-up"

    prs.save(path)
    return path


def _make_test_png(path: Path) -> Path:
    """Tiny solid-color PNG via raw bytes to avoid Pillow dependency."""
    # 8x8 red PNG, baked
    png_bytes = bytes.fromhex(
        "89504e470d0a1a0a0000000d4948445200000008000000080806000000c40fbe"
        "8b0000001b49444154789c63fcffff3f032003200a30c00100c40c0440c30202"
        "00f5170106e7e7c0c30000000049454e44ae426082"
    )
    path.write_bytes(png_bytes)
    return path


@pytest.fixture
def deck(tmp_path: Path) -> Path:
    return _make_three_slide_deck(tmp_path / "deck.pptx")


# ─────────────────────────────────────────────────────────────────────────────
# Text edits
# ─────────────────────────────────────────────────────────────────────────────

def test_replace_text_in_slide_counts_substitutions(deck: Path):
    prs = Presentation(deck)
    n = replace_text_in_slide(prs.slides[1], "Q1", "Q2")
    assert n == 1
    assert "Q2 revenue" in prs.slides[1].shapes[0].text_frame.text


def test_replace_text_in_deck_applies_across_all_slides(deck: Path, tmp_path: Path):
    out = tmp_path / "after.pptx"
    replace_text_in_deck(deck, {"Q1": "Q2"}, output_path=out)
    prs = Presentation(out)
    texts = [
        sh.text_frame.text
        for sl in prs.slides for sh in sl.shapes if sh.has_text_frame
    ]
    assert any("Q2 revenue" in t for t in texts)
    assert any("Q2 wrap-up" in t for t in texts)
    assert not any("Q1" in t for t in texts)


def test_replace_text_overwrite_in_place(deck: Path):
    replace_text_in_deck(deck, {"Q1": "Q2"})
    prs = Presentation(deck)
    texts = [
        sh.text_frame.text
        for sl in prs.slides for sh in sl.shapes if sh.has_text_frame
    ]
    assert any("Q2" in t for t in texts)
    assert not any("Q1" in t for t in texts)


# ─────────────────────────────────────────────────────────────────────────────
# Image swap
# ─────────────────────────────────────────────────────────────────────────────

def test_replace_picture_preserves_geometry(deck: Path, tmp_path: Path):
    new_img = _make_test_png(tmp_path / "new_logo.png")
    before = inspect_pptx(str(deck), include_raw_xml=False)
    logo_before = next(
        el for el in before["slides"][1]["elements"] if el.get("name") == "Logo"
    )
    box_before = logo_before["pos"]

    replace_picture(deck, slide_index=1, new_image_path=new_img, shape_name="Logo")

    after = inspect_pptx(str(deck), include_raw_xml=False)
    pics = [
        el for el in after["slides"][1]["elements"]
        if str(el.get("type", "")).startswith("PICTURE")
    ]
    assert pics, "expected a picture on slide 1 after swap"
    assert pics[0]["pos"] == box_before


def test_replace_picture_unknown_shape_raises(deck: Path, tmp_path: Path):
    new_img = _make_test_png(tmp_path / "x.png")
    with pytest.raises(LookupError):
        replace_picture(deck, slide_index=1, new_image_path=new_img,
                        shape_name="DoesNotExist")


# ─────────────────────────────────────────────────────────────────────────────
# Slide management (OOXML-aware)
# ─────────────────────────────────────────────────────────────────────────────

def _slide_count(path: Path) -> int:
    return len(Presentation(path).slides)


def test_duplicate_slide_appends(deck: Path, tmp_path: Path):
    out = tmp_path / "dup.pptx"
    duplicate_slide(deck, 0, output_path=out)
    assert _slide_count(out) == 4
    diff = diff_decks(str(deck), str(out))
    assert "slideCount" in " ".join(diff.get("differences", []))


def test_duplicate_slide_at_position(deck: Path, tmp_path: Path):
    out = tmp_path / "dup_at.pptx"
    duplicate_slide(deck, 0, position=1, output_path=out)
    prs = Presentation(out)
    assert len(prs.slides) == 4
    titles = [
        sl.shapes[0].text_frame.text if sl.shapes else ""
        for sl in prs.slides
    ]
    assert titles[0] == "Cover Title"
    assert titles[1] == "Cover Title"


def test_delete_slide_removes_one(deck: Path, tmp_path: Path):
    out = tmp_path / "del.pptx"
    delete_slide(deck, 1, output_path=out)
    prs = Presentation(out)
    assert len(prs.slides) == 2
    titles = [
        sl.shapes[0].text_frame.text if sl.shapes else ""
        for sl in prs.slides
    ]
    assert "Cover Title" in titles
    assert "Thank you — Q1 wrap-up" in titles
    # Deleted slide's body text gone
    all_text = " ".join(
        sh.text_frame.text for sl in prs.slides
        for sh in sl.shapes if sh.has_text_frame
    )
    assert "Q1 revenue grew strongly" not in all_text


def test_delete_slide_no_orphan_parts(deck: Path, tmp_path: Path):
    out = tmp_path / "del.pptx"
    delete_slide(deck, 1, output_path=out)
    with zipfile.ZipFile(out) as zf:
        names = zf.namelist()
    slide_parts = sorted(
        n for n in names if n.startswith("ppt/slides/slide") and n.endswith(".xml")
    )
    # Contiguous numbering starting from slide1
    assert slide_parts == ["ppt/slides/slide1.xml", "ppt/slides/slide2.xml"]


def test_reorder_slides_swaps(deck: Path, tmp_path: Path):
    out = tmp_path / "reord.pptx"
    reorder_slides(deck, [2, 1, 0], output_path=out)
    prs = Presentation(out)
    titles = [sl.shapes[0].text_frame.text for sl in prs.slides]
    assert titles[0] == "Thank you — Q1 wrap-up"
    assert titles[2] == "Cover Title"


def test_reorder_validates_permutation(deck: Path, tmp_path: Path):
    with pytest.raises(ValueError):
        reorder_slides(deck, [0, 0, 1], output_path=tmp_path / "bad.pptx")


# ─────────────────────────────────────────────────────────────────────────────
# Theme recolor
# ─────────────────────────────────────────────────────────────────────────────

def _theme_xml(path: Path) -> str:
    with zipfile.ZipFile(path) as zf:
        return zf.read("ppt/theme/theme1.xml").decode("utf-8")


def test_recolor_theme_writes_new_accent(deck: Path, tmp_path: Path):
    out = tmp_path / "recol.pptx"
    recolor_theme(
        deck,
        palette={"accent1": "#255BE3", "accent2": "0F1632"},
        output_path=out,
    )
    xml = _theme_xml(out)
    assert "255BE3" in xml
    assert "0F1632" in xml


def test_recolor_theme_rejects_unknown_slot(deck: Path, tmp_path: Path):
    with pytest.raises(KeyError):
        recolor_theme(deck, {"bogus": "#FFFFFF"},
                      output_path=tmp_path / "x.pptx")


def test_recolor_theme_rejects_bad_hex(deck: Path, tmp_path: Path):
    with pytest.raises(ValueError):
        recolor_theme(deck, {"accent1": "#ZZZ"},
                      output_path=tmp_path / "x.pptx")


def test_recolor_theme_replace_hardcoded(deck: Path, tmp_path: Path):
    # Inject a known srgbClr into a shape, then verify global swap
    prs = Presentation(deck)
    shape = prs.slides[0].shapes[0]
    from pptx.dml.color import RGBColor
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x44, 0x72, 0xC4)
    prs.save(deck)

    out = tmp_path / "swap.pptx"
    recolor_theme(
        deck,
        palette={"accent1": "#255BE3"},
        also_replace_hardcoded={"4472C4": "255BE3"},
        output_path=out,
    )
    with zipfile.ZipFile(out) as zf:
        slide1 = zf.read("ppt/slides/slide1.xml").decode("utf-8")
    assert "4472C4" not in slide1.upper()
    assert "255BE3" in slide1.upper()


def test_theme_slots_constant_is_complete():
    assert set(THEME_SLOTS) == {
        "dk1", "lt1", "dk2", "lt2",
        "accent1", "accent2", "accent3", "accent4", "accent5", "accent6",
        "hlink", "folHlink",
    }


# ─────────────────────────────────────────────────────────────────────────────
# Iteration helpers
# ─────────────────────────────────────────────────────────────────────────────

def test_iter_named_shapes_filters(deck: Path):
    prs = Presentation(deck)
    slide = prs.slides[1]
    bodies = list(iter_named_shapes(slide, prefix="Body"))
    assert len(bodies) == 1
    assert bodies[0].name == "Body"


def test_find_shape_by_id_returns_match(deck: Path):
    prs = Presentation(deck)
    slide = prs.slides[1]
    expected = next(sh for sh in slide.shapes if sh.name == "Logo")
    found = find_shape_by_id(slide, expected.shape_id)
    assert found is not None
    assert found.name == "Logo"


def test_find_shape_by_id_missing_returns_none(deck: Path):
    prs = Presentation(deck)
    assert find_shape_by_id(prs.slides[0], 999999) is None


# ─────────────────────────────────────────────────────────────────────────────
# Phase 1 primitives
# ─────────────────────────────────────────────────────────────────────────────

def _id_by_name(deck_path: Path, slide_idx: int, name: str) -> int:
    prs = Presentation(deck_path)
    sh = next(s for s in prs.slides[slide_idx].shapes if s.name == name)
    return sh.shape_id


def test_move_shape_absolute(deck: Path):
    sid = _id_by_name(deck, 1, "Logo")
    move_shape(deck, 1, shape_id=sid, left=2.0, top=3.0)
    prs = Presentation(deck)
    sh = next(s for s in prs.slides[1].shapes if s.shape_id == sid)
    assert abs(sh.left / 914400 - 2.0) < 0.01
    assert abs(sh.top / 914400 - 3.0) < 0.01


def test_move_shape_delta(deck: Path):
    sid = _id_by_name(deck, 1, "Logo")
    prs = Presentation(deck)
    before = next(s for s in prs.slides[1].shapes if s.shape_id == sid)
    before_left = before.left
    move_shape(deck, 1, shape_id=sid, dx=1.0)
    prs2 = Presentation(deck)
    after = next(s for s in prs2.slides[1].shapes if s.shape_id == sid)
    assert after.left - before_left > 0  # moved right


def test_move_shape_requires_input(deck: Path):
    sid = _id_by_name(deck, 1, "Logo")
    with pytest.raises(ValueError):
        move_shape(deck, 1, shape_id=sid)


def test_resize_shape(deck: Path):
    sid = _id_by_name(deck, 1, "Logo")
    resize_shape(deck, 1, shape_id=sid, width=3.0, height=2.0)
    prs = Presentation(deck)
    sh = next(s for s in prs.slides[1].shapes if s.shape_id == sid)
    assert abs(sh.width / 914400 - 3.0) < 0.01
    assert abs(sh.height / 914400 - 2.0) < 0.01


def test_set_shape_fill_solid(deck: Path):
    sid = _id_by_name(deck, 0, "Title")
    set_shape_fill(deck, 0, shape_id=sid, color="#E63946")
    prs = Presentation(deck)
    sh = next(s for s in prs.slides[0].shapes if s.shape_id == sid)
    assert str(sh.fill.fore_color.rgb).upper() == "E63946"


def test_set_shape_line(deck: Path, tmp_path: Path):
    sid = _id_by_name(deck, 0, "Title")
    set_shape_line(deck, 0, shape_id=sid, color="#000000",
                   width=2.0, dash="dash")
    # No exception is the assertion; load to confirm save valid
    Presentation(deck)


def test_set_text_style_bold_color_size(deck: Path):
    sid = _id_by_name(deck, 1, "Body")
    set_text_style(deck, 1, shape_id=sid, bold=True,
                   color="#255BE3", size=14)
    prs = Presentation(deck)
    sh = next(s for s in prs.slides[1].shapes if s.shape_id == sid)
    run = sh.text_frame.paragraphs[0].runs[0]
    assert run.font.bold is True
    assert run.font.size.pt == 14
    assert str(run.font.color.rgb).upper() == "255BE3"


def test_add_textbox(deck: Path):
    add_textbox(deck, 0, left=0.5, top=6.5, width=5, height=0.4,
                text="footnote", italic=True, color="#666666",
                size=10, name="FN")
    prs = Presentation(deck)
    fns = [s for s in prs.slides[0].shapes if s.name == "FN"]
    assert len(fns) == 1
    assert fns[0].text_frame.text == "footnote"


def test_add_line(deck: Path):
    add_line(deck, 0, x1=0.5, y1=6.5, x2=12, y2=6.5,
             color="#CCCCCC", width=1.0, name="Rule")
    prs = Presentation(deck)
    rules = [s for s in prs.slides[0].shapes if s.name == "Rule"]
    assert len(rules) == 1


def test_add_shape_box_capsule(deck: Path):
    add_shape_box(deck, 0, kind="capsule", left=10, top=0.2,
                  width=2.5, height=0.4, fill="#255BE3", text="DRAFT",
                  text_color="#FFFFFF", bold=True, name="Badge")
    prs = Presentation(deck)
    badges = [s for s in prs.slides[0].shapes if s.name == "Badge"]
    assert len(badges) == 1
    assert "DRAFT" in badges[0].text_frame.text


def test_add_shape_box_rejects_unknown_kind(deck: Path):
    with pytest.raises(ValueError):
        add_shape_box(deck, 0, kind="trapezoid", left=0, top=0,
                      width=1, height=1)


def test_delete_shape_by_name(deck: Path):
    add_textbox(deck, 0, left=0.5, top=6, width=2, height=0.4,
                text="zzz", name="ToRemove")
    delete_shape(deck, 0, shape_name="ToRemove")
    prs = Presentation(deck)
    assert not any(s.name == "ToRemove" for s in prs.slides[0].shapes)


def test_align_shapes_left(deck: Path):
    add_textbox(deck, 0, left=2, top=3, width=2, height=0.4,
                text="A", name="A")
    add_textbox(deck, 0, left=5, top=3.5, width=2, height=0.4,
                text="B", name="B")
    add_textbox(deck, 0, left=8, top=4, width=2, height=0.4,
                text="C", name="C")
    prs = Presentation(deck)
    ids = [s.shape_id for s in prs.slides[0].shapes
           if s.name in ("A", "B", "C")]
    align_shapes(deck, 0, ids, "left")
    prs2 = Presentation(deck)
    lefts = [s.left for s in prs2.slides[0].shapes
             if s.name in ("A", "B", "C")]
    assert len(set(lefts)) == 1


def test_distribute_shapes_x(deck: Path):
    add_textbox(deck, 0, left=1, top=3, width=1, height=0.4,
                text="A", name="A")
    add_textbox(deck, 0, left=5, top=3, width=1, height=0.4,
                text="B", name="B")
    add_textbox(deck, 0, left=9, top=3, width=1, height=0.4,
                text="C", name="C")
    prs = Presentation(deck)
    ids = [s.shape_id for s in prs.slides[0].shapes
           if s.name in ("A", "B", "C")]
    distribute_shapes(deck, 0, ids, "x")
    Presentation(deck)


def test_distribute_requires_three(deck: Path):
    add_textbox(deck, 0, left=1, top=3, width=1, height=0.4,
                text="A", name="A")
    add_textbox(deck, 0, left=5, top=3, width=1, height=0.4,
                text="B", name="B")
    prs = Presentation(deck)
    ids = [s.shape_id for s in prs.slides[0].shapes
           if s.name in ("A", "B")]
    with pytest.raises(ValueError):
        distribute_shapes(deck, 0, ids, "x")


# ─────────────────────────────────────────────────────────────────────────────
# Phase 4 macros
# ─────────────────────────────────────────────────────────────────────────────

def test_emphasize_applies_bold_color(deck: Path):
    sid = _id_by_name(deck, 1, "Body")
    emphasize(deck, 1, shape_id=sid, brand_color="#E63946", scale=1.5)
    prs = Presentation(deck)
    sh = next(s for s in prs.slides[1].shapes if s.shape_id == sid)
    run = sh.text_frame.paragraphs[0].runs[0]
    assert run.font.bold is True
    assert str(run.font.color.rgb).upper() == "E63946"


def test_de_emphasize_applies_gray(deck: Path):
    sid = _id_by_name(deck, 1, "Body")
    de_emphasize(deck, 1, shape_id=sid, mute_color="#888888", scale=0.8)
    prs = Presentation(deck)
    sh = next(s for s in prs.slides[1].shapes if s.shape_id == sid)
    run = sh.text_frame.paragraphs[0].runs[0]
    assert str(run.font.color.rgb).upper() == "888888"


def test_add_footnote_creates_two_shapes(deck: Path):
    add_footnote(deck, 0, "(1) Bloomberg consensus")
    prs = Presentation(deck)
    names = [s.name for s in prs.slides[0].shapes]
    assert "Footnote" in names
    assert "FootnoteRule" in names


def test_add_section_label_top_right(deck: Path):
    add_section_label(deck, 0, "DRAFT", position="top-right",
                      fill="#E63946")
    prs = Presentation(deck)
    badges = [s for s in prs.slides[0].shapes if s.name == "SectionLabel"]
    assert len(badges) == 1
    assert "DRAFT" in badges[0].text_frame.text
    # Top-right should have left near slide_width - badge_width
    slide_w = prs.slide_width / 914400
    assert badges[0].left / 914400 > slide_w / 2


def test_add_section_label_rejects_unknown_position(deck: Path):
    with pytest.raises(ValueError):
        add_section_label(deck, 0, "X", position="middle")


def test_make_callout_creates_capsule_and_arrow(deck: Path):
    sid = _id_by_name(deck, 1, "Logo")
    make_callout(deck, 1, target_shape_id=sid, text="Important",
                 side="left")
    prs = Presentation(deck)
    names = [s.name for s in prs.slides[1].shapes]
    assert "Callout" in names
    assert "CalloutArrow" in names


def test_make_callout_rejects_bad_side(deck: Path):
    sid = _id_by_name(deck, 1, "Logo")
    with pytest.raises(ValueError):
        make_callout(deck, 1, target_shape_id=sid, text="X",
                     side="diagonal")


def test_highlight_row_skips_shapes_without_fill(deck: Path):
    # Add a textbox + a connector line; highlight_row must handle the line
    add_textbox(deck, 0, left=1, top=2, width=2, height=0.4,
                text="Cell", name="Cell")
    add_line(deck, 0, x1=1, y1=2.5, x2=3, y2=2.5,
             color="#000000", name="Skip")
    prs = Presentation(deck)
    ids = [s.shape_id for s in prs.slides[0].shapes
           if s.name in ("Cell", "Skip")]
    # Should not raise even though Skip is a connector
    highlight_row(deck, 0, ids, fill_color="#255BE3")


def test_match_brand_style_extracts_palette(deck: Path, tmp_path: Path):
    # Create a ref deck with a known palette via recolor_theme
    ref = tmp_path / "ref.pptx"
    shutil.copy(deck, ref)
    recolor_theme(
        ref,
        palette={"accent1": "#E63946", "accent2": "#1D3557"},
    )
    target = tmp_path / "target.pptx"
    shutil.copy(deck, target)
    match_brand_style(target, ref)
    with zipfile.ZipFile(target) as zf:
        theme_xml = zf.read("ppt/theme/theme1.xml").decode("utf-8")
    assert "E63946" in theme_xml
    assert "1D3557" in theme_xml


# ─────────────────────────────────────────────────────────────────────────────
# Phase 5 — Edit log
# ─────────────────────────────────────────────────────────────────────────────

def test_edit_log_writes_entry_per_op(deck: Path):
    clear_edit_history(deck)
    add_textbox(deck, 0, left=1, top=1, width=1, height=0.5,
                text="x", name="X")
    sid = _id_by_name(deck, 0, "X")
    move_shape(deck, 0, shape_id=sid, dx=0.5)
    set_shape_fill(deck, 0, shape_id=sid, color="#FF0000")
    hist = get_edit_history(deck)
    ops = [e["op"] for e in hist]
    assert ops == ["add_textbox", "move_shape", "set_shape_fill"]


def test_edit_log_macro_does_not_log_inner_primitives(deck: Path):
    clear_edit_history(deck)
    sid = _id_by_name(deck, 1, "Body")
    emphasize(deck, 1, shape_id=sid)
    hist = get_edit_history(deck)
    ops = [e["op"] for e in hist]
    assert ops == ["emphasize"]
    # set_text_style should NOT appear because it ran nested under emphasize
    assert "set_text_style" not in ops


def test_edit_log_filter(deck: Path):
    clear_edit_history(deck)
    add_textbox(deck, 0, left=1, top=1, width=1, height=0.5, name="A")
    add_footnote(deck, 0, "footnote")
    sid = _id_by_name(deck, 0, "A")
    emphasize(deck, 0, shape_id=sid)

    macros = get_edit_history(deck, op_filter=["emphasize", "add_footnote"])
    assert len(macros) == 2
    assert {e["op"] for e in macros} == {"emphasize", "add_footnote"}


def test_clear_edit_history(deck: Path):
    add_textbox(deck, 0, left=1, top=1, width=1, height=0.5, name="A")
    clear_edit_history(deck)
    assert get_edit_history(deck) == []


def test_edit_log_path_uses_deck_extension(deck: Path):
    add_textbox(deck, 0, left=1, top=1, width=1, height=0.5, name="A")
    expected = deck.parent / (deck.name + ".edit_log.jsonl")
    assert expected.exists()


# ─────────────────────────────────────────────────────────────────────────────
# Native-table operations
# ─────────────────────────────────────────────────────────────────────────────

def _make_table_deck(path: Path) -> Path:
    """Build a 3-slide deck with one native PPTX table on slide 1."""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    prs.slides.add_slide(blank)
    s1 = prs.slides.add_slide(blank)
    table_shape = s1.shapes.add_table(
        rows=3, cols=4, left=Inches(1), top=Inches(1),
        width=Inches(10), height=Inches(3),
    )
    table_shape.name = "Matrix"
    headers = ["Metric", "Falcon", "Peer A", "Peer B"]
    for ci, h in enumerate(headers):
        table_shape.table.cell(0, ci).text_frame.text = h
    rows = [
        ["Revenue", "1250", "1800", "2100"],
        ["EBITDA",  "225",  "396",  "525"],
    ]
    for ri, row in enumerate(rows, start=1):
        for ci, v in enumerate(row):
            table_shape.table.cell(ri, ci).text_frame.text = v

    prs.slides.add_slide(blank)
    prs.save(path)
    return path


@pytest.fixture
def table_deck(tmp_path: Path) -> Path:
    return _make_table_deck(tmp_path / "tdeck.pptx")


def test_swap_table_columns(table_deck: Path):
    swap_table_columns(table_deck, 1, shape_name="Matrix",
                       col_a=1, col_b=2)
    prs = Presentation(table_deck)
    table = next(s for s in prs.slides[1].shapes if s.has_table).table
    headers = [table.cell(0, c).text_frame.text for c in range(4)]
    assert headers == ["Metric", "Peer A", "Falcon", "Peer B"]
    # Values stayed bonded to columns
    revenue_row = [table.cell(1, c).text_frame.text for c in range(4)]
    assert revenue_row == ["Revenue", "1800", "1250", "2100"]


def test_swap_table_columns_same_index_noop(table_deck: Path):
    swap_table_columns(table_deck, 1, shape_name="Matrix",
                       col_a=1, col_b=1)
    prs = Presentation(table_deck)
    table = next(s for s in prs.slides[1].shapes if s.has_table).table
    headers = [table.cell(0, c).text_frame.text for c in range(4)]
    assert headers == ["Metric", "Falcon", "Peer A", "Peer B"]


def test_swap_table_columns_out_of_range(table_deck: Path):
    with pytest.raises(IndexError):
        swap_table_columns(table_deck, 1, shape_name="Matrix",
                           col_a=0, col_b=99)


def test_swap_table_columns_wrong_shape_type(deck: Path):
    sid = _id_by_name(deck, 1, "Logo")
    with pytest.raises(ValueError):
        swap_table_columns(deck, 1, shape_id=sid, col_a=0, col_b=1)


def test_move_table_column(table_deck: Path):
    # Move Falcon (col 1) to position 3 (last column)
    move_table_column(table_deck, 1, shape_name="Matrix",
                      col_from=1, col_to=3)
    prs = Presentation(table_deck)
    table = next(s for s in prs.slides[1].shapes if s.has_table).table
    headers = [table.cell(0, c).text_frame.text for c in range(4)]
    assert headers == ["Metric", "Peer A", "Peer B", "Falcon"]


def test_swap_table_rows(table_deck: Path):
    swap_table_rows(table_deck, 1, shape_name="Matrix",
                    row_a=1, row_b=2)
    prs = Presentation(table_deck)
    table = next(s for s in prs.slides[1].shapes if s.has_table).table
    row1 = [table.cell(1, c).text_frame.text for c in range(4)]
    assert row1[0] == "EBITDA"


def test_copy_style(deck: Path):
    add_shape_box(deck, 0, kind="rect", left=1, top=2, width=2, height=1,
                  fill="#E63946", line="#000000", line_width=2,
                  name="Source")
    add_shape_box(deck, 0, kind="rect", left=4, top=2, width=2, height=1,
                  fill="#FFFFFF", name="Target")
    prs = Presentation(deck)
    src_id = next(s.shape_id for s in prs.slides[0].shapes
                  if s.name == "Source")
    tgt_id = next(s.shape_id for s in prs.slides[0].shapes
                  if s.name == "Target")
    copy_style(deck, 0, source_shape_id=src_id, target_shape_id=tgt_id)
    prs2 = Presentation(deck)
    tgt = next(s for s in prs2.slides[0].shapes if s.shape_id == tgt_id)
    assert str(tgt.fill.fore_color.rgb).upper() == "E63946"
