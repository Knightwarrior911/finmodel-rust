import json
from pptx import Presentation
from src.source_ledger import SourceLedger
from src.audit_pptx import annotate_pptx_with_sources


def _deck(tmp_path):
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[5])
    p = tmp_path / "d.pptx"; prs.save(p); return p


def _cache(tmp_path):
    led = SourceLedger()
    led.record_assumption("assumptions", "terminal_growth_rate", None,
                          value=0.025, rationale="GDP proxy", basis="house default")
    p = tmp_path / "c.json"
    p.write_text(json.dumps({"__ledger__": led.to_json()}), encoding="utf-8")
    return p


def test_sources_added_to_notes(tmp_path):
    deck = _deck(tmp_path)
    res = annotate_pptx_with_sources(str(deck), cache_path=str(_cache(tmp_path)))
    assert res["notes_added"] == 1
    prs = Presentation(str(deck))
    notes = prs.slides[0].notes_slide.notes_text_frame.text
    assert "Sources & Assumptions" in notes
    assert "terminal_growth_rate" in notes


def test_no_ledger_no_notes(tmp_path):
    deck = _deck(tmp_path)
    cp = tmp_path / "empty.json"; cp.write_text("{}", encoding="utf-8")
    res = annotate_pptx_with_sources(str(deck), cache_path=str(cp))
    assert res["notes_added"] == 0


def test_none_cache_path_safe(tmp_path):
    deck = _deck(tmp_path)
    assert annotate_pptx_with_sources(str(deck), cache_path=None)["notes_added"] == 0
