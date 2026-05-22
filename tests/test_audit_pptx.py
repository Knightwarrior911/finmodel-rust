"""Tests for src.audit_pptx — number parsing + deck source-linking."""
from __future__ import annotations

import json

import pytest

from src.audit_pptx import _parse_number, annotate_pptx_with_links


@pytest.mark.parametrize("text,expected", [
    ("172,664", 172664.0),
    ("1,234.5", 1234.5),
    ("(45)", -45.0),
    ("45.2%", 0.452),
    ("$1,068,000", 1068000.0),
    ("Revenue grew to 172,664 MSEK", 172664.0),
    ("n/a", None),
    ("", None),
])
def test_parse_number(text, expected):
    assert _parse_number(text) == expected


def test_annotate_pptx_links_market_value(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches
    from src.citations import MarketCitation

    cache = tmp_path / "X.json"
    cache.write_text(json.dumps({
        "years_found": ["2024"],
        "__citations__": [MarketCitation(
            1068000.0, "MSFT market cap",
            "https://finance.yahoo.com/quote/MSFT", "yfinance", "2026-05-22"
        ).to_json()],
    }), encoding="utf-8")

    # Deck with a 1x2 table: label cell "MSFT", value cell "1,068,000".
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tbl = slide.shapes.add_table(1, 2, Inches(1), Inches(1), Inches(5), Inches(1)).table
    tbl.cell(0, 0).text = "MSFT"
    tbl.cell(0, 1).text = "1,068,000"
    deck = tmp_path / "d.pptx"
    prs.save(str(deck))

    res = annotate_pptx_with_links(deck, cache_path=cache)
    assert res["linked_market"] == 1

    prs2 = Presentation(str(deck))
    tbl2 = next(s.table for s in prs2.slides[0].shapes if s.has_table)
    run = tbl2.cell(0, 1).text_frame.paragraphs[0].runs[0]
    assert run.hyperlink.address == "https://finance.yahoo.com/quote/MSFT"


def test_annotate_pptx_no_cache_is_noop(tmp_path):
    from pptx import Presentation
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[5])
    deck = tmp_path / "d.pptx"
    prs.save(str(deck))
    assert annotate_pptx_with_links(deck, cache_path=None)["total"] == 0
